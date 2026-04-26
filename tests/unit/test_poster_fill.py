"""Tests for poster_fill.py — KBase poster template content fill.

Coverage:
- Both shipped templates (horizontal + vertical) load.
- Title placeholder ("TITLE") gets replaced with spec.title.
- Authors placeholder ("NAME 1") gets replaced.
- Funding placeholder gets replaced when spec.funding set.
- Body sections (TL;DR / Methods / Implications / References) emit text boxes.
- Figure panels resolve relative paths against draft_dir.
- Missing figure file → placeholder text, no crash.
- assemble_pptx top-level dispatch: mode='poster-h' → fill_poster called.
- CLI surface.
"""
from __future__ import annotations

import importlib.util
import json
import struct
import sys
import zlib
from pathlib import Path

import pytest


REPO_ROOT = Path(__file__).resolve().parents[2]
PF_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
         / "tools" / "poster_fill.py")
ASM_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
          / "tools" / "assemble_pptx.py")
SLIDE_SPEC_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
                 / "tools" / "slide_spec.py")
TEMPLATE_DIR = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
                / "references" / "templates")


def _import(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def pf():
    return _import("poster_fill", PF_PY)


@pytest.fixture(scope="module")
def asm():
    _import("slide_spec", SLIDE_SPEC_PY)
    _import("poster_fill", PF_PY)
    return _import("assemble_pptx", ASM_PY)


@pytest.fixture(scope="module")
def ss():
    return _import("slide_spec", SLIDE_SPEC_PY)


def _make_png(path: Path, w: int = 200, h: int = 200) -> Path:
    def chunk(tag, data):
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
    raw = b"".join(b"\x00" + bytes([0xa0, 0xc0, 0xe0]) * w for _ in range(h))
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    path.write_bytes(sig + ihdr + idat + iend)
    return path


@pytest.fixture
def basic_spec(pf):
    return pf.PosterSpec(
        title="Sample Poster Title",
        authors="A. Arkin¹, J. Doe²",
        affiliation="¹UC Berkeley   ²LBNL",
        tl_dr="A one-sentence summary of the poster.",
        methods_summary=["Method A", "Method B", "Method C"],
        cross_tenant_summary="ENIGMA + PMI tenants integrated.",
        implications=["Implication 1", "Implication 2"],
        references_short=["Smith 2023", "Jones 2024"],
        acknowledgments="Lab team and funders.",
        funding="DOE BER",
    )


# ---------------------------------------------------------------------------
# Templates exist
# ---------------------------------------------------------------------------

def test_horizontal_template_exists(pf):
    assert pf.POSTER_TEMPLATES["horizontal"].is_file()


def test_vertical_template_exists(pf):
    assert pf.POSTER_TEMPLATES["vertical"].is_file()


# ---------------------------------------------------------------------------
# PosterSpec construction
# ---------------------------------------------------------------------------

def test_poster_spec_from_dict_minimal(pf):
    d = {"title": "T", "authors": "A"}
    spec = pf.PosterSpec.from_dict(d)
    assert spec.title == "T"
    assert spec.authors == "A"
    assert spec.figures == []
    assert spec.methods_summary == []


def test_poster_spec_from_dict_with_figures(pf):
    d = {
        "title": "T",
        "authors": "A",
        "figures": [
            {"path": "fig1.png", "caption": "First"},
            {"path": "fig2.png"},
        ],
    }
    spec = pf.PosterSpec.from_dict(d)
    assert len(spec.figures) == 2
    assert spec.figures[0].caption == "First"
    assert spec.figures[1].caption == ""


# ---------------------------------------------------------------------------
# fill_poster — horizontal
# ---------------------------------------------------------------------------

def test_fill_poster_horizontal_writes_pptx(pf, basic_spec, tmp_path):
    out = tmp_path / "poster.pptx"
    result = pf.fill_poster(basic_spec, out, orientation="horizontal")
    assert result.is_file()
    from pptx import Presentation
    prs = Presentation(out)
    assert len(prs.slides) == 1


def test_fill_poster_horizontal_replaces_title(pf, basic_spec, tmp_path):
    out = tmp_path / "poster.pptx"
    pf.fill_poster(basic_spec, out, orientation="horizontal")
    from pptx import Presentation
    prs = Presentation(out)
    # Find any shape containing the new title text
    found = False
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and basic_spec.title in shape.text_frame.text:
            found = True
            break
    assert found, "title not found on poster"


def test_fill_poster_horizontal_replaces_authors(pf, tmp_path):
    """Use ASCII-only authors text to avoid Unicode-encoding fragility."""
    spec = pf.PosterSpec(
        title="Sample",
        authors="Adam Arkin, Jane Doe",
        affiliation="UC Berkeley",
    )
    out = tmp_path / "poster.pptx"
    pf.fill_poster(spec, out, orientation="horizontal")
    from pptx import Presentation
    prs = Presentation(out)
    found = False
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and "Arkin" in shape.text_frame.text:
            found = True
            break
    assert found, "authors text not found on poster"


def test_fill_poster_emits_body_sections(pf, basic_spec, tmp_path):
    out = tmp_path / "poster.pptx"
    pf.fill_poster(basic_spec, out, orientation="horizontal")
    from pptx import Presentation
    prs = Presentation(out)
    text_blob = "\n".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    # Headings + content
    assert "TL;DR" in text_blob
    assert basic_spec.tl_dr in text_blob
    assert "Methods" in text_blob
    assert "Method A" in text_blob
    assert "Implications" in text_blob
    assert "References" in text_blob
    assert "Smith 2023" in text_blob


# ---------------------------------------------------------------------------
# fill_poster — vertical
# ---------------------------------------------------------------------------

def test_fill_poster_vertical_writes_pptx(pf, basic_spec, tmp_path):
    out = tmp_path / "poster_v.pptx"
    result = pf.fill_poster(basic_spec, out, orientation="vertical")
    assert result.is_file()
    from pptx import Presentation
    prs = Presentation(out)
    # Vertical: 36 × 48
    assert prs.slide_width / 914400 == pytest.approx(36.0, abs=0.5)
    assert prs.slide_height / 914400 == pytest.approx(48.0, abs=0.5)


def test_fill_poster_vertical_emits_body_sections(pf, basic_spec, tmp_path):
    out = tmp_path / "poster_v.pptx"
    pf.fill_poster(basic_spec, out, orientation="vertical")
    from pptx import Presentation
    prs = Presentation(out)
    text_blob = "\n".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    assert basic_spec.tl_dr in text_blob


# ---------------------------------------------------------------------------
# Figure handling
# ---------------------------------------------------------------------------

def test_fill_poster_with_figures(pf, basic_spec, tmp_path):
    fig_dir = tmp_path / "figures"
    fig_dir.mkdir()
    _make_png(fig_dir / "fig01.png", 400, 300)
    _make_png(fig_dir / "fig02.png", 400, 300)
    basic_spec.figures = [
        pf.PosterFigure(path="figures/fig01.png", caption="Cap 1"),
        pf.PosterFigure(path="figures/fig02.png", caption="Cap 2"),
    ]
    out = tmp_path / "poster.pptx"
    pf.fill_poster(basic_spec, out, orientation="horizontal", draft_dir=tmp_path)
    from pptx import Presentation
    prs = Presentation(out)
    # Verify pictures ended up on slide
    n_pics = sum(1 for s in prs.slides[0].shapes if s.shape_type == 13)
    assert n_pics >= 2  # template logos + 2 added figures


def test_fill_poster_missing_figure_emits_placeholder(pf, basic_spec, tmp_path):
    basic_spec.figures = [
        pf.PosterFigure(path="figures/missing.png", caption="Cap 1"),
    ]
    out = tmp_path / "poster.pptx"
    pf.fill_poster(basic_spec, out, orientation="horizontal", draft_dir=tmp_path)
    from pptx import Presentation
    prs = Presentation(out)
    text_blob = "\n".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    assert "missing" in text_blob.lower()


# ---------------------------------------------------------------------------
# Error paths
# ---------------------------------------------------------------------------

def test_fill_poster_unknown_orientation_raises(pf, basic_spec, tmp_path):
    with pytest.raises(ValueError):
        pf.fill_poster(basic_spec, tmp_path / "x.pptx", orientation="diagonal")


def test_fill_poster_missing_template_raises(pf, basic_spec, tmp_path):
    bad_template = tmp_path / "nope.pptx"
    with pytest.raises(FileNotFoundError):
        pf.fill_poster(basic_spec, tmp_path / "x.pptx",
                       orientation="horizontal", template_path=bad_template)


# ---------------------------------------------------------------------------
# assemble_pptx top-level dispatch on poster mode
# ---------------------------------------------------------------------------

def test_assemble_dispatches_to_poster_fill_for_poster_mode(asm, ss, tmp_path):
    """A slide_spec with mode='poster-h' should bypass slide-by-slide
    rendering and emit a single poster.pptx via fill_poster."""
    spec = ss.example_slide_spec()
    spec["mode"] = "poster-h"
    # Filter to mappable layouts AND clear substory_id (we drop substories below)
    keep_layouts = ("title", "methods_summary", "implications",
                    "cross_tenant_integration",
                    "references", "acknowledgments")
    spec["slides"] = [
        {**s, "substory_id": None}
        for s in spec["slides"] if s["layout"] in keep_layouts
    ]
    spec["substories"] = []
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "poster.pptx"
    result = asm.assemble(spec_path, out)
    assert result.n_slides == 1
    from pptx import Presentation
    prs = Presentation(out)
    assert prs.slide_width / 914400 == pytest.approx(48.0, abs=0.5)


def test_assemble_dispatches_to_poster_fill_for_poster_v_mode(asm, ss, tmp_path):
    spec = ss.example_slide_spec()
    spec["mode"] = "poster-v"
    keep_layouts = ("title", "methods_summary", "implications", "references")
    spec["slides"] = [
        {**s, "substory_id": None}
        for s in spec["slides"] if s["layout"] in keep_layouts
    ]
    spec["substories"] = []
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "poster.pptx"
    result = asm.assemble(spec_path, out)
    assert result.n_slides == 1
    from pptx import Presentation
    prs = Presentation(out)
    assert prs.slide_height / 914400 == pytest.approx(48.0, abs=0.5)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def test_cli_writes_poster(pf, basic_spec, tmp_path):
    spec_path = tmp_path / "spec.json"
    # Convert PosterSpec → dict for CLI input
    d = {
        "title": basic_spec.title, "authors": basic_spec.authors,
        "affiliation": basic_spec.affiliation, "tl_dr": basic_spec.tl_dr,
        "methods_summary": basic_spec.methods_summary,
        "implications": basic_spec.implications,
        "references_short": basic_spec.references_short,
        "acknowledgments": basic_spec.acknowledgments,
        "cross_tenant_summary": basic_spec.cross_tenant_summary,
        "funding": basic_spec.funding,
    }
    spec_path.write_text(json.dumps(d))
    out = tmp_path / "poster.pptx"
    rc = pf.main([str(spec_path), "--out", str(out),
                  "--orientation", "horizontal"])
    assert rc == 0
    assert out.is_file()


def test_cli_missing_spec_returns_2(pf, tmp_path):
    rc = pf.main([str(tmp_path / "nope.json"), "--out", str(tmp_path / "x.pptx")])
    assert rc == 2

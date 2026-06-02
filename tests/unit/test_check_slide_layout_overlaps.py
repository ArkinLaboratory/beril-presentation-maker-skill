"""v0.8.0 Tier G.10-A: deterministic bounding-box overlap detector.

Pure-geometry tests at the helper level — no real .pptx needed
for most cases (the geometry primitives + classification + allow-
list logic are decoupled from python-pptx). A handful of
integration tests construct minimal Presentation objects via
python-pptx to verify the end-to-end check.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

TOOLS_DIR = (
    Path(__file__).resolve().parents[2]
    / "src/beril_presentation_maker/skill/tools"
)
sys.path.insert(0, str(TOOLS_DIR))

import check_slide_layout_overlaps as overlaps_mod  # noqa: E402
from check_slide_layout_overlaps import (  # noqa: E402
    Rect, overlaps, intersection, contained_in,
    _classify_pair, load_allow_list,
    detect_container_breaches, detect_pairwise_overlaps,
    OverlapFinding, OverlapReport, render_markdown,
)


# ---------------------------------------------------------------------------
# Geometry primitives
# ---------------------------------------------------------------------------

def test_rect_right_bottom_area() -> None:
    r = Rect(left=100, top=200, width=300, height=400)
    assert r.right == 400
    assert r.bottom == 600
    assert r.area() == 300 * 400


def test_rect_with_zero_dimensions_area_is_zero() -> None:
    """Defensive: a degenerate rect has zero area, not negative."""
    assert Rect(0, 0, 0, 0).area() == 0
    assert Rect(0, 0, 100, 0).area() == 0


# ---------------------------------------------------------------------------
# overlaps() predicate
# ---------------------------------------------------------------------------

def test_overlaps_disjoint_horizontal() -> None:
    a = Rect(0, 0, 100, 100)
    b = Rect(200, 0, 100, 100)
    assert overlaps(a, b, pad_emu=0) is False


def test_overlaps_disjoint_vertical() -> None:
    a = Rect(0, 0, 100, 100)
    b = Rect(0, 200, 100, 100)
    assert overlaps(a, b, pad_emu=0) is False


def test_overlaps_clearly_intersecting() -> None:
    a = Rect(0, 0, 200, 200)
    b = Rect(100, 100, 200, 200)
    assert overlaps(a, b, pad_emu=0) is True


def test_overlaps_padding_tolerance_suppresses_small_touches() -> None:
    """Two rects that touch at the edge (overlap by 1 EMU) should
    NOT trigger when padding > 1 (the v0.8.0 default of 36000)."""
    a = Rect(0, 0, 100, 100)
    b = Rect(99, 0, 100, 100)  # overlaps by 1 EMU horizontally
    assert overlaps(a, b, pad_emu=0) is True
    assert overlaps(a, b, pad_emu=100) is False  # padding swallows


def test_overlaps_padding_does_not_suppress_real_overlap() -> None:
    """Padding should NOT swallow genuine overlaps when the overlap
    is larger than the padding. Use EMU-scale dimensions (the real
    deployment scale: 9.1M EMU slide width, 36000 EMU padding)."""
    a = Rect(0, 0, 2000000, 2000000)
    b = Rect(1000000, 1000000, 2000000, 2000000)  # 1M×1M overlap
    assert overlaps(a, b, pad_emu=36000) is True


# ---------------------------------------------------------------------------
# intersection() + contained_in()
# ---------------------------------------------------------------------------

def test_intersection_disjoint_is_none() -> None:
    assert intersection(Rect(0, 0, 100, 100), Rect(200, 0, 100, 100)) is None


def test_intersection_basic() -> None:
    a = Rect(0, 0, 200, 200)
    b = Rect(100, 100, 200, 200)
    isect = intersection(a, b)
    assert isect == Rect(100, 100, 100, 100)


def test_contained_in_strict() -> None:
    """A rect entirely inside another is contained."""
    inner = Rect(10, 10, 80, 80)
    outer = Rect(0, 0, 100, 100)
    assert contained_in(inner, outer, pad_emu=0) is True


def test_contained_in_breach_is_false() -> None:
    """A rect that pokes outside is not contained."""
    inner = Rect(50, 50, 100, 100)  # extends to (150, 150)
    outer = Rect(0, 0, 100, 100)
    assert contained_in(inner, outer, pad_emu=0) is False


# ---------------------------------------------------------------------------
# Pair classification
# ---------------------------------------------------------------------------

def test_classify_image_plus_text() -> None:
    assert _classify_pair("image", "text") == "image_text_overlap"
    assert _classify_pair("text", "image") == "image_text_overlap"


def test_classify_image_plus_title() -> None:
    assert _classify_pair("image", "title") == "image_text_overlap"


def test_classify_title_plus_text() -> None:
    assert _classify_pair("title", "text") == "footer_title_collision"


def test_classify_text_plus_text() -> None:
    assert _classify_pair("text", "text") == "text_box_overlap"


def test_classify_other_plus_other_is_uninteresting() -> None:
    assert _classify_pair("other", "other") is None


def test_classify_other_plus_text_is_uninteresting() -> None:
    """A non-classified shape (decoration, group) overlapping a text
    box is NOT a finding — we only flag the four named kinds."""
    assert _classify_pair("other", "text") is None


# ---------------------------------------------------------------------------
# Container-breach detection
# ---------------------------------------------------------------------------

class _FakeShape:
    """Lightweight stand-in for python-pptx shapes in tests."""
    def __init__(self, name: str = "shape"):
        self.name = name


def test_breach_inside_canvas_no_finding() -> None:
    shapes = [(_FakeShape("A"), Rect(100, 100, 500, 500))]
    findings = detect_container_breaches(
        shapes, slide_id=1, layout_name="claim_evidence",
        slide_width=9144000, slide_height=5143500,
        pad_emu=36000,
    )
    assert findings == []


def test_breach_extends_off_right_edge_fires_p0() -> None:
    shapes = [(_FakeShape("Wide"),
               Rect(8000000, 1000000, 2000000, 1000000))]
    # right edge = 10M; slide width = 9144000; breach by 856000 EMU
    findings = detect_container_breaches(
        shapes, slide_id=5, layout_name="data_figure",
        slide_width=9144000, slide_height=5143500,
        pad_emu=36000,
    )
    assert len(findings) == 1
    f = findings[0]
    assert f.kind == "container_breach"
    assert f.severity == "P0"
    assert f.slide_id == 5
    assert f.layout_name == "data_figure"
    assert f.shape_a == "Wide"
    assert f.overlap_area_emu > 0


def test_breach_padding_tolerance_swallows_small_overruns() -> None:
    """A shape that pokes out by less than pad_emu on each side
    should NOT trigger (rendering slop, not a real breach)."""
    # Slide is 100×100; shape extends to (101, 101) — 1 EMU breach.
    shapes = [(_FakeShape("Tiny"), Rect(0, 0, 101, 101))]
    findings = detect_container_breaches(
        shapes, slide_id=1, layout_name="x",
        slide_width=100, slide_height=100,
        pad_emu=10,
    )
    assert findings == []


# ---------------------------------------------------------------------------
# Pairwise-overlap detection
# ---------------------------------------------------------------------------

class _FakeTextShape(_FakeShape):
    """A python-pptx-shaped fake that classifies as 'text'."""
    shape_type = 17  # TEXT_BOX


class _FakeImageShape(_FakeShape):
    shape_type = 13  # PICTURE


class _FakeTitleShape(_FakeShape):
    shape_type = 17  # TEXT_BOX
    # title detection via the name (per shape_role)


def test_no_overlap_no_findings() -> None:
    shapes = [
        (_FakeTextShape("Body 1"), Rect(0, 0, 1000, 1000)),
        (_FakeTextShape("Body 2"), Rect(2000, 0, 1000, 1000)),
    ]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=1, layout_name="x", pad_emu=10, allow_zones=[],
    )
    assert findings == []


def test_text_text_overlap_classified() -> None:
    shapes = [
        (_FakeTextShape("Body 1"), Rect(0, 0, 1000, 1000)),
        (_FakeTextShape("Body 2"), Rect(500, 500, 1000, 1000)),
    ]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=3, layout_name="layout_a",
        pad_emu=10, allow_zones=[],
    )
    assert len(findings) == 1
    f = findings[0]
    assert f.kind == "text_box_overlap"
    assert f.severity == "P1"
    assert f.slide_id == 3
    assert f.layout_name == "layout_a"
    assert f.overlap_fraction == pytest.approx(0.25, abs=1e-3)


def test_image_text_overlap_classified() -> None:
    shapes = [
        (_FakeImageShape("Picture 1"), Rect(0, 0, 1000, 1000)),
        (_FakeTextShape("Body 1"), Rect(500, 500, 1000, 1000)),
    ]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=2, layout_name="big_idea",
        pad_emu=10, allow_zones=[],
    )
    assert len(findings) == 1
    assert findings[0].kind == "image_text_overlap"


def test_title_text_overlap_classified_as_footer_title_collision() -> None:
    shapes = [
        (_FakeTitleShape("Title 1"), Rect(0, 0, 1000, 1000)),
        (_FakeTextShape("Body 1"), Rect(500, 500, 1000, 1000)),
    ]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=4, layout_name="claim_evidence",
        pad_emu=10, allow_zones=[],
    )
    assert len(findings) == 1
    assert findings[0].kind == "footer_title_collision"


def test_allow_list_suppresses_overlap_inside_allowed_zone() -> None:
    """An overlap entirely inside an allow-zone rect must be
    suppressed (intentional design overlap, e.g., footer watermark
    behind a chrome band)."""
    shapes = [
        (_FakeTextShape("A"), Rect(100, 100, 500, 500)),
        (_FakeTextShape("B"), Rect(150, 150, 400, 400)),
    ]
    # Allow-zone that contains the entire overlap rect (150,150)-(600,600)
    allow = [Rect(0, 0, 1000, 1000)]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=1, layout_name="x",
        pad_emu=0, allow_zones=allow,
    )
    assert findings == []


def test_allow_list_does_not_suppress_overlap_outside_zone() -> None:
    """Sanity: an overlap that's NOT inside the allow-zone still
    fires."""
    shapes = [
        (_FakeTextShape("A"), Rect(100, 100, 500, 500)),
        (_FakeTextShape("B"), Rect(300, 300, 500, 500)),
    ]
    # Allow-zone that misses the overlap (overlap is at 300-600, allow at 0-100)
    allow = [Rect(0, 0, 100, 100)]
    findings = detect_pairwise_overlaps(
        shapes, slide_id=1, layout_name="x",
        pad_emu=0, allow_zones=allow,
    )
    assert len(findings) == 1


# ---------------------------------------------------------------------------
# Allow-list loading
# ---------------------------------------------------------------------------

def test_load_allow_list_missing_file_is_empty() -> None:
    assert load_allow_list(None) == {}
    assert load_allow_list(Path("/no/such/file.json")) == {}


def test_load_allow_list_valid_json(tmp_path: Path) -> None:
    target = tmp_path / "allow.json"
    target.write_text(json.dumps({
        "footer_band_layout": [
            {"left": 0, "top": 5000000, "width": 9144000, "height": 200000}
        ],
        "title_only": [],  # empty list → no zones
    }), encoding="utf-8")
    loaded = load_allow_list(target)
    assert "footer_band_layout" in loaded
    assert len(loaded["footer_band_layout"]) == 1
    assert loaded["footer_band_layout"][0] == Rect(0, 5000000, 9144000, 200000)


def test_load_allow_list_malformed_json_is_empty(tmp_path: Path) -> None:
    target = tmp_path / "broken.json"
    target.write_text("not a json", encoding="utf-8")
    assert load_allow_list(target) == {}


def test_load_allow_list_skips_malformed_zones(tmp_path: Path) -> None:
    """If one zone in a layout is malformed (missing keys), the rest
    of that layout's zones still load."""
    target = tmp_path / "mixed.json"
    target.write_text(json.dumps({
        "layout_x": [
            {"left": 0, "top": 0, "width": 100, "height": 100},
            {"left": 200},  # missing keys — skip this one
            {"left": 300, "top": 0, "width": 100, "height": 100},
        ],
    }), encoding="utf-8")
    loaded = load_allow_list(target)
    assert len(loaded["layout_x"]) == 2  # 1 bad zone dropped


# ---------------------------------------------------------------------------
# Markdown rendering
# ---------------------------------------------------------------------------

def test_markdown_zero_findings() -> None:
    rep = OverlapReport(pptx_path="/path/to/x.pptx")
    md = render_markdown(rep)
    assert "Total findings:** 0" in md
    assert "No layout overlaps" in md


def test_markdown_groups_by_severity() -> None:
    rep = OverlapReport(pptx_path="/x.pptx", findings=[
        OverlapFinding(
            kind="container_breach", severity="P0", slide_id=5,
            layout_name="data_figure", shape_a="X", shape_b="",
            overlap_area_emu=1000, overlap_fraction=0.1,
            message="X extends past canvas",
        ),
        OverlapFinding(
            kind="text_box_overlap", severity="P1", slide_id=8,
            layout_name="claim_evidence", shape_a="A", shape_b="B",
            overlap_area_emu=2000, overlap_fraction=0.2,
            message="A overlaps B",
        ),
    ])
    md = render_markdown(rep)
    assert "P0 — container breaches (1)" in md
    assert "P1 — overlaps (1)" in md
    assert "slide 5" in md
    assert "slide 8" in md


# ---------------------------------------------------------------------------
# Integration: real python-pptx Presentation
# ---------------------------------------------------------------------------

def test_check_pptx_on_minimal_deck_no_overlaps(tmp_path: Path) -> None:
    """A minimal one-slide deck with non-overlapping shapes should
    produce 0 findings."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Two non-overlapping text boxes
    slide.shapes.add_textbox(Emu(914400), Emu(914400),
                             Emu(2000000), Emu(1000000)).text_frame.text = "A"
    slide.shapes.add_textbox(Emu(914400), Emu(3000000),
                             Emu(2000000), Emu(1000000)).text_frame.text = "B"
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    rep = overlaps_mod.check_pptx(pptx_path)
    assert rep.findings == []


def test_check_pptx_detects_text_overlap(tmp_path: Path) -> None:
    """A deck with two text boxes that genuinely overlap should
    produce one text_box_overlap finding."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_textbox(Emu(914400), Emu(914400),
                             Emu(3000000), Emu(2000000)).text_frame.text = "A"
    slide.shapes.add_textbox(Emu(2000000), Emu(2000000),
                             Emu(3000000), Emu(2000000)).text_frame.text = "B"
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    rep = overlaps_mod.check_pptx(pptx_path)
    kinds = [f.kind for f in rep.findings]
    assert "text_box_overlap" in kinds


def test_check_pptx_detects_container_breach(tmp_path: Path) -> None:
    """A shape that extends past the slide edge fires container_breach P0."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Slide width default 9144000 EMU; place a wide shape that ends at 11M
    slide.shapes.add_textbox(Emu(8000000), Emu(1000000),
                             Emu(3000000), Emu(1000000)).text_frame.text = "X"
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    rep = overlaps_mod.check_pptx(pptx_path)
    breaches = [f for f in rep.findings if f.kind == "container_breach"]
    assert len(breaches) == 1
    assert breaches[0].severity == "P0"

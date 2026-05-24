"""Tests for assemble_pptx.py — slide_spec.json → .pptx rendering.

Coverage targets:
- Round-trip: example_slide_spec → 15 slides on output.
- Each layout: handler runs without error; placeholders filled.
- Figure path resolution: relative paths resolve against draft_dir.
- Missing figure: warning collected, slide still renders.
- Speaker notes: written to slide.notes_slide.
- Validator pre-flight: invalid spec → AssemblyError before any slide writes.
- Master template: defaults to the shipped one; can be overridden.
- CLI surface: validate, --out, --strict, --format pdf fallback.
"""
from __future__ import annotations

import importlib.util
import json
import struct
import sys
import zlib
from pathlib import Path

import pytest


REPO_ROOT = Path(__file__).resolve().parents[2]
SLIDE_SPEC_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
                 / "tools" / "slide_spec.py")
ASSEMBLE_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
               / "tools" / "assemble_pptx.py")
MASTER_PPTX = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
               / "references" / "templates" / "kbase-presentation-master.pptx")


def _import(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def ss():
    return _import("slide_spec", SLIDE_SPEC_PY)


@pytest.fixture(scope="module")
def asm():
    # Ensure slide_spec is loaded first (assemble_pptx imports it via
    # its own importlib).
    _import("slide_spec", SLIDE_SPEC_PY)
    return _import("assemble_pptx", ASSEMBLE_PY)


# ---------------------------------------------------------------------------
# v0.3.2.1: figure path resolver under v0.3.1+ layout
# ---------------------------------------------------------------------------

def test_derive_actual_draft_dir_strips_working(asm, tmp_path):
    """v0.3.1+: caller passes draft_N/working/, helper returns draft_N/."""
    draft = tmp_path / "talks" / "draft_5"
    working = draft / "working"
    working.mkdir(parents=True)
    assert asm._derive_actual_draft_dir(working) == draft


def test_derive_actual_draft_dir_passthrough_for_legacy(asm, tmp_path):
    """v0.3.0 legacy: caller passes draft_N/ directly, no transform."""
    draft = tmp_path / "talks" / "draft_5"
    draft.mkdir(parents=True)
    assert asm._derive_actual_draft_dir(draft) == draft


def test_derive_project_dir_v031_layout(asm, tmp_path):
    """v0.3.1: walk draft_N/working/ → projects/<id>/."""
    project = tmp_path / "projects" / "demo"
    working = project / "talks" / "draft_5" / "working"
    working.mkdir(parents=True)
    assert asm._derive_project_dir(working) == project


def test_derive_project_dir_v030_legacy(asm, tmp_path):
    """v0.3.0 legacy: walk draft_N/ → projects/<id>/ still works."""
    project = tmp_path / "projects" / "demo"
    draft = project / "talks" / "draft_5"
    draft.mkdir(parents=True)
    assert asm._derive_project_dir(draft) == project


def test_resolve_asset_path_finds_project_figure_under_v031(asm, tmp_path):
    """End-to-end: figure path 'figures/X.png' in spec resolves to
    project_dir/figures/X.png even when caller passed draft_N/working/.
    This is the v0.3.2.1 fix that closes the smoke-test bug where 3
    figure assets failed to render."""
    project = tmp_path / "projects" / "demo"
    figdir = project / "figures"
    figdir.mkdir(parents=True)
    (figdir / "fig1.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 32)

    # v0.3.1 layout: draft_dir is the actual draft_N/, but in the assembler
    # main path slide_spec_path.parent = draft_N/working/. Verify the
    # transform.
    working = project / "talks" / "draft_5" / "working"
    working.mkdir(parents=True)

    actual_draft = asm._derive_actual_draft_dir(working)
    warnings: list[str] = []
    found = asm._resolve_asset_path(
        "figures/fig1.png", actual_draft, warnings, "test")
    assert found == figdir / "fig1.png"
    assert warnings == []


requires_master = pytest.mark.skipif(
    not MASTER_PPTX.is_file(),
    reason=f"Master not built; run build_master.py first ({MASTER_PPTX})",
)


def _make_tiny_png(path: Path, size: int = 16) -> Path:
    """Write a minimal valid PNG (single-color square) for figure-asset tests."""
    # 8-byte signature + IHDR + IDAT + IEND
    def chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR",
                 struct.pack(">IIBBBBB", size, size, 8, 2, 0, 0, 0))  # 8-bit RGB
    raw = b"".join(b"\x00" + b"\xa0\xc0\xe0" * size for _ in range(size))
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    path.write_bytes(sig + ihdr + idat + iend)
    return path


# ---------------------------------------------------------------------------
# Smoke-level: example_slide_spec → assembled .pptx
# ---------------------------------------------------------------------------

@requires_master
def test_assemble_example_spec_smoke(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))

    out = tmp_path / "slides.pptx"
    result = asm.assemble(spec_path, out)

    from pptx import Presentation
    assert result.n_slides == 16  # v0.3.2: + data_table
    assert out.is_file()
    prs = Presentation(out)
    assert len(prs.slides) == 16  # v0.3.2: + data_table
    layouts_used = {s.slide_layout.name for s in prs.slides}
    # data_table is aliased to data_figure's master layout — so the master-
    # layout names cover 15 of the 16 spec layouts (data_table reuses
    # data_figure's master).
    expected_master_layouts = set(ss.LAYOUTS) - {"data_table"}
    assert layouts_used == expected_master_layouts


@requires_master
def test_assemble_returns_warnings_for_missing_figures(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    result = asm.assemble(spec_path, out)
    # data_figure and concept_illustration's figures aren't present in the
    # example spec → 2 missing-asset warnings. (Earlier expected ≥3 — that
    # was implicitly counting the workflow_diagram NoneType crash dressed
    # up as a warning. The 2026-04-26 fix to assemble_pptx.py registers
    # diagram_render in sys.modules before exec_module, so the diagram
    # now renders cleanly and emits no warning.)
    assert len(result.warnings) >= 2
    warning_text = " | ".join(result.warnings)
    assert "data_figure" in warning_text
    assert "concept_illustration" in warning_text
    # The diagram render bug from before should NOT appear
    assert "diagram render failed" not in warning_text
    # Slide still renders
    assert out.is_file()


@requires_master
def test_strict_mode_raises_on_warning(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    with pytest.raises(asm.AssemblyError):
        asm.assemble(spec_path, out, strict=True)


# ---------------------------------------------------------------------------
# Validator pre-flight — refuse to render bad spec
# ---------------------------------------------------------------------------

@requires_master
def test_invalid_spec_rejected_before_render(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec["mode"] = "talk-90"   # invalid
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    with pytest.raises(asm.AssemblyError) as exc:
        asm.assemble(spec_path, out)
    assert "schema validation" in str(exc.value)
    assert not out.exists()   # nothing written on failure


def test_missing_slide_spec_file_raises(asm, tmp_path):
    with pytest.raises(asm.AssemblyError):
        asm.assemble(tmp_path / "nope.json", tmp_path / "out.pptx")


# ---------------------------------------------------------------------------
# Figure path resolution
# ---------------------------------------------------------------------------

@requires_master
def test_relative_figure_path_resolves_against_draft_dir(ss, asm, tmp_path):
    # Place the PNG at <tmp>/figures/fig01.png
    fig_dir = tmp_path / "figures"
    fig_dir.mkdir(exist_ok=True)
    _make_tiny_png(fig_dir / "fig01.png")

    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [
            ss.example_slide("data_figure", slide_id=1, substory_id=None),
        ],
    }
    spec["slides"][0]["content"]["figure"] = "figures/fig01.png"

    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    result = asm.assemble(spec_path, out)
    # Relative path must have resolved successfully — no asset warning
    asset_warnings = [w for w in result.warnings if "asset not found" in w]
    assert asset_warnings == []


@requires_master
def test_absolute_figure_path_works(ss, asm, tmp_path):
    fig = _make_tiny_png(tmp_path / "fig.png")

    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide("data_figure", slide_id=1, substory_id=None)],
    }
    spec["slides"][0]["content"]["figure"] = str(fig.resolve())

    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    result = asm.assemble(spec_path, out)
    asset_warnings = [w for w in result.warnings if "asset not found" in w]
    assert asset_warnings == []


# ---------------------------------------------------------------------------
# v0.3.5: data_figure caption shrink-to-fit (belt-and-suspenders for the
# slide_spec validator's 280-char cap; protects the y=5.00 brand strip
# against any caption that slips through).
# ---------------------------------------------------------------------------

@requires_master
def test_data_figure_caption_textbox_shrink_to_fit(ss, asm, tmp_path):
    """v0.3.5 regression: the data_figure caption textbox MUST have
    auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE so the FONT shrinks to
    keep text inside the 0.65-in box at y=4.18..4.83. Without this,
    long captions (e.g., the 410-char gene_function_ecological_agora
    draft_1 slide-21 case) overflow visually into the y=5.00 brand
    strip even though the box is fixed-size and word_wrap is on
    (python-pptx renders text outside box bounds when no auto_size is
    set). The slide_spec validator's 280-char cap is the primary
    defense; this is the third layer for bypassed/old/edge-case specs."""
    from pptx.enum.text import MSO_AUTO_SIZE

    fig = _make_tiny_png(tmp_path / "fig.png")
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide("data_figure", slide_id=1, substory_id=None)],
    }
    spec["slides"][0]["content"]["figure"] = str(fig.resolve())
    spec["slides"][0]["content"]["caption"] = (
        "Caption that lives well inside the 280-char cap; this test only "
        "pins the auto_size setting on the textbox, not the cap itself."
    )

    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    asm.assemble(spec_path, out)

    from pptx import Presentation
    prs = Presentation(out)
    assert len(prs.slides) == 1
    rendered = prs.slides[0]

    # Find the freeform textbox at y≈4.18in (the caption — distinguishes
    # it from the title placeholder at the top, the data_source band at
    # y≈4.83, and the figure picture).
    from pptx.util import Emu
    target_top_emu = Emu.from_inches(4.18) if hasattr(Emu, "from_inches") else int(4.18 * 914400)
    captions = []
    for shape in rendered.shapes:
        if not shape.has_text_frame:
            continue
        # Filter for shapes positioned in the caption band.
        if shape.top is None:
            continue
        # 4.18 in ≈ 3,824,352 EMU. Allow ±0.05 in (45,720 EMU) tolerance.
        if abs(int(shape.top) - target_top_emu) <= 45_720:
            captions.append(shape)

    assert captions, (
        "expected at least one freeform textbox at y≈4.18in (caption band); "
        f"got top values: {[shape.top for shape in rendered.shapes if shape.has_text_frame]}"
    )
    # The caption box specifically — text contains 'Caption that lives'.
    cap_box = next((s for s in captions if "Caption that lives" in s.text_frame.text), None)
    assert cap_box is not None, "caption box not found by text content"
    assert cap_box.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, (
        f"data_figure caption must use TEXT_TO_FIT_SHAPE so font shrinks to "
        f"keep text out of the brand strip at y=5.00; got {cap_box.text_frame.auto_size!r}"
    )


def test_add_textbox_shrink_to_fit_unit(asm, tmp_path):
    """Unit-level: _add_textbox(shrink_to_fit=True) sets the right
    MSO_AUTO_SIZE on the resulting text frame. Avoids needing the
    master template (handler-level), so the contract gets a non-skipped
    test even in environments without the .pptx fixture."""
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)

    asm._add_textbox(slide, "x", 1.0, 1.0, 4.0, 0.5,
                     word_wrap=True, shrink_to_fit=True)
    asm._add_textbox(slide, "y", 1.0, 2.0, 4.0, 0.5,
                     word_wrap=True, auto_size=True)
    asm._add_textbox(slide, "z", 1.0, 3.0, 4.0, 0.5)

    # Slides have textbox shapes in insertion order; pull text frames.
    tfs = [s.text_frame for s in slide.shapes if s.has_text_frame
           and s.text_frame.text in ("x", "y", "z")]
    by_text = {tf.text: tf for tf in tfs}
    assert by_text["x"].auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert by_text["y"].auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # Default for python-pptx-created textboxes IS SHAPE_TO_FIT_TEXT
    # (the library inserts <a:spAutoFit/> by default). Pinning this
    # surfaces the surprising default — without an explicit
    # shrink_to_fit=True call, the caption box would silently grow
    # instead of shrinking the font, exactly the v0.3.2.8 / v0.3.5
    # failure mode. If python-pptx changes this default, this test
    # fails loudly and forces a re-evaluation of the caption render.
    assert by_text["z"].auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def test_add_textbox_shrink_to_fit_wins_over_auto_size(asm, tmp_path):
    """If both auto_size and shrink_to_fit are passed, shrink_to_fit
    takes precedence — that's the documented contract in the docstring."""
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    asm._add_textbox(slide, "both", 1.0, 1.0, 4.0, 0.5,
                     word_wrap=True, auto_size=True, shrink_to_fit=True)
    tf = next(s.text_frame for s in slide.shapes
              if s.has_text_frame and s.text_frame.text == "both")
    assert tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# ---------------------------------------------------------------------------
# Per-layout placeholder filling
# ---------------------------------------------------------------------------

@requires_master
@pytest.mark.parametrize("layout", [
    "title", "section_divider", "big_idea", "big_number",
    "claim_evidence", "two_column_compare", "data_figure",
    "workflow_diagram", "methods_summary", "concept_illustration",
    "cross_tenant_integration", "implications", "acknowledgments",
    "references", "qa_anticipated",
])
def test_each_layout_renders(ss, asm, tmp_path, layout):
    """Build a one-slide spec per layout. Assert it renders without raising
    and that a non-empty title is set on the output slide."""
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide(layout, slide_id=1, substory_id=None)],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    result = asm.assemble(spec_path, out)
    assert result.n_slides == 1

    from pptx import Presentation
    prs = Presentation(out)
    assert len(prs.slides) == 1
    rendered = prs.slides[0]
    assert rendered.slide_layout.name == layout
    # Title placeholder always set (acknowledgments/references hard-code it,
    # the rest pull from content).
    if rendered.shapes.title:
        assert rendered.shapes.title.text != ""


@requires_master
def test_slide_level_title_autofit_landed(ss, asm, tmp_path):
    """2026-04-26 fix #63: slide-level normAutofit must be set on title
    placeholders (layout-level autofit alone doesn't trigger at render).
    Verify every slide except big_number/big_idea has slide-level
    normAutofit fontScale=80000 + anchor=t on its title placeholder.
    """
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    asm.assemble(spec_path, out)

    from pptx import Presentation
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    prs = Presentation(out)
    failures = []
    intentional_no_autofit = {"big_number", "big_idea"}

    for i, slide in enumerate(prs.slides, 1):
        if slide.slide_layout.name in intentional_no_autofit:
            continue
        if not slide.shapes.title:
            continue
        title = slide.shapes.title
        body_pr = title.element.find(f"{{{P_NS}}}txBody/{{{A_NS}}}bodyPr")
        if body_pr is None:
            failures.append(
                f"slide {i} ({slide.slide_layout.name}): title has no bodyPr"
            )
            continue
        norm = body_pr.find(f"{{{A_NS}}}normAutofit")
        if norm is None:
            failures.append(
                f"slide {i} ({slide.slide_layout.name}): "
                f"missing slide-level <a:normAutofit/>"
            )
            continue
        if norm.get("fontScale") != "80000":
            failures.append(
                f"slide {i} ({slide.slide_layout.name}): "
                f"fontScale={norm.get('fontScale')!r}, expected '80000'"
            )
        if body_pr.get("anchor") != "t":
            failures.append(
                f"slide {i} ({slide.slide_layout.name}): "
                f"anchor={body_pr.get('anchor')!r}, expected 't'"
            )

    assert not failures, "\n  ".join(failures)


@requires_master
def test_methods_summary_renders_tools_versions(ss, asm, tmp_path):
    """2026-04-26 fix #59: tools_versions now renders as a footer band
    (was silently dropped before). Build a methods_summary slide with
    real tool/version pairs and assert the formatted footer text appears
    on the rendered slide.
    """
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [{
            "id": 1,
            "layout": "methods_summary",
            "content": {
                "title": "Methods overview",
                "bullets": [
                    "Quality-trim with fastp at Q20",
                    "Annotate with RAST default parameters",
                    "Cross-validate against Morgan Price gold standard",
                    "FDR correction via Benjamini-Hochberg",
                    "Recovery rate computed across n=142 loci",
                ],
                "tools_versions": [
                    {"tool": "RAST", "version": "2.0"},
                    {"tool": "fastp", "version": "0.23"},
                    {"tool": "DRAM", "version": "1.4"},
                ],
            },
        }],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    asm.assemble(spec_path, out)

    from pptx import Presentation
    prs = Presentation(out)
    slide = prs.slides[0]

    # Walk all text-bearing shapes; check for the tools_versions footer
    all_text = " | ".join(
        shape.text_frame.text for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "RAST 2.0" in all_text, f"RAST version missing; got: {all_text}"
    assert "fastp 0.23" in all_text, f"fastp version missing; got: {all_text}"
    assert "DRAM 1.4" in all_text, f"DRAM version missing; got: {all_text}"
    # The "see speaker notes" fallback should NOT appear when tools_versions present
    assert "see speaker notes" not in all_text, (
        "fallback footer should be suppressed when tools_versions populated"
    )


@requires_master
def test_methods_summary_falls_back_to_speaker_notes_hint(ss, asm, tmp_path):
    """When tools_versions is absent, the see-notes footer renders."""
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [{
            "id": 1,
            "layout": "methods_summary",
            "content": {
                "title": "Methods overview",
                "bullets": ["b1", "b2", "b3", "b4", "b5"],
                # No tools_versions
            },
        }],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    asm.assemble(spec_path, out)
    from pptx import Presentation
    slide = Presentation(out).slides[0]
    all_text = " | ".join(
        shape.text_frame.text for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "see speaker notes" in all_text


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------

@requires_master
def test_speaker_notes_written_to_notes_slide(ss, asm, tmp_path):
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide("claim_evidence", slide_id=1, substory_id=None)],
    }
    notes_text = "Speaker notes for slide 1: Smith 2023 showed 90% accuracy."
    spec["slides"][0]["speaker_notes"] = notes_text

    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    asm.assemble(spec_path, out)

    from pptx import Presentation
    prs = Presentation(out)
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert notes_text in notes


# ---------------------------------------------------------------------------
# Title-rule exemptions: acknowledgments + references hard-code their titles
# ---------------------------------------------------------------------------

@requires_master
def test_acknowledgments_title_hardcoded(ss, asm, tmp_path):
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide("acknowledgments", slide_id=1, substory_id=None)],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    asm.assemble(spec_path, out)

    from pptx import Presentation
    prs = Presentation(out)
    assert prs.slides[0].shapes.title.text == "Acknowledgments"


@requires_master
def test_references_title_hardcoded(ss, asm, tmp_path):
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [ss.example_slide("references", slide_id=1, substory_id=None)],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    asm.assemble(spec_path, out)

    from pptx import Presentation
    prs = Presentation(out)
    assert prs.slides[0].shapes.title.text == "References"


# ---------------------------------------------------------------------------
# Master override
# ---------------------------------------------------------------------------

@requires_master
def test_master_path_override(ss, asm, tmp_path):
    """Custom master_path is honored. Use the shipped master copied to a
    side path so we know the override is active."""
    custom = tmp_path / "custom_master.pptx"
    custom.write_bytes(MASTER_PPTX.read_bytes())

    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    result = asm.assemble(spec_path, out, master_path=custom)
    assert result.n_slides == 16  # v0.3.2: + data_table


def test_default_master_path_resolves(asm):
    p = asm.default_master_path()
    assert p.is_file()
    assert p.name == "kbase-presentation-master.pptx"


# ---------------------------------------------------------------------------
# CLI surface
# ---------------------------------------------------------------------------

@requires_master
def test_cli_assemble_clean(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    rc = asm.main([str(spec_path), "--out", str(out)])
    assert rc == 0
    assert out.is_file()


@requires_master
def test_cli_assemble_invalid_spec_returns_2(ss, asm, tmp_path):
    spec = ss.example_slide_spec()
    spec["mode"] = "talk-90"
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    rc = asm.main([str(spec_path), "--out", str(out)])
    assert rc == 2


@requires_master
def test_cli_pdf_fallback_when_no_libreoffice(ss, asm, tmp_path, monkeypatch):
    """If soffice not on PATH, --format pdf must emit pptx + clear message,
    not crash."""
    spec = ss.example_slide_spec()
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"

    # Force render_pdf to return None (simulates missing soffice)
    monkeypatch.setattr(asm, "render_pdf", lambda p: None)
    rc = asm.main([str(spec_path), "--out", str(out), "--format", "pdf"])
    assert rc == 0
    assert out.is_file()  # pptx still emitted


# ---------------------------------------------------------------------------
# Empty / minimal handling
# ---------------------------------------------------------------------------

@requires_master
def test_layout_handlers_dispatch_covers_full_vocabulary(ss, asm):
    """Every named layout has a handler registered. Pins the dispatch
    table so a future layout addition can't slip through unnoticed."""
    assert set(asm.LAYOUT_HANDLERS.keys()) == set(ss.LAYOUTS)


def test_graphite_gray_is_slate_dark(asm):
    """M4a Tier E round 2 (2026-05-23): GRAPHITE_GRAY_RGB is the
    secondary-text color used by step_captions, source_footers,
    citation bands, AI-disclosure footers, table data cells, etc.
    The pre-round-2 value (157, 146, 135) — a light tan-gray drawn
    from the KBase brand-token `graphite_gray` — washed out against
    the master's cream background. Round 2 bumped it to slate-dark
    (80, 75, 70) — a warm-tinted near-black that holds contrast on
    cream + post-watermark-strip flat backgrounds while keeping the
    visual hierarchy between primary (body) and secondary text.

    Pin: the value must be substantially darker than (140, 140, 140).
    """
    r, g, b = asm.GRAPHITE_GRAY_RGB
    # Each channel below 140 = visibly darker than 'gray'
    assert r < 140 and g < 140 and b < 140, (
        f"GRAPHITE_GRAY_RGB={asm.GRAPHITE_GRAY_RGB} — secondary text "
        f"must be substantially darker than the pre-round-2 value "
        f"(157, 146, 135) for contrast against the cream master."
    )


# ---------------------------------------------------------------------------
# M4a Tier A — explicit-fontScale shrink-to-fit + footer-safety geometry
# ---------------------------------------------------------------------------

def test_fontscale_for_chars_full_for_short_content(asm):
    """Content at or below `full_below` renders at 100% — no shrink."""
    assert asm._fontscale_for_chars(50, full_below=100) == asm.FONTSCALE_FULL
    assert asm._fontscale_for_chars(100, full_below=100) == asm.FONTSCALE_FULL


def test_fontscale_for_chars_floor_for_long_content(asm):
    """Content past the longest ladder cap clamps at the 60% floor (DQ3)."""
    # Default ladder: full_below=200, (400, 90), (700, 80), (1100, 70), else floor
    assert asm._fontscale_for_chars(2000) == asm.FONTSCALE_FLOOR
    assert asm.FONTSCALE_FLOOR == 60000  # DQ3 — 60% pinned


def test_fontscale_for_chars_ladder_steps(asm):
    """Adaptive ladder picks the first cap that fits."""
    assert asm._fontscale_for_chars(300) == 90000   # in (200, 400]
    assert asm._fontscale_for_chars(500) == 80000   # in (400, 700]
    assert asm._fontscale_for_chars(900) == 70000   # in (700, 1100]


def test_fit_textbox_appends_warning_when_clamped_at_floor(asm, tmp_path):
    """DQ3: content beyond the longest ladder cap clamps at the floor
    AND appends a soft-warning so the operator + Tier-C visual-QA pass
    both see the slot is at the edge of legibility."""
    from pptx import Presentation
    from pptx.util import Inches as _In
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(_In(0.5), _In(0.5), _In(4.0), _In(0.4))
    tb.text_frame.text = "x" * 1500   # past 1100, hits floor
    warnings = []
    scale = asm._fit_textbox(tb, warnings=warnings, where="synthetic")
    assert scale == asm.FONTSCALE_FLOOR
    assert any("floor" in w and "synthetic" in w for w in warnings), warnings


def test_fit_textbox_no_warning_for_short_content(asm):
    """Content that doesn't reach the floor produces no warning."""
    from pptx import Presentation
    from pptx.util import Inches as _In
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(_In(0.5), _In(0.5), _In(4.0), _In(0.4))
    tb.text_frame.text = "short"
    warnings = []
    scale = asm._fit_textbox(tb, warnings=warnings, where="synthetic")
    assert scale == asm.FONTSCALE_FULL
    assert warnings == []


@requires_master
def test_overflow_prone_slots_carry_explicit_fontscale(ss, asm, tmp_path):
    """M4a Tier A AC: a slot-busting synthetic deck (long content in every
    Tier-A targeted slot) assembles AND the targeted freeform textboxes
    each carry an explicit <a:normAutofit fontScale="..."> on their
    bodyPr — the only shrink-to-fit LibreOffice computes at render. A
    bare <a:normAutofit/> would silently fall back to no-shrink and the
    text would spill (the M3 Tier-E render defect class).

    Slots checked:
      - big_number subtitle + sub_pointer + source_footer
      - workflow_diagram step_caption[i] + tool_version_footer
      - data_table caption + footnote
      - methods_summary tools_versions footer
    """
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _has_explicit_fontscale(shape) -> tuple[bool, str | None]:
        """Return (has_normAutofit_with_explicit_fontScale, scale_str)."""
        tx_body = shape.text_frame._txBody
        body_pr = tx_body.find(f"{{{DML}}}bodyPr")
        if body_pr is None:
            return (False, None)
        norm = body_pr.find(f"{{{DML}}}normAutofit")
        if norm is None:
            return (False, None)
        return (norm.get("fontScale") is not None, norm.get("fontScale"))

    # Slot-busting content: long enough to trigger shrink on every Tier-A
    # textbox. Stay within slide_spec validator caps (data_figure caption
    # 280 chars is the only hard cap; the rest are advisory in Tier B).
    long_subtitle = "fitness scores integrated across 1,400 genomes " \
                    "from the DOE-funded BERDL pipeline running on KBase " \
                    "with Shewanella as the model organism (270 chars to busy autofit)"
    long_sub_pointer = "Top decile by ensemble score across three independent " \
                       "ML predictors plus phylogenetic conservation gates (180+ chars)"
    long_source_footer = ("REPORT.md §4.2 cited from primary sources "
                          "Smith2023, Jones2024, Lee2025; data DOI 10.5281/zenodo.example "
                          "(stretched to push the source_footer slot past its full_below cap)")
    long_step_caption = ("Run the workflow under controlled BERDL parameters "
                         "with version pins from tools.lock to ensure reproducibility "
                         "across the three pipeline steps")
    long_tool_version_footer = (
        "RAST 2.0.1 · fastp 0.23.4 · DRAM 1.4.6 · GTDB-Tk 2.3.0 · "
        "checkm2 1.0.2 · diamond 2.1.8 · spades 3.15.5 · prodigal 2.6.3 (push past 200 chars)"
    )
    long_data_table_caption = ("Top-decile candidates ranked by ensemble score; "
                               "ML+conservation+phenotype-gated subset from the "
                               "n=347 candidate pool (REPORT.md §4.2 full ranking, "
                               "with cross-validation on a 120-genome holdout — past 280 chars to test floor)")
    long_data_table_footnote = ("Full ranking (n=347) in REPORT.md §4.2; scores "
                                "are the geometric mean of three predictors; "
                                "phenotype evidence sourced from Shewanella growth panels (>240 chars)")

    slides = [
        {
            "id": 1, "substory_id": None, "layout": "big_number",
            "content": {
                "headline": "27,000,000",
                "subtitle": long_subtitle,
                "sub_pointer": long_sub_pointer,
                "source_footer": long_source_footer,
            },
        },
        {
            "id": 2, "substory_id": None, "layout": "workflow_diagram",
            "content": {
                "title": "Slot-busting workflow.",
                "diagram": {
                    "kind": "boxes_and_arrows",
                    "nodes": [
                        {"id": "n1", "label": "Long node label that overruns its box badly",
                         "shape": "rounded", "x": 0.5, "y": 1.4, "w": 1.5, "h": 0.8},
                        {"id": "n2", "label": "End", "shape": "rounded",
                         "x": 7.0, "y": 1.4, "w": 1.5, "h": 0.8},
                    ],
                    "edges": [{"from": "n1", "to": "n2", "kind": "straight",
                               "label": "step"}],
                },
                "step_caption": [long_step_caption, long_step_caption, long_step_caption],
                "tool_version_footer": long_tool_version_footer,
            },
        },
        {
            "id": 3, "substory_id": None, "layout": "data_table",
            "content": {
                "title": "Slot-busting table.",
                "columns": ["A", "B", "C"],
                "rows": [["a1", "b1", "c1"], ["a2", "b2", "c2"]],
                "caption": long_data_table_caption,
                "footnote": long_data_table_footnote,
            },
        },
        {
            "id": 4, "substory_id": None, "layout": "methods_summary",
            "content": {
                "title": "Slot-busting methods.",
                "bullets": ["b1", "b2", "b3", "b4", "b5"],
                "tools_versions": [
                    {"tool": f"tool_{i}", "version": f"{i}.0.0"}
                    for i in range(8)
                ],
            },
        },
    ]
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": slides,
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    result = asm.assemble(spec_path, out)
    assert result.n_slides == 4

    from pptx import Presentation
    prs = Presentation(out)
    failures = []

    # --- Slide 1 (big_number): subtitle + sub_pointer + source_footer ---
    bn = prs.slides[0]
    # The three textboxes are added freeform AFTER the title placeholder.
    # Walk shapes; the textboxes are the ones whose text matches our content.
    bn_subtitle = bn_sub_pointer = bn_source_footer = None
    for shp in bn.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text
        if txt.startswith("fitness scores"):
            bn_subtitle = shp
        elif txt.startswith("Top decile"):
            bn_sub_pointer = shp
        elif txt.startswith("REPORT.md §4.2 cited"):
            bn_source_footer = shp
    for name, shp in [("big_number subtitle", bn_subtitle),
                       ("big_number sub_pointer", bn_sub_pointer),
                       ("big_number source_footer", bn_source_footer)]:
        assert shp is not None, f"{name}: textbox not found in rendered slide"
        ok, scale = _has_explicit_fontscale(shp)
        if not ok:
            failures.append(f"{name}: no explicit fontScale on normAutofit "
                            "(LibreOffice will not shrink the text)")

    # --- Slide 2 (workflow_diagram): step_captions + tool_version_footer ---
    wf = prs.slides[1]
    wf_step_captions = []
    wf_tvf = None
    for shp in wf.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text
        if txt.startswith("Run the workflow"):
            wf_step_captions.append(shp)
        elif txt.startswith("RAST 2.0.1"):
            wf_tvf = shp
    assert len(wf_step_captions) == 3, \
        f"expected 3 step_caption textboxes, got {len(wf_step_captions)}"
    for i, shp in enumerate(wf_step_captions):
        ok, scale = _has_explicit_fontscale(shp)
        if not ok:
            failures.append(f"workflow_diagram step_caption[{i}]: "
                            "no explicit fontScale on normAutofit")
    assert wf_tvf is not None, "workflow_diagram tool_version_footer: not found"
    ok, _ = _has_explicit_fontscale(wf_tvf)
    if not ok:
        failures.append("workflow_diagram tool_version_footer: "
                        "no explicit fontScale on normAutofit")

    # --- Slide 3 (data_table): caption + footnote ---
    dt = prs.slides[2]
    dt_caption = dt_footnote = None
    for shp in dt.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text
        if txt.startswith("Top-decile candidates"):
            dt_caption = shp
        elif txt.startswith("Full ranking"):
            dt_footnote = shp
    for name, shp in [("data_table caption", dt_caption),
                       ("data_table footnote", dt_footnote)]:
        assert shp is not None, f"{name}: textbox not found"
        ok, _ = _has_explicit_fontscale(shp)
        if not ok:
            failures.append(f"{name}: no explicit fontScale on normAutofit")

    # --- Slide 4 (methods_summary): tools_versions footer ---
    ms = prs.slides[3]
    ms_tvf = None
    for shp in ms.shapes:
        if not shp.has_text_frame:
            continue
        if "tool_0 0.0.0" in shp.text_frame.text:
            ms_tvf = shp
            break
    assert ms_tvf is not None, "methods_summary tools_versions: textbox not found"
    ok, _ = _has_explicit_fontscale(ms_tvf)
    if not ok:
        failures.append("methods_summary tools_versions: "
                        "no explicit fontScale on normAutofit")

    assert not failures, "\n  ".join([""] + failures)


@requires_master
def test_advisory_soft_warning_flows_to_assembly_warnings(ss, asm, tmp_path):
    """M4a Tier B (DQ4): a long big_number subtitle emits a soft-warning
    validator issue; the assembler does NOT raise (renderer Tier A
    absorbs the overflow) but surfaces the warning through
    AssemblyResult.warnings so the operator + Tier-C visual-QA pass see
    it. Mirrors how _fit_textbox clamp-warnings flow."""
    long_subtitle = "x" * (ss.BIG_NUMBER_SUBTITLE_MAX_CHARS + 50)
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [{
            "id": 1, "substory_id": None, "layout": "big_number",
            "content": {"headline": "42", "subtitle": long_subtitle},
        }],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    result = asm.assemble(spec_path, out)  # must not raise
    assert result.n_slides == 1
    assert any("subtitle" in w and "advisory cap" in w
               for w in result.warnings), \
        f"expected soft-warning surfaced in AssemblyResult.warnings: {result.warnings}"


@requires_master
def test_strict_mode_treats_soft_warnings_as_failures(ss, asm, tmp_path):
    """--strict is the explicit opt-in to fail on any warning. Soft-
    warnings flow through .warnings, so --strict raises on them too —
    by design, the same way it raises on missing-figure warnings."""
    long_subtitle = "x" * (ss.BIG_NUMBER_SUBTITLE_MAX_CHARS + 50)
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": [{
            "id": 1, "substory_id": None, "layout": "big_number",
            "content": {"headline": "42", "subtitle": long_subtitle},
        }],
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    with pytest.raises(asm.AssemblyError, match="strict"):
        asm.assemble(spec_path, out, strict=True)


@requires_master
def test_overflow_prone_slots_geometry_clears_footer(ss, asm, tmp_path):
    """M4a Tier A AC: the same slot-busting deck assembles with all text-
    bearing freeform textboxes ending at or above FOOTER_SAFE_BOTTOM
    (4.92in). M3 E-6 set this constant; M3 E-7 added the width check
    (zero-width placeholders rendered text one char per line). This test
    pins both checks on a committed synthetic deck so regressions surface
    in the suite, not on the next live render.
    """
    # Reuse the same slot-busting deck as the fontScale test
    slides = [
        {
            "id": 1, "substory_id": None, "layout": "big_number",
            "content": {
                "headline": "27,000,000",
                "subtitle": "x" * 250,
                "sub_pointer": "y" * 180,
                "source_footer": "z" * 200,
            },
        },
        {
            "id": 2, "substory_id": None, "layout": "workflow_diagram",
            "content": {
                "title": "wf",
                "diagram": {
                    "kind": "boxes_and_arrows",
                    "nodes": [
                        {"id": "n1", "label": "n1", "shape": "rounded",
                         "x": 0.5, "y": 1.4, "w": 1.5, "h": 0.8},
                        {"id": "n2", "label": "n2", "shape": "rounded",
                         "x": 7.0, "y": 1.4, "w": 1.5, "h": 0.8},
                    ],
                    "edges": [{"from": "n1", "to": "n2", "kind": "straight"}],
                },
                "step_caption": ["s" * 150] * 3,
                "tool_version_footer": "t" * 200,
            },
        },
        {
            "id": 3, "substory_id": None, "layout": "data_table",
            "content": {
                "title": "dt",
                "columns": ["A", "B"],
                "rows": [["a", "b"]],
                "caption": "c" * 280,
                "footnote": "f" * 200,
            },
        },
        {
            "id": 4, "substory_id": None, "layout": "methods_summary",
            "content": {
                "title": "ms",
                "bullets": ["b1", "b2", "b3", "b4", "b5"],
                "tools_versions": [
                    {"tool": "t", "version": "1.0"} for _ in range(8)
                ],
            },
        },
        {
            "id": 5, "substory_id": None, "layout": "qa_anticipated",
            "content": {
                "question": "q",
                "answer_summary": "a" * 800,
                "answer_detail": "d" * 200,
                "evidence_pointer": "Substory 1",
            },
        },
    ]
    spec = {
        "schema_version": ss.SCHEMA_VERSION,
        "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "throughline": {"id": "TL1", "punchline": "x", "tier_evidence": "STRONG"},
        "substories": [],
        "slides": slides,
    }
    spec_path = tmp_path / "slide_spec.json"
    spec_path.write_text(json.dumps(spec))
    out = tmp_path / "slides.pptx"
    asm.assemble(spec_path, out)

    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation(out)
    failures = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            # Skip pictures, tables, connectors, the slide title (logos /
            # decorative banners are baked into the master and can sit
            # in the footer band by design — they're the FOOTER itself).
            if not shp.has_text_frame:
                continue
            # Skip the slide title placeholder
            if shp == slide.shapes.title:
                continue
            # Width must be > 0 (E-7 regression: placeholders with only
            # top/height set zero-widthed text into one-char-per-line).
            width_in = (shp.width or 0) / 914400  # EMU → in
            height_in = (shp.height or 0) / 914400
            top_in = (shp.top or 0) / 914400
            left_in = (shp.left or 0) / 914400
            if width_in <= 0.1:
                failures.append(
                    f"slide {slide_idx} ({slide.slide_layout.name}): "
                    f"text shape '{shp.text_frame.text[:40]}...' has "
                    f"width={width_in:.2f}in (E-7 zero-width regression)"
                )
            # Bottom must clear FOOTER_SAFE_BOTTOM = 4.92in
            bottom_in = top_in + height_in
            if bottom_in > asm.FOOTER_SAFE_BOTTOM + 0.01:  # 0.01 tolerance
                failures.append(
                    f"slide {slide_idx} ({slide.slide_layout.name}): "
                    f"text shape '{shp.text_frame.text[:40]}...' bottom "
                    f"{bottom_in:.2f}in > FOOTER_SAFE_BOTTOM "
                    f"{asm.FOOTER_SAFE_BOTTOM}in (logo strip starts ~5.00)"
                )
            # Negative or off-slide positions
            if left_in < -0.01 or top_in < -0.01:
                failures.append(
                    f"slide {slide_idx} ({slide.slide_layout.name}): "
                    f"text shape '{shp.text_frame.text[:40]}...' negative "
                    f"position (left={left_in:.2f}, top={top_in:.2f})"
                )

    assert not failures, "\n  ".join([""] + failures)

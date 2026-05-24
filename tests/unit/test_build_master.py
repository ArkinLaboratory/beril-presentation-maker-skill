"""Tests for tools/build_master.py — derived master template integrity.

These tests verify the v0.1.0-master-draft phase output:

- Build is reproducible (idempotency).
- 15 named layouts present in the derived master.
- Exactly 1 slide master in the output (no Google Slides duplicate).
- Brand-tokens JSON has the expected schema and the canonical Style
  Guide June 2022 values.
- Every named layout has at least one TITLE placeholder.
- Round-trip: a deck assembled using each layout opens cleanly.

Source .potx is gitignored (user-supplied). Tests skip with a clear
message when source is absent.
"""
from __future__ import annotations

import importlib.util
import json
from pathlib import Path

import pytest


REPO_ROOT = Path(__file__).resolve().parents[2]
SOURCE_POTX = REPO_ROOT / "reference" / "master-template-source" / "KBase 2026 and beyond.potx"
DEST_PPTX = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
             / "references" / "templates" / "kbase-presentation-master.pptx")
BRAND_JSON = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
              / "references" / "kbase-brand-tokens.json")
BUILD_SCRIPT = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
                / "tools" / "build_master.py")


def _import_build_master():
    """Import build_master.py as a module (it lives outside the package
    Python namespace so we load it by file path)."""
    spec = importlib.util.spec_from_file_location("build_master", BUILD_SCRIPT)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load build_master from {BUILD_SCRIPT}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


# ---------------------------------------------------------------------------
# Fixtures + skip-marks
# ---------------------------------------------------------------------------

requires_source = pytest.mark.skipif(
    not SOURCE_POTX.is_file(),
    reason=f"Source .potx absent at {SOURCE_POTX} — gitignored, user-supplied",
)

requires_built_master = pytest.mark.skipif(
    not DEST_PPTX.is_file(),
    reason=f"Built master absent at {DEST_PPTX} — run build_master.py first",
)


# ---------------------------------------------------------------------------
# Constant tests (no .potx required)
# ---------------------------------------------------------------------------

def test_named_vocabulary_is_15_unique():
    bm = _import_build_master()
    assert len(bm.NAMED_VOCABULARY) == 15
    assert len(set(bm.NAMED_VOCABULARY)) == 15


def test_layout_renames_keys_match_vocabulary_size():
    bm = _import_build_master()
    assert len(bm.LAYOUT_RENAMES) == 15
    assert set(bm.LAYOUT_RENAMES.values()) == set(bm.NAMED_VOCABULARY)


def test_brand_tokens_constant_has_expected_keys():
    bm = _import_build_master()
    bt = bm.BRAND_TOKENS
    assert bt["version"] == "1.0"
    assert "KBase Style Guide" in bt["source"] or "Style Guidelines" in bt["source"]
    assert set(bt["palette"]["primary"].keys()) == {
        "microbe_orange", "grass_green", "freshwater_blue",
        "golden_yellow", "spring_green", "ocean_blue",
    }
    assert set(bt["palette"]["secondary"].keys()) == {
        "cyanobacteria_teal", "lupine_purple", "frost_blue",
        "rainier_cherry_red", "graphite_gray",
    }
    # Spot-check a canonical hex from KBase Style Guide June 2022
    assert bt["palette"]["primary"]["microbe_orange"]["hex"] == "#F78E1E"
    assert bt["palette"]["primary"]["freshwater_blue"]["hex"] == "#007DC3"
    # Forbidden contrast pair documented per Style Guide §5
    pairs = [tuple(w["pair"]) for w in bt["palette"]["contrast_warnings"]]
    assert ("spring_green", "golden_yellow") in pairs


def test_brand_tokens_typography_min_sizes():
    """SPEC §6.3 mandates 24-pt body and 36-pt content title minima.
    Brand tokens must reflect those sizes."""
    bm = _import_build_master()
    sizes = bm.BRAND_TOKENS["typography"]["sizes_pt"]
    assert sizes["body"] >= 24
    assert sizes["content_title"] >= 36


# ---------------------------------------------------------------------------
# Built-master tests (require build to have run successfully)
# ---------------------------------------------------------------------------

@requires_built_master
def test_derived_master_has_one_slide_master():
    from pptx import Presentation
    prs = Presentation(DEST_PPTX)
    assert len(prs.slide_masters) == 1, (
        f"derived master should have exactly 1 slide master "
        f"(Google Slides duplicate removed), got {len(prs.slide_masters)}"
    )


@requires_built_master
def test_derived_master_has_15_named_layouts():
    from pptx import Presentation
    bm = _import_build_master()
    prs = Presentation(DEST_PPTX)
    layouts = list(prs.slide_masters[0].slide_layouts)
    found = sorted(l.name for l in layouts)
    expected = sorted(bm.NAMED_VOCABULARY)
    assert found == expected, (
        f"layout name set mismatch:\n"
        f"  found:    {found}\n"
        f"  expected: {expected}"
    )


@requires_built_master
def test_every_layout_has_title_placeholder():
    from pptx import Presentation
    prs = Presentation(DEST_PPTX)
    failures = []
    for layout in prs.slide_masters[0].slide_layouts:
        has_title = any(
            ph.placeholder_format.idx == 0
            for ph in layout.placeholders
        )
        if not has_title:
            failures.append(layout.name)
    assert not failures, f"layouts missing TITLE placeholder: {failures}"


@requires_built_master
def test_derived_master_ships_no_slides():
    """The master ships as a template with 0 slides; the assembler creates
    fresh slides per draft. SPEC §14 / D-007."""
    from pptx import Presentation
    prs = Presentation(DEST_PPTX)
    assert len(prs.slides) == 0, (
        f"derived master should ship empty (0 slides), found {len(prs.slides)}"
    )


@requires_built_master
def test_brand_tokens_json_on_disk_matches_constant():
    bm = _import_build_master()
    on_disk = json.loads(BRAND_JSON.read_text(encoding="utf-8"))
    # JSON-roundtrip the constant for comparison
    expected = json.loads(json.dumps(bm.BRAND_TOKENS))
    assert on_disk == expected


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _find_sp(layout_element, target_spec):
    """Re-implements _find_target_shape's lookup for test verification.
    We don't import the build_master function so the test is independent of
    the production target-resolution logic (and would catch its bugs)."""
    from lxml import etree as _et
    drawable_locals = {"sp", "pic", "grpSp", "graphicFrame"}

    if "by_ph" in target_spec:
        ph_type, ph_idx = target_spec["by_ph"]
        for sp in layout_element.iter(f"{{{P_NS}}}sp"):
            ph = sp.find(f".//{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
            if ph is None:
                continue
            this_type = ph.get("type", "body")
            this_idx = ph.get("idx", "0")
            if ph_type == "title":
                if this_type in ("title", "ctrTitle") or this_idx == "0":
                    return sp
            else:
                if this_type == ph_type and (ph_idx is None or this_idx == ph_idx):
                    return sp
        return None

    if "by_shape_index" in target_spec:
        sptree = layout_element.find(f".//{{{P_NS}}}cSld/{{{P_NS}}}spTree")
        if sptree is None:
            return None
        drawables = [c for c in sptree
                     if _et.QName(c).localname in drawable_locals]
        n = target_spec["by_shape_index"]
        return drawables[n] if n < len(drawables) else None

    if "by_kind_index" in target_spec:
        kind, n = target_spec["by_kind_index"]
        sptree = layout_element.find(f".//{{{P_NS}}}cSld/{{{P_NS}}}spTree")
        if sptree is None:
            return None
        matching = []
        for c in sptree:
            local = _et.QName(c).localname
            if local != kind:
                continue
            if local == "sp":
                # Skip placeholders — by_kind_index targets decorative shapes.
                ph = c.find(f".//{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
                if ph is not None:
                    continue
            matching.append(c)
        return matching[n] if n < len(matching) else None

    return None


def _check_shape_edit_applied(layout_element, shape_edit):
    """Assert that all change attributes in `shape_edit` are present on the
    target shape in `layout_element`. Helper for test_layout_fix_*."""
    sp = _find_sp(layout_element, shape_edit)
    assert sp is not None, f"target shape not found for {shape_edit}"

    if "xfrm" in shape_edit:
        off = sp.find(f".//{{{A_NS}}}xfrm/{{{A_NS}}}off")
        ext = sp.find(f".//{{{A_NS}}}xfrm/{{{A_NS}}}ext")
        assert off is not None and ext is not None, "xfrm/off/ext missing"
        for k, expected in shape_edit["xfrm"].items():
            attr = {"off_x": ("off", "x"), "off_y": ("off", "y"),
                    "ext_cx": ("ext", "cx"), "ext_cy": ("ext", "cy")}[k]
            elem = off if attr[0] == "off" else ext
            actual = int(elem.get(attr[1]))
            assert actual == expected, f"{k}: expected {expected}, got {actual}"

    if "body_pr" in shape_edit:
        body_pr = sp.find(f".//{{{P_NS}}}txBody/{{{A_NS}}}bodyPr")
        assert body_pr is not None, "<a:bodyPr> missing"
        if "anchor" in shape_edit["body_pr"]:
            assert body_pr.get("anchor") == shape_edit["body_pr"]["anchor"]
        if "auto_fit_kind" in shape_edit["body_pr"]:
            kind = shape_edit["body_pr"]["auto_fit_kind"]
            assert body_pr.find(f"{{{A_NS}}}{kind}") is not None
            for other in ("normAutofit", "noAutofit", "spAutoFit"):
                if other != kind:
                    assert body_pr.find(f"{{{A_NS}}}{other}") is None, (
                        f"unexpected <a:{other}> still present"
                    )

    if "lvl1_ppr" in shape_edit:
        lvl1 = sp.find(f".//{{{A_NS}}}lstStyle/{{{A_NS}}}lvl1pPr")
        assert lvl1 is not None
        if "algn" in shape_edit["lvl1_ppr"]:
            assert lvl1.get("algn") == shape_edit["lvl1_ppr"]["algn"]

    if "def_rpr" in shape_edit:
        lvl1 = sp.find(f".//{{{A_NS}}}lstStyle/{{{A_NS}}}lvl1pPr")
        assert lvl1 is not None
        def_rpr = lvl1.find(f"{{{A_NS}}}defRPr")
        assert def_rpr is not None
        if "sz" in shape_edit["def_rpr"]:
            assert def_rpr.get("sz") == str(shape_edit["def_rpr"]["sz"])
        if "b" in shape_edit["def_rpr"]:
            assert def_rpr.get("b") == str(shape_edit["def_rpr"]["b"])


@requires_built_master
@pytest.mark.parametrize("layout_name", [
    "big_number", "big_idea", "section_divider",
    "two_column_compare", "concept_illustration",
])
def test_layout_fix_applied(layout_name):
    """All shape_edits in LAYOUT_FIXES[layout_name] are present on the
    derived master. Pins Adam's 2026-04-26 visual-review fixes so a future
    regenerated master can't drift back to the source .potx defaults.
    """
    from pptx import Presentation
    bm = _import_build_master()

    prs = Presentation(DEST_PPTX)
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
    layout = layouts[layout_name]

    fix = bm.LAYOUT_FIXES[layout_name]
    for shape_edit in fix["shape_edits"]:
        _check_shape_edit_applied(layout.element, shape_edit)


@requires_built_master
def test_title_autofit_universal_sweep():
    """Every layout EXCEPT big_number/big_idea has explicit normAutofit
    with fontScale + lnSpcReduction on its title placeholder, plus
    anchor=t. Regression for the 2026-04-26 visual review where 19/21
    slides had 1.6-5.0x title overrun because PowerPoint doesn't honor
    bare normAutofit at render time.
    """
    from pptx import Presentation
    bm = _import_build_master()
    prs = Presentation(DEST_PPTX)

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    intentional = bm._LAYOUTS_WITH_INTENTIONAL_NO_AUTOFIT_TITLE
    failures = []

    for layout in prs.slide_masters[0].slide_layouts:
        # Find the title placeholder
        title_sp = None
        for sp in layout.element.iter(f"{{{P_NS}}}sp"):
            ph = sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
            if ph is not None and ph.get("type", "body") in ("title", "ctrTitle"):
                title_sp = sp
                break
        if title_sp is None:
            failures.append(f"{layout.name}: no title placeholder")
            continue
        body_pr = title_sp.find(f"{{{P_NS}}}txBody/{{{A_NS}}}bodyPr")
        if body_pr is None:
            failures.append(f"{layout.name}: title has no bodyPr")
            continue

        norm = body_pr.find(f"{{{A_NS}}}normAutofit")
        no_af = body_pr.find(f"{{{A_NS}}}noAutofit")
        sp_af = body_pr.find(f"{{{A_NS}}}spAutoFit")

        if layout.name in intentional:
            # Intentional: must have noAutofit (font size pinned)
            if no_af is None:
                failures.append(
                    f"{layout.name}: intentional-no-autofit layout missing "
                    f"<a:noAutofit/> (got norm={norm is not None}, sp={sp_af is not None})"
                )
            continue

        # Universal: must have normAutofit with explicit fontScale + lnSpcReduction
        if norm is None:
            failures.append(
                f"{layout.name}: missing <a:normAutofit/> "
                f"(got noAutofit={no_af is not None}, spAutoFit={sp_af is not None})"
            )
            continue
        font_scale = norm.get("fontScale")
        ln_spc = norm.get("lnSpcReduction")
        if font_scale is None:
            failures.append(f"{layout.name}: normAutofit lacks fontScale attr")
        if ln_spc is None:
            failures.append(f"{layout.name}: normAutofit lacks lnSpcReduction attr")
        anchor = body_pr.get("anchor")
        if anchor != "t":
            failures.append(
                f"{layout.name}: anchor must be 't' (top — overflow grows "
                f"downward only), got {anchor!r}"
            )

    assert not failures, (
        "title autofit sweep failures:\n  " + "\n  ".join(failures)
    )


@requires_built_master
def test_two_column_compare_has_two_body_placeholders():
    """The two_column_compare layout is the only one with 2 BODY placeholders;
    failure here means we mapped the wrong source layout."""
    from pptx import Presentation
    prs = Presentation(DEST_PPTX)
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
    layout = layouts["two_column_compare"]
    body_count = sum(
        1 for ph in layout.placeholders
        if ph.placeholder_format.type is not None
        and "BODY" in str(ph.placeholder_format.type)
    )
    assert body_count == 2, (
        f"two_column_compare should have exactly 2 BODY placeholders, got {body_count}"
    )


# ---------------------------------------------------------------------------
# Idempotency (requires source .potx — skipped in CI without it)
# ---------------------------------------------------------------------------

@requires_source
def test_build_is_idempotent(tmp_path: Path):
    """Two builds from the same source .potx produce byte-identical zip
    contents (modulo zip-internal timestamps; we compare unzipped XML)."""
    bm = _import_build_master()
    out1 = tmp_path / "build1.pptx"
    out2 = tmp_path / "build2.pptx"
    bt1 = tmp_path / "tokens1.json"
    bt2 = tmp_path / "tokens2.json"

    bm.build_master(SOURCE_POTX, out1, bt1, verbose=False)
    bm.build_master(SOURCE_POTX, out2, bt2, verbose=False)

    import zipfile
    with zipfile.ZipFile(out1) as z1, zipfile.ZipFile(out2) as z2:
        names1 = sorted(z1.namelist())
        names2 = sorted(z2.namelist())
        assert names1 == names2, "zip member set differs across rebuilds"
        # Compare each member's contents
        for name in names1:
            assert z1.read(name) == z2.read(name), (
                f"member '{name}' differs across rebuilds"
            )

    # JSON outputs are also byte-identical
    assert bt1.read_bytes() == bt2.read_bytes()


@requires_source
def test_build_emits_expected_report_keys(tmp_path: Path):
    bm = _import_build_master()
    out = tmp_path / "build.pptx"
    bt = tmp_path / "tokens.json"
    report = bm.build_master(SOURCE_POTX, out, bt, verbose=False)
    assert set(report.keys()) >= {
        "source", "dest_master", "dest_brand_tokens",
        "renamed_layouts", "deleted_layouts", "deleted_masters",
        "deleted_slides_count", "applied_fixes",
        "dest_master_sha256", "dest_master_size_bytes",
    }
    # Must rename exactly 15 layouts and delete the 17 unused ones
    assert len(report["renamed_layouts"]) == 15
    assert len(report["deleted_layouts"]) == 17
    # The Google Slides duplicate master should have been removed
    assert len(report["deleted_masters"]) == 1
    # The big_number layout fix must have been applied
    assert "big_number" in report["applied_fixes"]


@requires_source
def test_build_strips_full_slide_watermarks(tmp_path: Path):
    """M4a Tier E round 2 (2026-05-23): every layout that carried a
    full-slide plant-stem watermark in the source .potx must have it
    removed in the derived master. Tier E round-1 visual-QA found the
    watermark visually competed with body content on 9 of 28 slides
    (workflow_diagram, methods_summary, qa_anticipated, etc.). Strip
    is bounded to full-slide pictures (w >= 8in AND h >= 4.5in); DOE
    + KBase logos (small bottom-right pictures) survive.
    """
    bm = _import_build_master()
    out = tmp_path / "build.pptx"
    bt = tmp_path / "tokens.json"
    report = bm.build_master(SOURCE_POTX, out, bt, verbose=False)
    assert "stripped_watermarks" in report, "report missing stripped_watermarks key"
    # At least 10 layouts should be stripped (we observed 12 in the source
    # .potx; allow some variance if the brand template changes upstream)
    assert len(report["stripped_watermarks"]) >= 10, (
        f"expected ≥10 layouts stripped of full-slide watermarks; "
        f"got {report['stripped_watermarks']}"
    )
    # The four canonical content-bearing layouts MUST be in the stripped list
    stripped_names = {item["layout"] for item in report["stripped_watermarks"]}
    for required in ("workflow_diagram", "methods_summary",
                     "qa_anticipated", "data_figure"):
        assert required in stripped_names, (
            f"layout {required!r} must have its watermark stripped "
            f"(visual-QA round-1 found these were the worst offenders)"
        )


@requires_source
def test_build_preserves_logo_pictures(tmp_path: Path):
    """The watermark-strip uses a full-slide heuristic; the DOE + KBase
    logos (~1.5 x 0.5in, bottom-right) must NOT be removed."""
    from pptx import Presentation
    bm = _import_build_master()
    out = tmp_path / "build.pptx"
    bt = tmp_path / "tokens.json"
    bm.build_master(SOURCE_POTX, out, bt, verbose=False)
    prs = Presentation(out)
    EMU = 914400
    # Sample two layouts known to carry logos: workflow_diagram
    # (watermark stripped) and title (no watermark, has subtitle slot).
    for layout_name in ("workflow_diagram",):
        layout = next((l for l in prs.slide_masters[0].slide_layouts
                       if l.name == layout_name), None)
        assert layout is not None, f"layout {layout_name!r} missing"
        small_pics = [
            shp for shp in layout.shapes
            if str(shp.shape_type) == "PICTURE (13)"
            and shp.width and shp.height
            and (shp.width / EMU) < 8.0   # below the watermark heuristic
        ]
        # Workflow_diagram had 2 logos (DOE + KBase) in addition to the watermark
        assert len(small_pics) >= 2, (
            f"layout {layout_name!r}: expected ≥2 small pictures (logos) "
            f"to survive watermark strip; got {len(small_pics)}"
        )


@requires_source
def test_build_strips_all_source_slides(tmp_path: Path):
    """The source .potx contains 32 example slides (the Gazi 2026 deck);
    the derived master must ship zero slides."""
    bm = _import_build_master()
    out = tmp_path / "build.pptx"
    bt = tmp_path / "tokens.json"
    report = bm.build_master(SOURCE_POTX, out, bt, verbose=False)
    # The .potx Adam supplied has 32 slides
    assert report["deleted_slides_count"] >= 30, (
        f"expected ≥30 slides stripped from .potx (Gazi 2026 deck), "
        f"got {report['deleted_slides_count']}"
    )


# ---------------------------------------------------------------------------
# Round-trip: assemble a sample deck using each layout
# ---------------------------------------------------------------------------

@requires_built_master
def test_round_trip_one_slide_per_layout(tmp_path: Path):
    """Programmatically build a sample deck that uses every named layout
    once. Verify it saves and re-opens without error.

    This is the foundational test for the assembler that v0.1.0-extractors
    will build on. If a layout is structurally broken, this test catches it."""
    from pptx import Presentation
    bm = _import_build_master()

    prs = Presentation(DEST_PPTX)
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}

    for name in bm.NAMED_VOCABULARY:
        assert name in layouts, f"layout '{name}' missing in derived master"
        layout = layouts[name]
        slide = prs.slides.add_slide(layout)
        # Set the title placeholder if available
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = f"Sample slide for layout '{name}'"
                break

    out = tmp_path / "round_trip.pptx"
    prs.save(out)
    # Re-open to confirm the file is structurally valid
    prs2 = Presentation(out)
    assert len(prs2.slides) == 15
    assert {s.slide_layout.name for s in prs2.slides} == set(bm.NAMED_VOCABULARY)

"""Cycle 1 — pre-handoff deliverable validation tests.

Coverage organization:
  - One CLEAN fixture (matches a post-hotfix caulobacter-shape draft):
    every gate passes (G3 may emit an advisory but never P0/P1).
  - One regression fixture per gate, drawn from the live caulobacter
    failure mode: each fixture trips its gate AND ONLY its gate's P0/P1
    findings, modulo G3 advisories which the brief calls out as never
    blocking.
  - DP9b user_intent helper: idempotent merge + explicit-conflict
    fail-loud + read/explicit CLI surface.
  - finalize_deliverable: auto-remediation paths run end-to-end against
    spec fixtures (Pillow-free; we only mutate the spec + skip the
    reassemble step in unit tests).

Each fixture builds a minimal-but-valid working layout under tmp_path:
  draft_N/
    audit/user_intent.json     (DP9b — explicit mode etc.)
    working/slide_spec.json    (the canonical spec)
    working/03_slides/         (per-substory + qa fragments where needed)
    working/05_image_requests/ (where image-completeness gates fire)
    working/05_images/         (where PNGs land)
    working/05_image_decisions.json
    deliverable/draft.pptx     (built on-the-fly for G5 figure-aspect)
"""
from __future__ import annotations

import json
import struct
import sys
import zlib
from pathlib import Path

import pytest

# Load the in-tree modules without requiring a reinstall (same pattern
# as tests/unit/test_v1_1_1_hotfixes.py).
_REPO_ROOT = Path(__file__).resolve().parents[2]
_SKILL_TOOLS = _REPO_ROOT / "src" / "beril_presentation_maker" / "skill" / "tools"
sys.path.insert(0, str(_SKILL_TOOLS))

import finalize_deliverable as fd  # noqa: E402
import user_intent as ui  # noqa: E402
import validate_deliverable as vd  # noqa: E402

# ===========================================================================
# Fixture builders
# ===========================================================================


def _write_json(path: Path, obj: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, indent=2) + "\n", encoding="utf-8")


def _make_minimal_png(width: int, height: int) -> bytes:
    """Return a syntactically valid PNG of the given pixel dimensions.

    Smallest valid PNG: 8-byte signature + IHDR + IDAT(zlib-empty) + IEND.
    Used to satisfy python-pptx's image.size property (which reads
    width/height from the IHDR chunk) without needing Pillow at test time.
    """
    sig = b"\x89PNG\r\n\x1a\n"
    # IHDR: width(4) height(4) bit_depth(1) color_type(1)
    #       compression(1) filter(1) interlace(1)
    ihdr = struct.pack(">II5B", width, height, 8, 2, 0, 0, 0)
    ihdr_chunk = _png_chunk(b"IHDR", ihdr)
    # IDAT: a single zlib-compressed scanline of one black pixel per row.
    # Build a raw row (filter byte 0 + width × 3 RGB bytes), repeat for
    # `height` rows, then zlib-deflate the lot.
    raw = b"".join(b"\x00" + b"\x00\x00\x00" * width for _ in range(height))
    idat = zlib.compress(raw, level=1)
    idat_chunk = _png_chunk(b"IDAT", idat)
    iend_chunk = _png_chunk(b"IEND", b"")
    return sig + ihdr_chunk + idat_chunk + iend_chunk


def _png_chunk(tag: bytes, data: bytes) -> bytes:
    return (
        struct.pack(">I", len(data))
        + tag
        + data
        + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
    )


def _build_clean_draft(tmp_path: Path) -> Path:
    """Build a minimal draft layout that passes all six gates.

    Mirrors a post-hotfix caulobacter shape: explicit user_intent at
    talk-30 (smallest valid budget); slide_spec with title slide
    carrying a real presenter; acknowledgments with contributors; one
    data_figure pointing at a real PNG; no orphan image requests; mode
    consistent across all artifacts.
    """
    # Build under projects/<proj>/talks/draft_1/ so _derive_project_dir
    # works (G1 dirname-leak check uses this).
    proj_dir = tmp_path / "projects" / "smoke_proj"
    draft_dir = proj_dir / "talks" / "draft_1"
    for d in (
        draft_dir / "audit",
        draft_dir / "working" / "03_slides",
        draft_dir / "working" / "05_image_requests",
        draft_dir / "working" / "05_images",
        draft_dir / "deliverable",
    ):
        d.mkdir(parents=True, exist_ok=True)

    # beril.yaml on the project dir.
    (proj_dir / "beril.yaml").write_text(
        "project_id: smoke_proj\n"
        "authors:\n"
        "  - name: Adam Arkin\n"
        "    affiliation: LBNL\n",
        encoding="utf-8",
    )

    # user_intent: explicit talk-30 + STRONG + peer.
    ui.write_user_intent(
        draft_dir,
        mode="talk-30", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=True,
        now="2026-06-07T15:00:00Z",
    )

    # A real PNG figure under the project dir (so the figure path
    # resolves via _derive_project_dir's fallback).
    figures_dir = proj_dir / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    (figures_dir / "fig01.png").write_bytes(_make_minimal_png(800, 600))

    # slide_spec.json (minimal — 25 slides to be inside talk-30's [25,32]
    # budget). One title, one ack, one data_figure (the rest are
    # section_dividers for count).
    slides = [
        {"position": 1, "layout": "title",
         "content": {"title": "Real Science Title",
                     "presenter": "Adam Arkin",
                     "affiliation": "LBNL",
                     "date": "2026-06-07"}},
        {"position": 2, "layout": "data_figure",
         "content": {"title": "Result", "figure": "figures/fig01.png",
                     "figure_caption": "A real caption."}},
        {"position": 25, "layout": "acknowledgments",
         "content": {"contributors": ["Adam Arkin · LBNL"]}},
    ]
    # Pad with section_dividers (any layout that doesn't require figures)
    # to hit the 25-slide floor.
    for i in range(3, 25):
        slides.append({"position": i, "layout": "section_divider",
                       "content": {"title": f"§{i}", "punchline": "x"}})
    slides.sort(key=lambda s: s["position"])
    spec = {
        "schema_version": "slide_spec.v1",
        "project_id": "smoke_proj",
        "mode": "talk-30",
        "audience": "peer",
        "tier": "STRONG",
        "slides": slides,
    }
    _write_json(draft_dir / "working" / "slide_spec.json", spec)
    return draft_dir


# ===========================================================================
# DP9b — user_intent persistence helper
# ===========================================================================


def test_dp9b_user_intent_write_then_read_roundtrip(tmp_path):
    """Basic write + read."""
    ui.write_user_intent(
        tmp_path,
        mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=False, audience_explicit=False,
    )
    assert ui.read_field(tmp_path, "mode") == "talk-45"
    assert ui.field_was_explicit(tmp_path, "mode") is True
    assert ui.read_field(tmp_path, "tier") == "STRONG"
    assert ui.field_was_explicit(tmp_path, "tier") is False


def test_dp9b_user_intent_idempotent_merge_keeps_explicit(tmp_path):
    """Process 1 sets mode=talk-45 explicit. Process 2 (continue) calls
    write with mode=talk-30 NOT explicit. Result: talk-45 wins."""
    ui.write_user_intent(
        tmp_path,
        mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    ui.write_user_intent(
        tmp_path,
        mode="talk-30", tier="STRONG", audience="peer",
        mode_explicit=False, tier_explicit=False, audience_explicit=False,
    )
    assert ui.read_field(tmp_path, "mode") == "talk-45"
    assert ui.field_was_explicit(tmp_path, "mode") is True


def test_dp9b_user_intent_explicit_conflict_returns_findings(tmp_path):
    """Process 1 sets mode=talk-45 explicit. Process 2 re-passes
    mode=talk-30 explicit. The merge keeps talk-45 AND returns a
    conflict finding — the CLI surfaces this as an exit-1 failure."""
    ui.write_user_intent(
        tmp_path,
        mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    _path, conflicts = ui.write_user_intent(
        tmp_path,
        mode="talk-30", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    assert len(conflicts) == 1
    assert "mode" in conflicts[0]
    assert "talk-45" in conflicts[0]
    assert "talk-30" in conflicts[0]
    # The existing-explicit value persisted.
    assert ui.read_field(tmp_path, "mode") == "talk-45"


def test_dp9b_user_intent_cli_read_with_fallback(tmp_path, capsys):
    """`read --field mode --fallback talk-30` on a missing file prints
    talk-30 and exits 0 — the bash hook depends on this."""
    rc = ui.main(["read", str(tmp_path), "--field", "mode",
                  "--fallback", "talk-30"])
    assert rc == 0
    assert capsys.readouterr().out.strip() == "talk-30"


def test_dp9b_user_intent_cli_explicit_reads_sentinel(tmp_path, capsys):
    ui.write_user_intent(
        tmp_path,
        mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=False, audience_explicit=False,
    )
    rc = ui.main(["explicit", str(tmp_path), "--field", "mode"])
    assert rc == 0
    assert capsys.readouterr().out.strip() == "1"
    rc = ui.main(["explicit", str(tmp_path), "--field", "tier"])
    assert rc == 0
    assert capsys.readouterr().out.strip() == "0"


def test_dp9b_user_intent_cli_write_conflict_exits_1(tmp_path, capsys):
    ui.write_user_intent(
        tmp_path,
        mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    rc = ui.main([
        "write", str(tmp_path),
        "--mode", "talk-30", "--tier", "STRONG", "--audience", "peer",
        "--mode-explicit", "1", "--tier-explicit", "1",
        "--audience-explicit", "0",
    ])
    assert rc == 1
    assert "CONFLICT" in capsys.readouterr().err


# ===========================================================================
# Clean fixture — every gate passes (P0/P1 == 0)
# ===========================================================================


def test_clean_fixture_passes_all_p0_p1(tmp_path):
    """The hotfixed-caulobacter shape: no P0 + no P1 findings. G5
    can't run without a real .pptx (skipped via the no-deck path
    returning P0 — but we strip that one expected finding here, since
    the test is about the SPEC-side gates being clean)."""
    draft_dir = _build_clean_draft(tmp_path)
    findings = vd.validate(draft_dir)
    # G5 will emit one P0 because there's no rendered deck. That's a
    # different test surface (covered by the dedicated G5 fixture
    # below). Strip it for this assertion.
    findings_no_g5 = [f for f in findings if f.gate != "figure_integrity"]
    blocking = [
        f for f in findings_no_g5
        if f.severity in (vd.SEVERITY_P0, vd.SEVERITY_P1)
    ]
    assert blocking == [], (
        f"clean fixture had {len(blocking)} P0/P1 finding(s): "
        f"{[(f.id, f.message) for f in blocking]}"
    )


# ===========================================================================
# G1 placeholder_or_leaked_template — three failure modes
# ===========================================================================


def test_g1_tbd_presenter_fails_loud(tmp_path):
    """The caulobacter TBD-presenter failure. P0; auto-remediable
    via populate_title_from_beril."""
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    title_slide = next(s for s in spec["slides"] if s["layout"] == "title")
    title_slide["content"]["presenter"] = "TBD"
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g1_placeholder_or_leaked_template(
        draft_dir, vd._load_slide_spec(draft_dir))
    ids = [f.id for f in findings]
    assert any("g1:presenter_tbd" in i for i in ids), (
        f"expected g1:presenter_tbd finding; got {ids}"
    )
    presenter_finding = next(f for f in findings if "presenter_tbd" in f.id)
    assert presenter_finding.severity == vd.SEVERITY_P0
    assert presenter_finding.remediation.kind == vd.REMEDIATION_AUTO
    assert presenter_finding.remediation.action == vd.AUTO_POPULATE_TITLE


def test_g1_dirname_leak_in_title_fails_loud(tmp_path):
    """The caulobacter `Lipida` typo: project dir is
    `caulobacter_fur_lipida_loss`; the LLM hallucinated `Lipida` into
    the title. P0; auto-remediable via strip_dirname_token."""
    proj_dir = tmp_path / "projects" / "caulobacter_fur_lipida_loss"
    draft_dir = proj_dir / "talks" / "draft_1"
    for d in (draft_dir / "audit", draft_dir / "working" / "03_slides",
              draft_dir / "working" / "05_image_requests",
              draft_dir / "working" / "05_images",
              draft_dir / "deliverable"):
        d.mkdir(parents=True, exist_ok=True)
    (proj_dir / "beril.yaml").write_text(
        "project_id: caulobacter_fur_lipida_loss\nauthors: []\n",
        encoding="utf-8",
    )
    ui.write_user_intent(
        draft_dir, mode="talk-30", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    spec = {
        "schema_version": "slide_spec.v1",
        "project_id": "caulobacter_fur_lipida_loss",
        "mode": "talk-30",
        "audience": "peer", "tier": "STRONG",
        "slides": [{
            "position": 1, "layout": "title",
            "content": {
                "title": "Caulobacter Lipida response: novel findings",
                "presenter": "Adam Arkin",
                "affiliation": "LBNL", "date": "2026-06-07",
            },
        }],
    }
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g1_placeholder_or_leaked_template(
        draft_dir, vd._load_slide_spec(draft_dir))
    leak_findings = [f for f in findings if "dirname_leak" in f.id]
    assert leak_findings, (
        f"expected g1:title_dirname_leak finding; got "
        f"{[f.id for f in findings]}"
    )
    leak = leak_findings[0]
    assert leak.severity == vd.SEVERITY_P0
    assert leak.remediation.action == vd.AUTO_STRIP_DIRNAME_TOKEN


def test_g1_acknowledgments_all_tbd_fires(tmp_path):
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    ack = next(s for s in spec["slides"] if s["layout"] == "acknowledgments")
    ack["content"]["contributors"] = ["TBD - production"]
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g1_placeholder_or_leaked_template(
        draft_dir, vd._load_slide_spec(draft_dir))
    ack_findings = [f for f in findings if "acknowledgments_tbd" in f.id]
    assert ack_findings
    assert ack_findings[0].severity == vd.SEVERITY_P1
    assert ack_findings[0].remediation.action == vd.AUTO_POPULATE_TITLE


# ===========================================================================
# G2 image_completeness — orphan request + unbound placeholder
# ===========================================================================


def test_g2_orphan_image_request_fails_loud(tmp_path):
    """The DP3 caulobacter failure: request written, no PNG, no
    manifest entry. P0; targeted-remediable via resume-from image_gen."""
    draft_dir = _build_clean_draft(tmp_path)
    req_dir = draft_dir / "working" / "05_image_requests"
    _write_json(req_dir / "S4-pos3_request.json",
                {"slide_id_target": "S4-pos3",
                 "image_prompt": "x", "worst_case_cost_usd": 0.04})

    findings = vd.check_g2_image_completeness(
        draft_dir, vd._load_slide_spec(draft_dir))
    orphan = [f for f in findings if "orphan_request" in f.id]
    assert orphan
    assert orphan[0].severity == vd.SEVERITY_P0
    assert orphan[0].remediation.kind == vd.REMEDIATION_TARGETED
    assert "--resume-from image_gen" in orphan[0].remediation.command


def test_g2_unbound_image_placeholder_fails_loud(tmp_path):
    """The deck has a concept_illustration slide whose image_path is
    still {TBD} — image_gen/bind never completed for it. P0."""
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    spec["slides"].insert(1, {
        "position": 99, "layout": "concept_illustration",
        "content": {"title": "Concept", "image_path": "{TBD}",
                    "image_prompt": "test"},
    })
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g2_image_completeness(
        draft_dir, vd._load_slide_spec(draft_dir))
    unbound = [f for f in findings if "unbound_image_placeholder" in f.id]
    assert unbound
    assert unbound[0].severity == vd.SEVERITY_P0


# ===========================================================================
# G3 slide_count_vs_budget — advisory only
# ===========================================================================


def test_g3_under_budget_emits_advisory(tmp_path):
    """37 slides at talk-45 (budget [35,48]) is IN-band → no finding.
    Reduce to 20 slides → under-band advisory."""
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    spec["slides"] = spec["slides"][:10]   # 10 slides; talk-30 floor=25
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g3_slide_count_vs_budget(
        draft_dir, vd._load_slide_spec(draft_dir), "talk-30")
    out_of_band = [f for f in findings if "out_of_band" in f.id]
    assert out_of_band
    assert out_of_band[0].severity == vd.SEVERITY_ADVISORY
    assert "under" in out_of_band[0].id


def test_g3_in_band_emits_nothing(tmp_path):
    """25-slide deck under talk-30 budget [25,32] → clean."""
    draft_dir = _build_clean_draft(tmp_path)
    findings = vd.check_g3_slide_count_vs_budget(
        draft_dir, vd._load_slide_spec(draft_dir), "talk-30")
    assert findings == []


# ===========================================================================
# G4 mode_vs_user_intent — the load-bearing new gate (DP9b)
# ===========================================================================


def test_g4_intent_matches_artifacts_passes(tmp_path):
    """User picked talk-30; every artifact is talk-30 → no finding."""
    draft_dir = _build_clean_draft(tmp_path)
    findings = vd.check_g4_mode_vs_user_intent(draft_dir)
    assert findings == []


def test_g4_intent_disagrees_with_image_decisions_fails_loud(tmp_path):
    """The caulobacter DP9b failure: user explicit talk-45, every
    artifact agrees on talk-30. v1.1.1 cross-artifact check would say
    'consistent'; G4 catches it because user_intent.json carries
    'talk-45' as the truth anchor."""
    proj_dir = tmp_path / "projects" / "fortyfive_proj"
    draft_dir = proj_dir / "talks" / "draft_1"
    for d in (draft_dir / "audit",
              draft_dir / "working" / "03_slides"):
        d.mkdir(parents=True, exist_ok=True)
    # User explicit pick: talk-45.
    ui.write_user_intent(
        draft_dir, mode="talk-45", tier="STRONG", audience="peer",
        mode_explicit=True, tier_explicit=True, audience_explicit=False,
    )
    # But every artifact agrees on talk-30 (the silent drop).
    _write_json(draft_dir / "working" / "slide_spec.json",
                {"mode": "talk-30", "slides": []})
    _write_json(draft_dir / "working" / "05_image_decisions.json",
                {"mode": "talk-30", "decisions": []})
    _write_json(draft_dir / "working" / "03_slides" / "qa_anticipated.json",
                {"mode": "talk-30", "items": []})

    findings = vd.check_g4_mode_vs_user_intent(draft_dir)
    assert len(findings) == 3
    for f in findings:
        assert f.severity == vd.SEVERITY_P0
        assert "talk-45" in f.message
        assert "talk-30" in f.message
        assert f.remediation.kind == vd.REMEDIATION_TARGETED


def test_g4_user_default_not_explicit_advisory_only(tmp_path):
    """If the user never explicitly picked, G4 surfaces an advisory
    rather than blocking — we don't punish operators for inheriting
    a default."""
    proj_dir = tmp_path / "projects" / "default_proj"
    draft_dir = proj_dir / "talks" / "draft_1"
    (draft_dir / "audit").mkdir(parents=True, exist_ok=True)
    ui.write_user_intent(
        draft_dir, mode="talk-30", tier="STRONG", audience="peer",
        mode_explicit=False, tier_explicit=False, audience_explicit=False,
    )
    findings = vd.check_g4_mode_vs_user_intent(draft_dir)
    assert len(findings) == 1
    assert findings[0].severity == vd.SEVERITY_ADVISORY
    assert "not_explicit" in findings[0].id


def test_g4_no_user_intent_advisory_for_legacy_drafts(tmp_path):
    """Pre-v1.2.0 drafts without user_intent.json get a single
    advisory — never block; the file just isn't there to compare."""
    draft_dir = tmp_path / "draft_legacy"
    draft_dir.mkdir(parents=True)
    findings = vd.check_g4_mode_vs_user_intent(draft_dir)
    assert len(findings) == 1
    assert findings[0].severity == vd.SEVERITY_ADVISORY
    assert findings[0].id == "g4:no_user_intent"


# ===========================================================================
# G5 figure_integrity (display vs. native aspect)
# ===========================================================================


def test_g5_no_deck_fails_loud(tmp_path):
    draft_dir = tmp_path / "draft"
    draft_dir.mkdir()
    findings = vd.check_g5_figure_integrity(draft_dir)
    assert len(findings) == 1
    assert findings[0].id == "g5:no_deck"
    assert findings[0].severity == vd.SEVERITY_P0
    assert findings[0].remediation.action == vd.AUTO_REASSEMBLE


def test_g5_correct_aspect_passes(tmp_path):
    """A real pptx round-trip with a non-stretched picture passes G5.
    Uses the same path the assembler uses."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches
    draft_dir = tmp_path / "draft"
    (draft_dir / "deliverable").mkdir(parents=True)
    png_path = tmp_path / "fig.png"
    png_path.write_bytes(_make_minimal_png(800, 600))

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # 800:600 = 4:3 native. Place at width=4 inches → height = 3 in
    # (preserve-aspect from python-pptx). Box ratio matches native.
    slide.shapes.add_picture(str(png_path), Inches(1), Inches(1),
                             width=Inches(4))
    deck_path = draft_dir / "deliverable" / "draft.pptx"
    prs.save(str(deck_path))

    findings = vd.check_g5_figure_integrity(draft_dir)
    aspect_findings = [f for f in findings if "aspect_skew" in f.id]
    assert aspect_findings == []


def test_g5_stretched_figure_fails_loud(tmp_path):
    """Simulate the DP4 pre-hotfix behavior by adding a picture with
    both width AND height that don't match the source aspect → G5 trips.
    """
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches
    draft_dir = tmp_path / "draft"
    (draft_dir / "deliverable").mkdir(parents=True)
    # Square source (1:1); display 4:1 (stretched 4×).
    png_path = tmp_path / "fig.png"
    png_path.write_bytes(_make_minimal_png(500, 500))

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png_path), Inches(1), Inches(1),
                             width=Inches(8), height=Inches(2))  # stretched
    deck_path = draft_dir / "deliverable" / "draft.pptx"
    prs.save(str(deck_path))

    findings = vd.check_g5_figure_integrity(draft_dir)
    aspect = [f for f in findings if "aspect_skew" in f.id]
    assert aspect
    assert aspect[0].severity == vd.SEVERITY_P0
    assert aspect[0].remediation.action == vd.AUTO_REASSEMBLE


# ===========================================================================
# G6 figure_path_resolution
# ===========================================================================


def test_g6_unresolved_figure_fails_loud(tmp_path):
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    # Replace the real figure path with a non-existent one.
    figure_slide = next(s for s in spec["slides"]
                        if s["layout"] == "data_figure")
    figure_slide["content"]["figure"] = "figures/does_not_exist.png"
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    findings = vd.check_g6_figure_path_resolution(
        draft_dir, vd._load_slide_spec(draft_dir))
    assert findings
    assert findings[0].severity == vd.SEVERITY_P0
    assert "does_not_exist.png" in findings[0].message


def test_g6_resolved_figure_passes(tmp_path):
    """The clean fixture has a real figure path → no G6 findings."""
    draft_dir = _build_clean_draft(tmp_path)
    findings = vd.check_g6_figure_path_resolution(
        draft_dir, vd._load_slide_spec(draft_dir))
    assert findings == []


def test_g6_tbd_placeholder_skipped(tmp_path):
    """{TBD} placeholders are G2's territory; G6 ignores them so a
    mid-pipeline check doesn't double-fire."""
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    spec["slides"].append({
        "position": 100, "layout": "concept_illustration",
        "content": {"title": "x", "image_path": "{TBD}",
                    "image_prompt": "p"},
    })
    _write_json(draft_dir / "working" / "slide_spec.json", spec)
    findings = vd.check_g6_figure_path_resolution(
        draft_dir, vd._load_slide_spec(draft_dir))
    # No G6 finding from the {TBD} placeholder (G2 covers it).
    assert all("g6:unresolved_figure:100" not in f.id for f in findings)


# ===========================================================================
# Schema shape — projectable tokens, no free-text in slots
# ===========================================================================


def test_schema_finding_fields_are_tokens(tmp_path):
    """Telemetry-readiness: gate, severity, remediation.kind are all
    drawn from the frozen vocabularies — no free-text in those fields."""
    draft_dir = _build_clean_draft(tmp_path)
    # Force-fire G1 + G2 + G4 to get a mix.
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    next(s for s in spec["slides"]
         if s["layout"] == "title")["content"]["presenter"] = "TBD"
    _write_json(draft_dir / "working" / "slide_spec.json", spec)
    _write_json(
        draft_dir / "working" / "05_image_requests" / "S1-pos1_request.json",
        {"slide_id_target": "S1-pos1", "image_prompt": "x",
         "worst_case_cost_usd": 0.04},
    )

    findings = vd.validate(draft_dir)
    for f in findings:
        assert f.gate in vd.GATES, f"gate {f.gate!r} not in vocab"
        assert f.severity in vd.SEVERITIES, (
            f"severity {f.severity!r} not in vocab"
        )
        assert f.remediation.kind in vd.REMEDIATION_KINDS, (
            f"remediation.kind {f.remediation.kind!r} not in vocab"
        )


def test_schema_payload_shape(tmp_path):
    """write_findings emits the deliverable-validation.v1 envelope
    with summary + findings[] in the documented shape."""
    draft_dir = _build_clean_draft(tmp_path)
    findings = vd.validate(draft_dir)
    out = vd.write_findings(draft_dir, findings)
    payload = json.loads(out.read_text(encoding="utf-8"))
    assert payload["schema_version"] == "deliverable-validation.v1"
    assert "summary" in payload
    s = payload["summary"]
    assert {"total", "by_gate", "by_severity", "by_remediation_kind",
            "blocking"} <= set(s.keys())
    # Aggregated counters use the frozen vocabs as keys.
    assert set(s["by_gate"].keys()) == set(vd.GATES)
    assert set(s["by_severity"].keys()) == set(vd.SEVERITIES)
    assert set(s["by_remediation_kind"].keys()) == set(vd.REMEDIATION_KINDS)


def test_readiness_exit_code_clean_is_zero():
    assert vd.readiness_exit_code([]) == 0


def test_readiness_exit_code_advisory_is_zero():
    f = vd.Finding(
        id="x", gate=vd.GATES[2], severity=vd.SEVERITY_ADVISORY,
        slide_id_or_target="deck", message="advisory",
        remediation=vd.Remediation(kind=vd.REMEDIATION_ADVISORY),
    )
    assert vd.readiness_exit_code([f]) == 0


def test_readiness_exit_code_p0_is_one():
    f = vd.Finding(
        id="x", gate=vd.GATES[0], severity=vd.SEVERITY_P0,
        slide_id_or_target="1", message="p0",
        remediation=vd.Remediation(kind=vd.REMEDIATION_AUTO),
    )
    assert vd.readiness_exit_code([f]) == 1


# ===========================================================================
# finalize_deliverable — auto-remediation paths
# ===========================================================================


def test_finalize_populate_title_from_beril_mutates_spec(tmp_path):
    """G1 presenter=TBD → finalize reads beril.yaml, sets presenter,
    writes spec back. We test the handler directly (skips reassemble)."""
    draft_dir = _build_clean_draft(tmp_path)
    spec = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    next(s for s in spec["slides"]
         if s["layout"] == "title")["content"]["presenter"] = "TBD"
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    ok, msg = fd._populate_from_beril_yaml(draft_dir)
    assert ok, msg
    spec2 = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    new_presenter = next(s for s in spec2["slides"]
                         if s["layout"] == "title")["content"]["presenter"]
    assert new_presenter == "Adam Arkin"


def test_finalize_strip_dirname_token_mutates_title(tmp_path):
    """G1 dirname-leak → finalize strips the token; writes spec back."""
    proj_dir = tmp_path / "projects" / "caulobacter_fur_lipida_loss"
    draft_dir = proj_dir / "talks" / "draft_1"
    for d in (draft_dir / "audit", draft_dir / "working" / "03_slides",
              draft_dir / "deliverable"):
        d.mkdir(parents=True, exist_ok=True)
    (proj_dir / "beril.yaml").write_text(
        "project_id: caulobacter_fur_lipida_loss\nauthors: []\n",
        encoding="utf-8",
    )
    spec = {
        "schema_version": "slide_spec.v1", "project_id": "x",
        "mode": "talk-30", "audience": "peer", "tier": "STRONG",
        "slides": [{
            "position": 1, "layout": "title",
            "content": {
                "title": "Caulobacter Lipida response: novel findings",
                "presenter": "Adam Arkin", "affiliation": "LBNL",
                "date": "2026-06-07",
            },
        }],
    }
    _write_json(draft_dir / "working" / "slide_spec.json", spec)

    ok, msg = fd._strip_dirname_token(draft_dir)
    assert ok, msg
    spec2 = json.loads((draft_dir / "working" / "slide_spec.json").read_text())
    new_title = next(s for s in spec2["slides"]
                     if s["layout"] == "title")["content"]["title"]
    # The distinctive ≥5-char tokens (caulobacter, lipida) are stripped.
    assert "caulobacter" not in new_title.lower()
    assert "lipida" not in new_title.lower()


def test_finalize_skips_reassemble_when_no_auto_actions(tmp_path):
    """A clean run (no auto-findings, no reassemble requested) doesn't
    invoke the reassemble handler. Pure-orchestration check."""
    draft_dir = _build_clean_draft(tmp_path)
    # First pass: write findings (will include G5 no-deck P0 +
    # auto/reassemble action). Drop that finding manually so the
    # orchestration path takes the no-auto branch.
    vd.write_findings(draft_dir, [])
    result = fd.finalize(draft_dir)
    assert result["actions_applied"] == []

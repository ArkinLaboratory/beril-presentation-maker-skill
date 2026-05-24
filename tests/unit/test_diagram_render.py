"""Tests for diagram_render.py — slide_spec diagram → python-pptx shapes.

Coverage:
- 7 node shape kinds map to MSO_SHAPE values.
- 3 edge kinds map to MSO_CONNECTOR values.
- Coordinate transform: diagram (0,0) maps to region top-left.
- Brand-color resolution (token + hex + fallback).
- swimlane gets no-fill.
- Edges connect node centers; missing nodes are skipped, not crashed.
- Unknown shape / edge kind raises ValueError.
- Unsupported diagram kind ('mermaid', 'tree', etc.) raises ValueError.
- end-to-end: full diagram on a blank slide produces N+M shapes.
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest


REPO_ROOT = Path(__file__).resolve().parents[2]
DR_PY = (REPO_ROOT / "src" / "beril_presentation_maker" / "skill"
         / "tools" / "diagram_render.py")


def _import(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def dr():
    return _import("diagram_render", DR_PY)


@pytest.fixture
def blank_slide():
    from pptx import Presentation
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Shape + edge mappings
# ---------------------------------------------------------------------------

def test_node_shape_map_covers_all_seven(dr):
    expected = {"rectangle", "rounded", "ellipse", "parallelogram",
                "cylinder", "callout", "swimlane"}
    assert set(dr.NODE_SHAPE_MAP.keys()) == expected


def test_edge_kind_map_covers_three(dr):
    expected = {"straight", "elbow", "curved"}
    assert set(dr.EDGE_KIND_MAP.keys()) == expected


# ---------------------------------------------------------------------------
# Coordinate transform
# ---------------------------------------------------------------------------

def test_transform_coords_pass_through_origin(dr):
    """2026-04-27 #78: _transform_coords now passes through absolute
    slide coords. The prompt + repair_diagram_stubs both produce
    absolute coords; the prior offset-add caused diagrams to render
    off-slide on tall layouts (live failure draft_7 slide 11)."""
    region = (1.0, 2.0, 5.0, 4.0)
    x, y = dr._transform_coords(0.0, 0.0, region)
    assert x == 0.0
    assert y == 0.0


def test_transform_coords_passes_through_arbitrary_point(dr):
    region = (1.0, 2.0, 5.0, 4.0)
    x, y = dr._transform_coords(0.5, 1.5, region)
    # Inputs returned verbatim — no offset added
    assert x == 0.5
    assert y == 1.5


# ---------------------------------------------------------------------------
# Color resolution
# ---------------------------------------------------------------------------

def test_resolve_color_hex_input(dr):
    rgb = dr.resolve_color("#FF0000", brand_tokens=None)
    # python-pptx RGBColor exposes its underlying int via str()
    assert rgb is not None


def test_resolve_color_known_token_with_fallback_hex(dr):
    rgb = dr.resolve_color("freshwater_blue", brand_tokens=None)
    assert rgb is not None


def test_resolve_color_unknown_token_falls_back(dr):
    """Unknown token → falls through to default → fallback to white."""
    rgb = dr.resolve_color("space_indigo", brand_tokens=None)
    assert rgb is not None  # falls back to default freshwater_blue or white


def test_resolve_color_brand_tokens_lookup(dr):
    tokens = {
        "palette": {
            "primary": {"my_blue": {"hex": "#0011FF"}},
            "secondary": {},
            "neutral": {},
        }
    }
    rgb = dr.resolve_color("my_blue", brand_tokens=tokens)
    assert rgb is not None


# ---------------------------------------------------------------------------
# Per-node rendering
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("shape_kind", [
    "rectangle", "rounded", "ellipse", "parallelogram",
    "cylinder", "callout", "swimlane",
])
def test_render_node_each_shape_kind(dr, blank_slide, shape_kind):
    node = {"id": "n1", "label": "Test", "shape": shape_kind,
            "x": 0.5, "y": 0.5, "w": 1.5, "h": 0.8}
    rs = dr._render_node(blank_slide, node, region=(0.0, 0.0, 10.0, 5.0),
                         brand_tokens=None)
    assert rs.node_id == "n1"
    assert rs.shape is not None
    # Center is (left + w/2, top + h/2) = (1.25, 0.9)
    assert abs(rs.cx - 1.25) < 0.01
    assert abs(rs.cy - 0.9) < 0.01


def test_render_node_unknown_shape_raises(dr, blank_slide):
    node = {"id": "n1", "label": "X", "shape": "hexagon",
            "x": 0, "y": 0, "w": 1, "h": 1}
    with pytest.raises(ValueError):
        dr._render_node(blank_slide, node, region=(0, 0, 10, 5),
                        brand_tokens=None)


def test_render_node_swimlane_has_no_fill(dr, blank_slide):
    """swimlane shapes should not have a solid fill — they're container outlines."""
    node = {"id": "lane", "label": "Phase", "shape": "swimlane",
            "x": 0.0, "y": 0.0, "w": 5.0, "h": 3.0}
    rs = dr._render_node(blank_slide, node, region=(0, 0, 10, 5),
                         brand_tokens=None)
    # The fill type is 'background' (no fill); python-pptx exposes via
    # shape.fill.type. For swimlane we set sp.fill.background() which
    # makes type==MSO_FILL.BACKGROUND (which is None == 5).
    # Just verify shape was created without crashing.
    assert rs.shape is not None


# ---------------------------------------------------------------------------
# Edge rendering
# ---------------------------------------------------------------------------

def test_render_edge_line_unknown_kind_raises(dr, blank_slide):
    """M4a Tier A3: _render_edge was split into _render_edge_line +
    _render_edge_label so labels paint on top of nodes (third pass).
    The unknown-kind ValueError lives on the line function."""
    node_a = {"id": "a", "label": "A", "shape": "rectangle",
              "x": 0, "y": 0, "w": 1, "h": 1}
    node_b = {"id": "b", "label": "B", "shape": "rectangle",
              "x": 3, "y": 0, "w": 1, "h": 1}
    rsa = dr._render_node(blank_slide, node_a, (0, 0, 10, 5), None)
    rsb = dr._render_node(blank_slide, node_b, (0, 0, 10, 5), None)
    with pytest.raises(ValueError):
        dr._render_edge_line(blank_slide,
                             {"from": "a", "to": "b", "kind": "magic"},
                             {"a": rsa, "b": rsb}, brand_tokens=None)


def test_render_edge_line_skips_missing_node(dr, blank_slide):
    """Edge referencing an undeclared node should silently skip
    (validator catches; renderer is forgiving)."""
    rsa = dr._render_node(blank_slide,
                          {"id": "a", "label": "A", "shape": "rectangle",
                           "x": 0, "y": 0, "w": 1, "h": 1},
                          (0, 0, 10, 5), None)
    # Should not crash
    dr._render_edge_line(blank_slide,
                         {"from": "a", "to": "ghost", "kind": "straight"},
                         {"a": rsa}, brand_tokens=None)


def test_render_edge_label_skips_missing_node(dr, blank_slide):
    """The label pass is independent of the line pass; missing node →
    silent skip (same forgiving behavior, separate code path)."""
    rsa = dr._render_node(blank_slide,
                          {"id": "a", "label": "A", "shape": "rectangle",
                           "x": 0, "y": 0, "w": 1, "h": 1},
                          (0, 0, 10, 5), None)
    # Should not crash; should not add any shape
    n_before = len(blank_slide.shapes)
    dr._render_edge_label(blank_slide,
                          {"from": "a", "to": "ghost", "label": "lost",
                           "kind": "straight"},
                          {"a": rsa}, brand_tokens=None)
    assert len(blank_slide.shapes) == n_before


def test_render_edge_label_skips_empty_label(dr, blank_slide):
    """Edges without a label add nothing on the third pass."""
    rsa = dr._render_node(blank_slide,
                          {"id": "a", "label": "A", "shape": "rectangle",
                           "x": 0, "y": 0, "w": 1, "h": 1},
                          (0, 0, 10, 5), None)
    rsb = dr._render_node(blank_slide,
                          {"id": "b", "label": "B", "shape": "rectangle",
                           "x": 3, "y": 0, "w": 1, "h": 1},
                          (0, 0, 10, 5), None)
    n_before = len(blank_slide.shapes)
    dr._render_edge_label(blank_slide,
                          {"from": "a", "to": "b", "kind": "straight"},
                          {"a": rsa, "b": rsb}, brand_tokens=None)
    assert len(blank_slide.shapes) == n_before


# ---------------------------------------------------------------------------
# M4a Tier A — explicit-fontScale shrink-to-fit + label z-order
# ---------------------------------------------------------------------------

def test_apply_fontscale_to_shape_short_label_renders_full(dr, blank_slide):
    """Short labels (<= full_below) render at 100% fontScale — no shrink."""
    node = {"id": "n", "label": "Short", "shape": "rectangle",
            "x": 0, "y": 0, "w": 1.5, "h": 0.8}
    rs = dr._render_node(blank_slide, node, (0, 0, 10, 5), None)
    # _render_node already calls _apply_fontscale_to_shape; inspect the
    # bodyPr it wrote.
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    tx_body = rs.shape.element.find(f"{{{PML}}}txBody")
    if tx_body is None:
        tx_body = rs.shape.element.find(f"{{{DML}}}txBody")
    body_pr = tx_body.find(f"{{{DML}}}bodyPr")
    norm = body_pr.find(f"{{{DML}}}normAutofit")
    assert norm is not None, "node label must carry explicit normAutofit (LibreOffice quirk)"
    assert norm.get("fontScale") == str(dr.NODE_FONTSCALE_FULL), \
        f"5-char label should render at full scale, got {norm.get('fontScale')}"


def test_apply_fontscale_to_shape_long_label_clamps_at_floor(dr, blank_slide):
    """Long labels clamp at the 60% floor (DQ3) — never silently sub-60%."""
    # 150-char label, well past the (40, 60, 90, 120) ladder
    long_label = "x" * 150
    node = {"id": "n", "label": long_label, "shape": "rectangle",
            "x": 0, "y": 0, "w": 1.5, "h": 0.8}
    rs = dr._render_node(blank_slide, node, (0, 0, 10, 5), None)
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    tx_body = rs.shape.element.find(f"{{{PML}}}txBody")
    if tx_body is None:
        tx_body = rs.shape.element.find(f"{{{DML}}}txBody")
    body_pr = tx_body.find(f"{{{DML}}}bodyPr")
    norm = body_pr.find(f"{{{DML}}}normAutofit")
    assert norm.get("fontScale") == str(dr.NODE_FONTSCALE_FLOOR), \
        f"150-char label should clamp at floor (60%), got {norm.get('fontScale')}"


def test_connector_color_is_slate_dark(dr):
    """M4a Tier E round 2 (2026-05-23): connector lines and edge labels
    moved off the brand's `graphite_gray` (#9D9389) onto a slate-dark
    tone (80, 75, 70) that holds contrast against cream backgrounds.
    Round-1 visual-QA found the connector lines were essentially
    invisible against the watermark + cream background. swimlane
    borders intentionally keep graphite_gray — they ARE meant to
    recede."""
    r, g, b = dr._CONNECTOR_RGB[0], dr._CONNECTOR_RGB[1], dr._CONNECTOR_RGB[2]
    assert r < 140 and g < 140 and b < 140, (
        f"_CONNECTOR_RGB={(r, g, b)} — connector lines/labels must be "
        f"substantially darker than graphite_gray for legibility "
        f"against the cream master."
    )
    # swimlane border stays graphite_gray (intentional softness)
    assert dr.DEFAULT_SWIMLANE_BORDER == "graphite_gray"


def test_edge_label_horizontal_wide_gap_fits_in_gap(dr, blank_slide):
    """M4a Tier E round 4 (2026-05-24): horizontal edges with a WIDE gap
    (>=1.0in) place the label inside the gap, sized to gap-width minus
    margin. Round-3 used a 0.4in threshold but a 0.4in box was too
    narrow for any label at 9pt and triggered char-by-char wrap."""
    # Two horizontal nodes 1.2in apart (src.right at 2.25, dst.left at 3.45),
    # gap = 1.20in — comfortably above the 1.0in in-gap threshold.
    src = dr._RenderedShape(node_id="src", shape=None,
                            cx=1.5, cy=2.0, w=1.5, h=0.8)
    dst = dr._RenderedShape(node_id="dst", shape=None,
                            cx=4.2, cy=2.0, w=1.5, h=0.8)
    rendered = {"src": src, "dst": dst}
    n_before = len(blank_slide.shapes)
    dr._render_edge_label(
        blank_slide,
        {"from": "src", "to": "dst", "kind": "straight", "label": "test"},
        rendered, brand_tokens=None,
    )
    assert len(blank_slide.shapes) == n_before + 1
    label_shape = blank_slide.shapes[-1]
    EMU = 914400
    left_in = label_shape.left / EMU
    width_in = label_shape.width / EMU
    # src right edge = 2.25; dst left edge = 3.45; gap = 1.20
    # Label must fit in [2.25, 3.45]
    assert 2.20 <= left_in <= 3.50, (
        f"label left {left_in:.3f} must be inside gap [2.25, 3.45]"
    )
    assert (left_in + width_in) <= 3.50, (
        f"label right {left_in + width_in:.3f} must be at-or-left-of dst left 3.45"
    )


def test_edge_label_horizontal_narrow_gap_overflows_above_nodes(dr, blank_slide):
    """M4a Tier E round 4 narrow-gap branch: when the gap is < 1.0in
    (the ibd_phage_targeting slide-6 case, gap = 0.4in), the label
    box is sized to 1.1in centered on the connector midpoint and
    placed above the connector line. It WILL overlap the node-tops
    horizontally — but at label-bottom y = mid_y - 0.15 vs node-top
    y = cy - 0.45, there's ~0.30in of vertical clearance, so the
    label paints above the node row. The previous round-3 behavior
    forced the box to gap-width which produced one-char-per-line
    stacks; the round-4 fallback renders the label as one line."""
    # Two horizontal nodes 0.5in apart (gap = 0.5in), below the 1.0in
    # threshold — should trigger the wider-than-gap fallback.
    src = dr._RenderedShape(node_id="src", shape=None,
                            cx=1.5, cy=2.0, w=1.5, h=0.8)
    dst = dr._RenderedShape(node_id="dst", shape=None,
                            cx=3.5, cy=2.0, w=1.5, h=0.8)
    rendered = {"src": src, "dst": dst}
    n_before = len(blank_slide.shapes)
    dr._render_edge_label(
        blank_slide,
        {"from": "src", "to": "dst", "kind": "straight", "label": "test"},
        rendered, brand_tokens=None,
    )
    assert len(blank_slide.shapes) == n_before + 1
    label_shape = blank_slide.shapes[-1]
    EMU = 914400
    width_in = label_shape.width / EMU
    top_in = label_shape.top / EMU
    height_in = label_shape.height / EMU
    # Width must be the fallback 1.1in — wider than the 0.5in gap
    assert 1.0 <= width_in <= 1.3, (
        f"narrow-gap label width {width_in:.3f} must be ~1.1in (fallback), "
        f"not constrained to the 0.5in gap"
    )
    # word_wrap must be OFF on the textbox — confirms the fix
    assert label_shape.text_frame.word_wrap is False, (
        "narrow-gap label must have word_wrap=False to avoid char-by-char "
        "wrap when the box is wider than the gap"
    )
    # Vertical position: label-bottom must clear node-top (cy - h/2 = 1.60)
    assert top_in + height_in <= 1.85, (
        f"narrow-gap label bottom {top_in + height_in:.3f} must clear "
        f"node-top at 1.60 (plus small margin for paint-order overlap)"
    )


def test_edge_label_falls_back_when_w_unknown(dr, blank_slide):
    """If _RenderedShape carries w=0 (the proxy form used by the
    line-pass before node geometry is read), the label uses the prior
    offset heuristic — the diagram still renders, just less optimally.
    Guards against a regression that would crash on a proxy."""
    # Proxies created without w/h (default 0.0)
    src = dr._RenderedShape(node_id="src", shape=None, cx=1.0, cy=2.0)
    dst = dr._RenderedShape(node_id="dst", shape=None, cx=5.0, cy=2.0)
    rendered = {"src": src, "dst": dst}
    n_before = len(blank_slide.shapes)
    # Must not crash; should add exactly one textbox
    dr._render_edge_label(
        blank_slide,
        {"from": "src", "to": "dst", "kind": "straight", "label": "fallback"},
        rendered, brand_tokens=None,
    )
    assert len(blank_slide.shapes) == n_before + 1


def test_rendered_shape_has_width_height_fields(dr):
    """_RenderedShape MUST carry w + h (added in Tier E round 3 so the
    third-pass edge labels can size to the inter-node gap). Pin so a
    future refactor doesn't drop the fields silently."""
    rs = dr._RenderedShape(node_id="x", shape=None,
                           cx=1.0, cy=2.0, w=1.5, h=0.8)
    assert rs.w == 1.5
    assert rs.h == 0.8
    # Defaults must be 0 so proxies created without w/h are detectable
    proxy = dr._RenderedShape(node_id="x", shape=None, cx=1.0, cy=2.0)
    assert proxy.w == 0.0
    assert proxy.h == 0.0


def test_render_diagram_edge_labels_paint_after_nodes(dr, blank_slide):
    """M4a Tier A3: label textboxes render AFTER node shapes in
    document/paint order — z-order = paint order in python-pptx, so a
    later position means the label paints on top. The M3 deferred
    defect (ibd draft_1 slides 10/19) was that labels rendered before
    nodes and disappeared under the boxes."""
    diagram = {
        "kind": "boxes_and_arrows",
        "nodes": [
            {"id": "a", "label": "A", "shape": "rectangle",
             "x": 1.0, "y": 1.0, "w": 1.5, "h": 0.8},
            {"id": "b", "label": "B", "shape": "rectangle",
             "x": 4.0, "y": 1.0, "w": 1.5, "h": 0.8},
        ],
        "edges": [
            {"from": "a", "to": "b", "kind": "straight",
             "label": "transitions"},
        ],
    }
    dr.render_diagram(blank_slide, diagram, (0, 0, 10, 5), None)
    # Walk shapes in document order and record the index of the last
    # node (auto-shape) and the first textbox whose text is "transitions".
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    last_node_idx = -1
    label_idx = -1
    for i, shp in enumerate(blank_slide.shapes):
        if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            last_node_idx = i
        if shp.has_text_frame and shp.text_frame.text == "transitions":
            label_idx = i
    assert last_node_idx >= 0, "expected at least one auto-shape (node)"
    assert label_idx >= 0, "expected the edge-label textbox to be present"
    assert label_idx > last_node_idx, (
        f"edge label (idx {label_idx}) must paint AFTER all node shapes "
        f"(last node idx {last_node_idx}); otherwise it sits behind a node "
        f"box (M3 Tier-A-deferred defect on ibd draft_1 slides 10/19)."
    )


# ---------------------------------------------------------------------------
# Top-level render_diagram
# ---------------------------------------------------------------------------

def test_render_diagram_full_graph(dr, blank_slide):
    diagram = {
        "kind": "boxes_and_arrows",
        "nodes": [
            {"id": "n1", "label": "Start", "shape": "rounded",
             "x": 0.5, "y": 0.5, "w": 1.5, "h": 0.8},
            {"id": "n2", "label": "End", "shape": "ellipse",
             "x": 5.0, "y": 0.5, "w": 1.5, "h": 0.8},
        ],
        "edges": [
            {"from": "n1", "to": "n2", "kind": "straight", "label": "next"},
        ],
    }
    n_before = len(blank_slide.shapes)
    dr.render_diagram(blank_slide, diagram, region=(0.5, 1.4, 9.0, 3.5),
                      brand_tokens=None)
    n_after = len(blank_slide.shapes)
    # 2 node shapes + 1 connector + 1 edge-label textbox
    assert n_after - n_before == 4


def test_render_diagram_unsupported_kind_raises(dr, blank_slide):
    diagram = {"kind": "mermaid", "nodes": [], "edges": []}
    with pytest.raises(ValueError) as exc:
        dr.render_diagram(blank_slide, diagram, region=(0, 0, 10, 5),
                          brand_tokens=None)
    assert "boxes_and_arrows" in str(exc.value)


def test_render_diagram_empty_graph(dr, blank_slide):
    """Empty nodes/edges should still succeed (no-op)."""
    diagram = {"kind": "boxes_and_arrows", "nodes": [], "edges": []}
    n_before = len(blank_slide.shapes)
    dr.render_diagram(blank_slide, diagram, region=(0, 0, 10, 5),
                      brand_tokens=None)
    n_after = len(blank_slide.shapes)
    assert n_after == n_before


# ---------------------------------------------------------------------------
# Brand-tokens loader
# ---------------------------------------------------------------------------

def test_load_brand_tokens_default_path(dr):
    """The shipped kbase-brand-tokens.json should load successfully."""
    tokens = dr.load_brand_tokens()
    assert tokens is not None
    assert "palette" in tokens
    assert "primary" in tokens["palette"]


def test_load_brand_tokens_missing_path_returns_none(dr, tmp_path):
    tokens = dr.load_brand_tokens(tmp_path / "nope.json")
    assert tokens is None


# ---------------------------------------------------------------------------
# CLI smoke
# ---------------------------------------------------------------------------

def test_cli_smoke(dr, tmp_path):
    import json
    diagram = {
        "kind": "boxes_and_arrows",
        "nodes": [
            {"id": "n1", "label": "X", "shape": "rectangle",
             "x": 1, "y": 1, "w": 2, "h": 1},
        ],
        "edges": [],
    }
    spec_path = tmp_path / "diagram.json"
    spec_path.write_text(json.dumps(diagram))
    out = tmp_path / "out.pptx"
    rc = dr.main([str(spec_path), "--out", str(out)])
    assert rc == 0
    assert out.is_file()

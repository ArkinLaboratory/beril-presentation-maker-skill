#!/usr/bin/env python3
"""assemble_pptx.py — render a validated slide_spec.json to a .pptx file.

Consumes slide_spec.json (per the contract in slide_spec.py) and writes
slides.pptx using the shipped KBase-branded master with 15 named layouts.

Pipeline:
  1. Load slide_spec.json.
  2. Validate against slide_spec.py (pre-flight; refuse to render invalid).
  3. Open the master template (kbase-presentation-master.pptx).
  4. For each slide in spec.slides:
       - Look up layout by name.
       - Add a new slide using that layout.
       - Dispatch to the layout-specific handler that fills placeholders +
         freeform shapes (figures, captions, footers).
       - Set speaker notes if present.
  5. Save to <out_path>.

PDF render (--format pdf) is delegated to LibreOffice (`soffice
--headless --convert-to pdf`). If LibreOffice is absent, the assembler
emits .pptx only and prints a clear message — pure-Python deps stay
the same. (D-016 pattern.)

Diagram rendering (workflow_diagram, cross_tenant_integration's optional
data_flow_diagram) is a STUB in this commit. The full implementation
ships in v0.1.0-extractors-c via diagram_render.py. The stub emits a
placeholder text box that says what the diagram should be — preserves
the slide structure and lets validators run end-to-end.

CLI:

    python3 assemble_pptx.py <slide_spec.json> --out <out.pptx>
                             [--master <override.pptx>]
                             [--format pptx|pdf]
                             [--strict]   # fail on any warning

Library:

    from assemble_pptx import assemble, AssemblyResult
    result = assemble("path/to/slide_spec.json", "out.pptx")
    print(result.warnings)

References:
  - SPEC §6 (slide vocabulary), §14.2 (slide_spec contract)
  - LAYOUT_FIXES in build_master.py for layout placeholder positions
  - slide_spec.py for the validator + types
"""
from __future__ import annotations

import argparse
import importlib.util
import json
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Master template discovery
# ---------------------------------------------------------------------------

# When this module is imported as a library, the master ships under
# `references/templates/kbase-presentation-master.pptx` next to it. When
# executed as a script, we resolve via the file path.
_THIS_DIR = Path(__file__).resolve().parent
_DEFAULT_MASTER = (_THIS_DIR.parent / "references" / "templates"
                   / "kbase-presentation-master.pptx")


def default_master_path() -> Path:
    """Path to the shipped master template."""
    if not _DEFAULT_MASTER.is_file():
        raise FileNotFoundError(
            f"Master template not found at {_DEFAULT_MASTER}. "
            f"Run tools/build_master.py to regenerate."
        )
    return _DEFAULT_MASTER


# ---------------------------------------------------------------------------
# slide_spec.py loader (sibling module)
# ---------------------------------------------------------------------------

def _load_slide_spec_module():
    """Load slide_spec.py from the same directory.

    We don't import via the package namespace because the skill/ tree
    ships as package_data and these tools are invoked as standalone
    scripts (per the orchestrator pattern in
    presentation_maker.sh — Phase 3).
    """
    path = _THIS_DIR / "slide_spec.py"
    spec = importlib.util.spec_from_file_location("slide_spec", path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"cannot load slide_spec from {path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["slide_spec"] = module  # for dataclass forward refs
    spec.loader.exec_module(module)
    return module


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------

@dataclass
class AssemblyResult:
    """Returned by assemble(). Holds the output path + any non-fatal warnings."""
    out_path: Path
    n_slides: int
    warnings: list[str] = field(default_factory=list)


class AssemblyError(Exception):
    """Raised when the assembler cannot proceed (e.g., master missing,
    layout not in master, validation failed)."""


# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

# v0.3.2: spec-layout → master-layout aliases. Some spec layouts reuse
# an existing master layout because their handler does its own freeform
# rendering (removes the body placeholder, adds shapes). The master
# layout only needs to provide the title placeholder + slide background.
#
# Aliasing avoids requiring a source-.potx update for every new spec
# layout we introduce.
SPEC_TO_MASTER_LAYOUT = {
    "data_table": "data_figure",   # title placeholder + body region; handler
                                   # removes body and renders its own table
}


def _get_layout_by_name(prs: Presentation, name: str):
    """Look up a layout by name in master 0. Raises AssemblyError if absent.

    Resolves spec-layout aliases (SPEC_TO_MASTER_LAYOUT) — e.g. `data_table`
    spec slides use the `data_figure` master layout under the hood.
    """
    resolved = SPEC_TO_MASTER_LAYOUT.get(name, name)
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == resolved:
            return layout
    available = sorted(l.name for l in prs.slide_masters[0].slide_layouts)
    if resolved != name:
        raise AssemblyError(
            f"layout '{name}' (aliased to '{resolved}') not in master template. "
            f"Available: {available}"
        )
    raise AssemblyError(
        f"layout '{name}' not in master template. Available: {available}"
    )


# OOXML namespaces — needed for slide-level bodyPr autofit fix.
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _ensure_slide_text_autofit(text_frame_element,
                                font_scale: int = 80000,
                                ln_spc_reduction: int = 20000,
                                anchor: str = "t") -> None:
    """Force normAutofit + anchor on a slide-level text frame's <a:bodyPr>.

    PowerPoint quirk: layout-level <a:normAutofit/> on a title placeholder
    governs initial layout but does NOT trigger runtime text-shrinking on
    rendered slides. The autofit MUST be set on the slide's own bodyPr
    for PowerPoint to compute fontScale at render time.

    The 2026-04-26 v0.1.1-visual master template fix (#53) wrote autofit
    on every layout's title placeholder, but verification showed slides
    1/5/7/10/12/15 still overran because the slide-level bodyPr was
    missing entirely (placeholders inherit text properties from layout
    BUT the bodyPr's autofit child is layout-init-only). This helper
    inserts the slide-level bodyPr autofit to make the inheritance
    actually apply at render.

    Pass an lxml Element (the <p:sp> for a placeholder OR the inner
    <p:txBody> directly). Looks up <a:bodyPr> under it, replaces any
    existing autofit child with <a:normAutofit fontScale="..." lnSpcReduction="..."/>,
    sets anchor attribute.
    """
    # If we got a <p:sp>, find <p:txBody>; if we got <p:txBody>, use it.
    if text_frame_element.tag.endswith("}sp"):
        tx_body = text_frame_element.find(f"{{{_PML_NS}}}txBody")
    else:
        tx_body = text_frame_element
    if tx_body is None:
        return
    body_pr = tx_body.find(f"{{{_DML_NS}}}bodyPr")
    if body_pr is None:
        return
    # Remove any existing autofit child(ren)
    for tag in ("normAutofit", "noAutofit", "spAutoFit"):
        for child in list(body_pr):
            if child.tag == f"{{{_DML_NS}}}{tag}":
                body_pr.remove(child)
    # Add normAutofit with explicit fontScale + lnSpcReduction
    from lxml import etree as _et
    af = _et.SubElement(body_pr, f"{{{_DML_NS}}}normAutofit")
    af.set("fontScale", str(font_scale))
    af.set("lnSpcReduction", str(ln_spc_reduction))
    # Set anchor (overflow-grows-down for top-anchored)
    body_pr.set("anchor", anchor)


def _set_title(slide, text: str) -> None:
    if not slide.shapes.title:
        return
    title_shape = slide.shapes.title
    title_shape.text = text
    # 2026-04-26 #63 fix: write slide-level normAutofit so PowerPoint
    # actually shrinks long titles at render. Layout-level autofit isn't
    # honored at render; the slide's own bodyPr must carry the autofit
    # element. Skip layouts where the master pins font size by design
    # (big_number, big_idea — those depend on prompt-side title-length
    # caps, not autofit, and forcing autofit here would override the
    # master's intentional pinning).
    layout_name = slide.slide_layout.name
    if layout_name not in ("big_number", "big_idea"):
        _ensure_slide_text_autofit(title_shape.element)


def _find_placeholder(slide, idx: int):
    """Return the placeholder with given idx, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_placeholder_text(slide, idx: int, text: str) -> bool:
    """Set a placeholder's text. Returns True if found, False otherwise."""
    ph = _find_placeholder(slide, idx)
    if ph is None:
        return False
    ph.text = text
    return True


def _set_placeholder_bullets(slide, idx: int, bullets: list[str]) -> bool:
    """Set placeholder text frame to the given list of bullets (one per
    paragraph). Returns True if placeholder found."""
    ph = _find_placeholder(slide, idx)
    if ph is None or not bullets:
        return False
    tf = ph.text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b
    return True


def _enable_normautofit(slide, idx: int,
                        *, font_scale: int = 80000,
                        ln_spc_reduction: int = 20000) -> bool:
    """Enable PowerPoint's normAutofit (text-shrinks-to-fit-shape) on a
    slide-level placeholder, with explicit fontScale + lnSpcReduction.

    Why this is needed: when assemble_pptx fills a placeholder via
    `tf.text = ...` + `tf.add_paragraph()`, python-pptx creates a
    slide-level <p:txBody> from scratch with no <a:bodyPr> autofit
    children. That overrides the layout's <a:normAutofit> setting from
    LAYOUT_FIXES. Without this helper, layouts whose master defines
    normAutofit on the body placeholder (methods_summary, qa_anticipated,
    references in v0.2.1) would render with no autofit at the slide
    level. Calling this AFTER `_set_placeholder_bullets` patches the
    slide-level body to match the layout's intent.

    PowerPoint requires explicit fontScale (defaults to 100000 = 100%
    when bare). 80000/20000 = 80% font shrink + 20% line-spacing
    reduction; same defaults as build_master's _apply_body_pr_change.

    Returns True if the placeholder + body_pr was found and patched.
    """
    ph = _find_placeholder(slide, idx)
    if ph is None:
        return False
    # Placeholder text bodies are <p:txBody> (presentationml namespace),
    # not <a:txBody>. The bodyPr child inside is <a:bodyPr> (drawingml).
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    txBody = ph._element.find(f".//{{{p_ns}}}txBody")
    if txBody is None:
        # Fall back to drawingml namespace for non-placeholder shapes
        txBody = ph._element.find(f".//{{{a_ns}}}txBody")
    if txBody is None:
        return False
    bodyPr = txBody.find(f"{{{a_ns}}}bodyPr")
    if bodyPr is None:
        return False
    # Remove any existing autofit children
    for tag in ("normAutofit", "noAutofit", "spAutoFit"):
        for child in list(bodyPr):
            if child.tag == f"{{{a_ns}}}{tag}":
                bodyPr.remove(child)
    # Insert normAutofit with explicit scale params
    from lxml import etree as _et
    autofit = _et.SubElement(bodyPr, f"{{{a_ns}}}normAutofit")
    autofit.set("fontScale", str(font_scale))
    autofit.set("lnSpcReduction", str(ln_spc_reduction))
    return True


def _remove_placeholder(slide, idx: int) -> bool:
    """Remove a placeholder from the slide entirely (not just clear its text).

    Setting `ph.text = ""` does NOT suppress an empty placeholder's "Click to
    add text" prompt in PowerPoint — the prompt is layout-defined and shows
    whenever the placeholder shape is on the slide. The only reliable way
    to hide it is to remove the placeholder element from the slide's spTree.

    Returns True if a placeholder with the given idx was removed, False if
    none was found.

    Use this for layouts where the body placeholder's region is occupied by
    a freeform figure (data_figure, concept_illustration) — the placeholder
    serves no purpose and its empty prompt visually overlaps the figure.
    """
    ph = _find_placeholder(slide, idx)
    if ph is None:
        return False
    sptree = slide.shapes._spTree
    sptree.remove(ph._element)
    return True


def _remove_decorative_banner(slide) -> bool:
    """Remove the FIRST non-placeholder <p:sp> from the slide's spTree.

    Used by `_fill_big_idea`'s centered-assertion mode (v0.2.2): the
    layout inherits a decorative top banner from LAYOUT_FIXES["big_idea"]
    that's appropriate when supporting_graphic is present, but visually
    weak when the slide is a pure assertion. Removing the banner lets
    the title float on a clean slide.

    The decorative banner is typically the first non-placeholder shape
    in document order. python-pptx's shape iteration walks the spTree;
    the first <p:sp> that has no <p:nvSpPr>/<p:nvPr>/<p:ph> child is
    the banner. Returns True if removed, False if no banner was found.
    """
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sptree = slide.shapes._spTree
    for child in list(sptree):
        # Match <p:sp> only (not <p:pic>, <p:grpSp>, etc.)
        if not child.tag.endswith("}sp"):
            continue
        # Skip placeholders (have a <p:ph> descendant under <p:nvSpPr>)
        ph = child.find(f"{{{p_ns}}}nvSpPr/{{{p_ns}}}nvPr/{{{p_ns}}}ph")
        if ph is not None:
            continue
        sptree.remove(child)
        return True
    return False


def _reposition_placeholder_to_center(slide, *, idx: int,
                                      left_in: float, top_in: float,
                                      width_in: float, height_in: float) -> bool:
    """Move + resize a placeholder to absolute slide-relative position.

    Used by `_fill_big_idea`'s centered-assertion mode to override the
    layout-defined title placement (which is in the top accent banner)
    and put the title at slide center. Returns True on success, False
    if no placeholder with the given idx was found.

    `idx=0` is the title placeholder convention.
    """
    ph = _find_placeholder(slide, idx)
    if ph is None:
        return False
    ph.left = Inches(left_in)
    ph.top = Inches(top_in)
    ph.width = Inches(width_in)
    ph.height = Inches(height_in)
    return True


def _set_title_font_size(slide, *, font_pt: int) -> bool:
    """Set the font size on the title placeholder's runs (idx=0).

    Used by `_fill_big_idea`'s centered-assertion mode to bump the title
    from the LAYOUT_FIXES default (36pt for banner-mode) to 48pt for the
    pull-quote treatment. Walks all paragraphs/runs in the title's text
    frame and sets each run's font size; this is necessary because the
    layout-level def_rpr.sz doesn't propagate when the slide-level
    placeholder body is rebuilt at fill time.
    """
    ph = _find_placeholder(slide, 0)
    if ph is None or not ph.has_text_frame:
        return False
    for paragraph in ph.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_pt)
    return True


def _enable_normautofit_on_title(slide,
                                 *, font_scale: int = 80000,
                                 ln_spc_reduction: int = 20000) -> bool:
    """Apply normAutofit to the title placeholder's body_pr.

    Same shape as `_enable_normautofit` but targets idx=0 (title) and
    used by big_idea's centered-assertion mode so 200+-char claims
    shrink to fit the centered band rather than overflowing.
    """
    return _enable_normautofit(slide, 0,
                               font_scale=font_scale,
                               ln_spc_reduction=ln_spc_reduction)


def _is_tbd_placeholder(text: str) -> bool:
    """True if `text` is a TBD-style placeholder that should be hidden
    from the rendered deck.

    Used by acknowledgments rendering (v0.2.2) to filter orchestrator-
    template defaults that leak through when a draft hasn't been edited
    by the user. Recognizes:
      - exact "TBD"
      - "TBD - populated by production orchestrator"
      - "TBD - <anything>"
      - case-insensitive variants
    """
    if not isinstance(text, str):
        return False
    stripped = text.strip().lower()
    if stripped == "tbd":
        return True
    if stripped.startswith("tbd -") or stripped.startswith("tbd—"):
        return True
    if stripped.startswith("tbd:"):
        return True
    return False


def _set_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def _add_picture(slide, image_path: Path,
                 left_in: float, top_in: float,
                 width_in: float, height_in: float):
    """Add a picture at the given inches-position. Returns the picture shape."""
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in), Inches(top_in),
        width=Inches(width_in), height=Inches(height_in),
    )


def _add_textbox(slide, text: str,
                 left_in: float, top_in: float,
                 width_in: float, height_in: float,
                 *, font_size_pt: int = 14, bold: bool = False,
                 color_rgb: tuple[int, int, int] | None = None,
                 align_center: bool = False,
                 word_wrap: bool = False,
                 auto_size: bool = False,
                 shrink_to_fit: bool = False):
    """Add a freeform text box at the given position.

    word_wrap=True enables long-text line wrapping inside the box (default
    is False — single-line, which silently truncates / spills off edge).
    Use word_wrap=True for caption / source / subtitle textboxes that
    take production-realistic content; use False for short labels where
    fixed-width is the design intent.

    auto_size=True enables python-pptx's MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    so the textbox grows to fit its content. Combined with word_wrap=True
    this gives a "fit to content" textbox. Default False (fixed size).

    shrink_to_fit=True enables python-pptx's MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    so the FONT shrinks to keep text inside a fixed-size box. Use this
    for caption boxes where overflow into adjacent bands (e.g., the
    KBase brand strip at y=5.00) is unacceptable. Mutually exclusive
    with auto_size; passing both honors shrink_to_fit (v0.3.5: belt-
    and-suspenders for data_figure captions; the slide_spec validator
    is the primary cap at 280 chars).
    """
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = tb.text_frame
    tf.text = text
    if word_wrap:
        tf.word_wrap = True
    if shrink_to_fit:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif auto_size:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = tf.paragraphs[0]
    if align_center:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        if color_rgb:
            run.font.color.rgb = RGBColor(*color_rgb)
    return tb


# ---------------------------------------------------------------------------
# Figure/asset position tables (per layout, when the layout has one figure
# region beyond placeholders). Inches.
# ---------------------------------------------------------------------------

# Slide is 10.0 × 5.625 inches (16:9). Logos sit at y=5.00–5.55 (bottom-
# right) so figure regions stay clear of that band. Body placeholders are
# removed via _remove_placeholder for layouts where the body region is
# occupied by a freeform figure — the placeholder's "Click to add text"
# prompt would otherwise show through.
FIGURE_REGIONS = {
    # Bullets in body placeholder (left half via runtime resize); figure
    # on right. 2026-04-29 (v0.2.2): figure H 3.50 → 3.15 in to leave
    # 0.40 in band (4.45..4.85) for the caption above the logo strip
    # at 5.00. Live failure draft_10 slide 18: caption ran into logos
    # because figure ended at y=4.80, logos at 5.00, leaving only 0.20 in.
    "claim_evidence":       (5.30, 1.30, 4.50, 3.15),
    # Title is in the accent banner (round 3); supporting graphic fills
    # the body area below banner, ABOVE the logos at y=5.00.
    "big_idea":             (1.00, 1.10, 8.00, 3.85),
    # Body placeholder removed; figure fills former body region.
    # 2026-05-03 (v0.3.2.8): figure H 3.10 → 2.85 (top 1.40 → 1.30, bottom
    # 4.50 → 4.15) to give a 0.65-in caption band + 0.15-in data_source
    # band. Live failure draft_2 slide 8: revise-loop produced ~410-char
    # caption that auto-size'd into the data_source's y=4.82 anchor;
    # texts overlapped. Geometry now budgets ~3 wrapped lines at 12pt
    # for the caption + 1 line at 10pt for data_source.
    "data_figure":          (0.50, 1.30, 9.00, 2.85),
    # Body placeholder removed; image on the right of the slide.
    "concept_illustration": (5.30, 1.30, 4.50, 3.70),
}

# Caption/footer regions (overlaid on slide)
CAPTION_BAND = (0.30, 4.95, 9.40, 0.40)   # bottom 0.4-inch strip
CITATION_BAND = (0.30, 5.20, 9.40, 0.30)  # very-bottom strip
AI_DISCLOSURE_BAND = (0.30, 5.30, 9.40, 0.20)  # 8pt graphite-gray

GRAPHITE_GRAY_RGB = (157, 146, 135)  # KBase secondary palette

# v0.3.2: KBase brand palette (full hex for data_table styling)
KBASE_BLUE_RGB   = (0x00, 0x7D, 0xC3)   # #007DC3 — table header bg, links
KBASE_GREEN_RGB  = (0x5E, 0x97, 0x32)   # #5E9732 — secondary accents
KBASE_ORANGE_RGB = (0xF7, 0x8E, 0x1E)   # #F78E1E — highlight rows
ROW_BAND_RGB     = (0xF2, 0xF2, 0xF2)   # alternating row band (light gray)
WHITE_RGB        = (0xFF, 0xFF, 0xFF)
TABLE_TEXT_RGB   = (0x33, 0x33, 0x33)   # dark gray, slightly off-black


# ---------------------------------------------------------------------------
# Path resolution
# ---------------------------------------------------------------------------

def _derive_project_dir(draft_dir: Path) -> Path | None:
    """Walk up from draft_dir to find the project_dir.

    Two layouts to handle:

    v0.3.0 (legacy):  projects/<id>/talks/draft_N/  → walk up 2 levels
    v0.3.1+:          projects/<id>/talks/draft_N/working/  → walk up 3
                       projects/<id>/talks/draft_N/  → walk up 2 (same as v0.3.0)

    The first form is hit when callers pass `slide_spec_path.parent` and the
    spec lives at draft_N/working/slide_spec.json (v0.3.1+ layout). The
    second form is hit when callers pass the draft_N/ directory directly
    (e.g., from the orchestrator side passing `OUTDIR`).

    Returns the project_dir Path or None if the standard structure isn't
    present.
    """
    # v0.3.1+ shape: caller passed draft_N/working/, walk up 3 levels.
    if draft_dir.name == "working" and draft_dir.parent.parent.name == "talks":
        candidate = draft_dir.parent.parent.parent
        if candidate.is_dir():
            return candidate
    # v0.3.0 / direct-draft shape: caller passed draft_N/, walk up 2 levels.
    if draft_dir.parent.name == "talks" and draft_dir.parent.parent.is_dir():
        return draft_dir.parent.parent
    return None


def _derive_actual_draft_dir(maybe_working: Path) -> Path:
    """Given the path the caller passes as `draft_dir`, return the actual
    draft_N/ directory.

    v0.3.1+: the caller may pass `draft_N/working/` (the parent of
    slide_spec.json); we want `draft_N/` for relative-path resolution
    against the rest of the layout. Detect by checking the leaf name.
    """
    if maybe_working.name == "working" and maybe_working.parent.is_dir():
        return maybe_working.parent
    return maybe_working


def _resolve_asset_path(rel_or_abs: str, draft_dir: Path,
                       warnings: list[str], where: str) -> Path | None:
    """Resolve a figure / image / logo path.

    Lookup order for relative paths (2026-04-27 fix #77):
      1. Resolve against draft_dir (same as before)
      2. Fall back to project_dir (../../ from draft_dir if it's
         under projects/<id>/talks/draft_N/) — this is the standard
         layout, and curate_figures.md emits paths relative to
         project_dir (e.g., 'figures/fig01.png'), not draft_dir.

    Returns the first existing path; None if neither exists.
    """
    candidate_path = Path(rel_or_abs)
    if candidate_path.is_absolute():
        if candidate_path.is_file():
            return candidate_path
        warnings.append(
            f"{where}: asset not found at {candidate_path} "
            f"(slide will render with a placeholder note)"
        )
        return None

    # Try draft_dir first (legacy behavior)
    cand_draft = draft_dir / candidate_path
    if cand_draft.is_file():
        return cand_draft

    # Fall back to project_dir (where curate_figures.md paths live)
    project_dir = _derive_project_dir(draft_dir)
    if project_dir is not None:
        cand_project = project_dir / candidate_path
        if cand_project.is_file():
            return cand_project

    # Not found anywhere — warn and return None
    tried = [str(cand_draft)]
    if project_dir is not None:
        tried.append(str(project_dir / candidate_path))
    warnings.append(
        f"{where}: asset not found (tried: {' | '.join(tried)}) "
        f"(slide will render with a placeholder note)"
    )
    return None


# ---------------------------------------------------------------------------
# Per-layout handlers
# ---------------------------------------------------------------------------
#
# Each handler signature:
#   _fill_<layout>(slide, content, draft_dir, warnings) -> None
#
# Handlers MAY add freeform shapes. They MUST not assume layouts have
# placeholders beyond what build_master.py guarantees (TITLE on every
# layout; BODY idx 1 on most; BODY idx 2 only on two_column_compare).

def _fill_title(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    parts: list[str] = []
    if content.get("subtitle"):
        parts.append(content["subtitle"])
    presenter_line = content["presenter"]
    if content.get("affiliation"):
        presenter_line = f"{presenter_line} · {content['affiliation']}"
    parts.append(presenter_line)
    if content.get("venue"):
        parts.append(content["venue"])
    parts.append(content["date"])
    _set_placeholder_text(slide, 1, "\n".join(parts))
    # 2026-04-26 #63 fix: subtitle placeholder also needs slide-level
    # autofit. With the v0.1.1 title-slide stub fix, the subtitle now
    # carries the throughline punchline (often 200+ chars) — the layout's
    # autofit doesn't trigger at render unless the slide's own bodyPr
    # has it.
    subtitle_ph = _find_placeholder(slide, 1)
    if subtitle_ph is not None:
        _ensure_slide_text_autofit(subtitle_ph.element)


def _fill_section_divider(slide, content, draft_dir, warnings):
    _set_title(slide, content["punchline"])
    if content.get("substory_number"):
        # Add a small footer with the substory number
        _add_textbox(slide, f"Substory {content['substory_number']}",
                     0.30, 5.10, 4.00, 0.30,
                     font_size_pt=12, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_big_idea(slide, content, draft_dir, warnings):
    """big_idea has two render modes:

    1. **Centered assertion (default).** When `supporting_graphic` is
       absent, the slide is a pull-quote — title centered vertically and
       horizontally on a clean (non-banner) slide, similar in feel to
       section_divider but without the substory label. Used for opening
       claims and key transitions.

    2. **Banner + image (legacy).** When `supporting_graphic` is present,
       the original LAYOUT_FIXES design lights up: title in a top accent
       banner, image below in the body region. This mode is forward-
       compatible with v0.3's `ai_image_prompt.v1` — when generated
       supporting graphics ship, big_idea slides can opt into them and
       this branch handles render.

    2026-04-29 (v0.2.2 fix, draft_10 slide 2):
      Mode 1 was added because mode 2 was the only render path and the
      LLM rarely emits supporting_graphic, leaving every big_idea slide
      with title-at-top + empty body. The dual-mode pattern matches
      `_fill_claim_evidence` (with-figure / without-figure split).
    """
    if content.get("supporting_graphic"):
        # Mode 2: banner + image (original layout)
        _set_title(slide, content["title"])
        path = _resolve_asset_path(content["supporting_graphic"], draft_dir,
                                   warnings, "big_idea.supporting_graphic")
        if path:
            _add_picture(slide, path, *FIGURE_REGIONS["big_idea"])
        return

    # Mode 1: centered assertion (default)
    # Remove the decorative top banner (it's the first non-placeholder
    # <p:sp> on the slide, inherited from the layout's LAYOUT_FIXES setup).
    _remove_decorative_banner(slide)
    # Reposition the title placeholder to slide-center, mirroring the
    # section_divider geometry minus the substory label.
    _reposition_placeholder_to_center(slide, idx=0,
                                      left_in=0.0, top_in=1.94,
                                      width_in=10.0, height_in=1.27)
    _set_title(slide, content["title"])
    # 48pt vs section_divider's 40pt — big_idea is opening-claim
    # emphasis; section_divider is transition cadence.
    _set_title_font_size(slide, font_pt=48)
    # normAutofit so 200+-char claims shrink instead of overflowing.
    _enable_normautofit_on_title(slide)


def _fill_big_number(slide, content, draft_dir, warnings):
    """big_number's TITLE placeholder is repositioned by build_master.py
    LAYOUT_FIXES to be the huge centered area. We place the headline + a
    smaller subtitle below as a separate text box (the placeholder font is
    66pt bold, and we want subtitle smaller)."""
    _set_title(slide, content["headline"])
    # Subtitle in a separate textbox below the title region.
    # Title region per LAYOUT_FIXES: off (660902,923453) ext (7840301,3286408)
    # = (0.72, 1.01, 8.57 × 3.59 in). Title bottom = 4.60; logos start
    # at 5.00. Only 0.40 in available for subtitle.
    #
    # 2026-04-28 (v0.2.1 fix, draft_9 slide 18):
    #   Subtitle was 20pt × H=0.40 without word_wrap. Production text
    #   (64 chars) needs ~1.5 lines at 20pt — overflowed off the right
    #   edge. Cannot grow box height (would overlap title or logos),
    #   so reduce font to 16pt: 64 chars / 8.57in × (1.6 × 72/16) cpi
    #   = ~1.04 lines, fits 0.40 in tall slot. word_wrap=True handles
    #   any longer subtitles via wrapping (slide_compose prompt should
    #   cap subtitle ≤45 chars in v0.3+).
    _add_textbox(slide, content["subtitle"],
                 0.72, 4.65, 8.57, 0.40,
                 font_size_pt=16, align_center=True,
                 word_wrap=True)
    if content.get("sub_pointer"):
        _add_textbox(slide, content["sub_pointer"],
                     0.72, 5.05, 8.57, 0.30,
                     font_size_pt=12, color_rgb=GRAPHITE_GRAY_RGB,
                     align_center=True)
    if content.get("source_footer"):
        _add_textbox(slide, content["source_footer"],
                     0.30, 5.35, 9.40, 0.20,
                     font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_claim_evidence(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    bullets = content["bullets"]
    if content.get("figure"):
        # Bullets in narrower body (left half), figure on right.
        #
        # 2026-04-28 (v0.2.1 fix, draft_9 slide 8):
        #   Previously the body placeholder kept its full ~9.32 in width
        #   while the figure went at L=5.30 W=4.50 — collision area
        #   ~15 in² of bullets text overlapping the figure. Now: when a
        #   figure is present, resize the body placeholder to the left
        #   half (~4.86 in wide ending at 5.20 in) before filling. Figure
        #   sits in FIGURE_REGIONS["claim_evidence"] = (5.30, 1.30, ...)
        #   on the right, no overlap.
        ph = _find_placeholder(slide, 1)
        if ph is not None:
            ph.left = Inches(0.34)
            ph.top = Inches(1.30)
            ph.width = Inches(4.86)  # ends at 5.20; figure starts at 5.30 (0.10 gap)
            ph.height = Inches(3.50)
        _set_placeholder_bullets(slide, 1, bullets)
        path = _resolve_asset_path(content["figure"], draft_dir, warnings,
                                    "claim_evidence.figure")
        if path:
            _add_picture(slide, path, *FIGURE_REGIONS["claim_evidence"])
            # Caption band below figure (only if figure rendered).
            #
            # 2026-04-29 (v0.2.2): drop auto_size (was growing box past
            # bottom logos); fix word_wrap to actually take. Live test
            # on draft_10 slide 18 showed caption truncated with "..."
            # — auto_size=SHAPE_TO_FIT_TEXT overrides word_wrap in
            # python-pptx (auto-size assumes single-line). Geometry:
            # figure H 3.50 → 3.15 (FIGURE_REGIONS update above) so
            # figure ends at y=4.45; caption sits at 4.50..4.85 in the
            # cleared 0.40 in band; logos start at 5.00.
            _add_textbox(slide, content["figure_caption"],
                         5.30, 4.50, 4.50, 0.35,
                         font_size_pt=11, color_rgb=GRAPHITE_GRAY_RGB,
                         word_wrap=True)
    else:
        _set_placeholder_bullets(slide, 1, bullets)
    if content.get("citations"):
        # short-form citation footer at bottom
        cite_text = " · ".join(content["citations"])
        _add_textbox(slide, cite_text, *CITATION_BAND,
                     font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_two_column_compare(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # Idx 1 = left column, Idx 2 = right column (per master inspection).
    for idx, col_key in [(1, "left_col"), (2, "right_col")]:
        col_title = content[f"{col_key}_title"]
        col_content = content[f"{col_key}_content"]
        ph = _find_placeholder(slide, idx)
        if ph is None:
            warnings.append(
                f"two_column_compare: layout missing BODY idx {idx}"
            )
            continue
        tf = ph.text_frame
        # Column-title as first paragraph (bold)
        tf.text = col_title
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].font.bold = True
        # Then the content
        if isinstance(col_content, str):
            p = tf.add_paragraph()
            p.text = col_content
        else:
            for line in col_content:
                p = tf.add_paragraph()
                p.text = line
    # 2026-04-29 (v0.2.2): normAutofit on both column placeholders. Live
    # test on draft_10 slide 19 showed the right column's last bullet
    # ("scores 0.875 for CRISPRi analysis") overflowing into the bottom
    # logo region. The two_column_compare body H ≈ 3.48 in starting at
    # T ≈ 0.98 — bottom 4.46, only 0.54 in clearance to logos at 5.00.
    # 4-5 bullets at 18pt blow through. normAutofit shrinks per-column.
    _enable_normautofit(slide, 1)
    _enable_normautofit(slide, 2)


def _fill_data_figure(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # Remove the body placeholder entirely — the body region is occupied
    # by the freeform figure + caption + data_source. Setting ph.text = ""
    # is insufficient: the layout-defined "Click to add text" prompt
    # still shows in PowerPoint when the placeholder is present.
    _remove_placeholder(slide, 1)
    path = _resolve_asset_path(content["figure"], draft_dir, warnings,
                                "data_figure.figure")
    if path:
        _add_picture(slide, path, *FIGURE_REGIONS["data_figure"])
    # Caption + data source budget under the figure (which now ends at
    # y=4.15 per FIGURE_REGIONS["data_figure"] update in v0.3.2.8):
    #   y=4.15..4.80 → caption (0.65 in, ~3 wrapped lines at 12pt)
    #   y=4.83..4.98 → data_source (0.15 in, ~1 line at 10pt)
    #   y=5.00       → logo strip
    #
    # auto_size=False is critical: it prevents the caption box from
    # growing downward and overlapping the data_source band when the
    # caption is long.
    #
    # shrink_to_fit=True (v0.3.5): MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    # shrinks the FONT (not the box) so any caption that would otherwise
    # spill below y=4.83 instead renders at a smaller point size and
    # stays inside the 0.65-in band. This is belt-and-suspenders against
    # the slide_spec validator's 280-char cap (DATA_FIGURE_CAPTION_MAX_CHARS):
    # if the validator is bypassed, an old draft is re-assembled, or
    # specific word lengths cause weird wrapping at edge cases, the
    # render side still keeps text out of the brand strip at y=5.00.
    #
    # History:
    # 2026-04-28 (v0.2.1 fix #2, draft_9 walk): word_wrap=True introduced.
    # 2026-05-03 (v0.3.2.8, draft_2 slide 8): drop auto_size, grow caption
    #   H to 0.65, shrink figure region. Live failure: revise-loop produced
    #   ~410-char caption; auto_size grew caption box past data_source's
    #   y=4.82 anchor, producing visual overlap.
    # 2026-05-05 (v0.3.5, gene_function_ecological_agora draft_1 slides
    #   21+23): even with auto_size=False the 410-char caption still
    #   spilled past the 0.65-in box bottom into the data_source / brand
    #   strip (text overflow renders outside box bounds when no auto_size
    #   is set). Layered fix: prompt cap + validator hard-fail + render
    #   shrink_to_fit. See slide_compose.v1.md / revise_slide.v1.md /
    #   slide_spec.DATA_FIGURE_CAPTION_MAX_CHARS.
    _add_textbox(slide, content["caption"], 0.50, 4.18, 9.00, 0.65,
                 font_size_pt=12, color_rgb=GRAPHITE_GRAY_RGB,
                 word_wrap=True, shrink_to_fit=True)
    if content.get("data_source"):
        _add_textbox(slide, content["data_source"],
                     0.50, 4.83, 9.00, 0.15,
                     font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB,
                     word_wrap=True)


def _fill_data_table(slide, content, draft_dir, warnings):
    """Render a `data_table` slide. v0.3.2.

    Layout zones:
      title:          slide title placeholder (top)
      table:          centered horizontally, immediately below title
                      (y=1.10), height bounded so it clears the bottom
                      logo strip at y=5.00.
      caption:        small textbox immediately below the table
      footnote:       even smaller textbox at slide bottom (above logos)

    Brand styling:
      header row:     KBase blue (#007DC3) bg, white text, bold
      odd data rows:  light gray band (#F2F2F2) bg
      even data rows: white bg
      highlight rows: KBase orange (#F78E1E) bg with white text (overrides
                      the alternating-band coloring)

    Validator (slide_spec._check_data_table) caps rows ≤ 12 and
    columns ≤ 6. The renderer assumes these caps held; widths and font
    sizes are tuned for that range.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    _set_title(slide, content["title"])
    _remove_placeholder(slide, 1)  # body region taken by the table

    columns = content["columns"]
    rows = content["rows"]
    n_cols = len(columns)
    n_rows = len(rows) + 1  # +1 for header
    highlight_rows = set(content.get("highlight_rows") or [])

    # Geometry — table fills the body region. Same horizontal envelope as
    # other body-region layouts (0.50..9.50, 9.0in wide). Height adapts to
    # row count so the table doesn't crash through the logo strip at y=5.00.
    table_left = Inches(0.50)
    table_top = Inches(1.10)
    table_width = Inches(9.00)
    # Max usable vertical band: 1.10 → 4.50 = 3.40 in. Leaves 0.50 in
    # below the table for caption + footnote, then logos at 5.00.
    max_table_h_in = 3.40
    # Header row gets a slightly taller line (16pt + padding); data rows
    # use 14pt. Cap row height so 12-row tables fit; smaller tables get
    # taller rows for legibility.
    target_row_h_in = min(0.34, max_table_h_in / n_rows)
    table_height = Inches(target_row_h_in * n_rows)

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=n_cols,
        left=table_left,
        top=table_top,
        width=table_width,
        height=table_height,
    )
    table = shape.table

    # Column widths: equal-fraction by default. Tighten the first column
    # if it's clearly an identifier (e.g., gene names) — heuristic: if the
    # first column header is short and the rest are wider, narrow it.
    # Skip this heuristic for now (equal widths render acceptably for
    # 2-6 cols). Future: per-column width hints from the spec.

    # --- Header row ---
    for j, header_text in enumerate(columns):
        cell = table.cell(0, j)
        _set_table_cell(
            cell, header_text,
            bg_rgb=KBASE_BLUE_RGB,
            text_rgb=WHITE_RGB,
            bold=True,
            font_pt=12,
            align=PP_ALIGN.LEFT,
        )

    # --- Data rows ---
    for i, row in enumerate(rows):
        is_highlight = i in highlight_rows
        is_odd_band = (i % 2 == 1)  # 0-based: row 0 white, row 1 banded, ...
        if is_highlight:
            bg = KBASE_ORANGE_RGB
            text = WHITE_RGB
            bold = True
        elif is_odd_band:
            bg = ROW_BAND_RGB
            text = TABLE_TEXT_RGB
            bold = False
        else:
            bg = WHITE_RGB
            text = TABLE_TEXT_RGB
            bold = False
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            _set_table_cell(
                cell, cell_text,
                bg_rgb=bg,
                text_rgb=text,
                bold=bold,
                font_pt=11,
                align=PP_ALIGN.LEFT,
            )

    # --- Caption (below table) ---
    table_bottom_in = 1.10 + target_row_h_in * n_rows
    caption = content.get("caption")
    if caption:
        _add_textbox(
            slide, caption,
            0.50, table_bottom_in + 0.05, 9.00, 0.30,
            font_size_pt=11, color_rgb=GRAPHITE_GRAY_RGB,
            word_wrap=True,
        )

    # --- Footnote (very bottom, above logo strip) ---
    footnote = content.get("footnote") or content.get("data_source")
    if footnote:
        # Place at y=4.80, height 0.18, above logos at 5.00.
        _add_textbox(
            slide, footnote,
            0.50, 4.80, 9.00, 0.18,
            font_size_pt=9, color_rgb=GRAPHITE_GRAY_RGB,
            word_wrap=True,
        )


def _set_table_cell(cell, text: str, *, bg_rgb, text_rgb,
                    bold: bool = False, font_pt: int = 11,
                    align=None):
    """Fill a python-pptx table cell with text + brand styling.

    Sets:
      - cell.text (single run)
      - cell.fill.solid() + .fore_color.rgb = bg_rgb
      - run.font.color.rgb = text_rgb
      - run.font.bold = bold
      - run.font.size = font_pt
      - paragraph alignment (if provided)

    Note: python-pptx requires assigning to cell.text first to get a
    single paragraph + run, then we mutate that run's font.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    cell.text = text or ""
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*bg_rgb)

    para = cell.text_frame.paragraphs[0]
    if align is not None:
        para.alignment = align
    if not para.runs:
        # Empty cell → no run to style. python-pptx auto-creates a run on
        # text assignment; defensive check.
        return
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*text_rgb)


def _fill_workflow_diagram(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # Remove the body placeholder — its region is occupied by the
    # rendered diagram (native python-pptx shapes).
    _remove_placeholder(slide, 1)
    # Render the diagram via diagram_render.py (v0.1.0-extractors-c).
    # Region: same as the body placeholder bounds, slightly inset.
    diagram = content["diagram"]
    try:
        from importlib import util as _util
        _spec = _util.spec_from_file_location(
            "diagram_render", _THIS_DIR / "diagram_render.py")
        _dr = _util.module_from_spec(_spec)
        # Register in sys.modules BEFORE exec_module so the @dataclass
        # decorator inside diagram_render can resolve cls.__module__
        # via sys.modules.get(...).__dict__ (Python's dataclass machinery
        # walks sys.modules, and unregistered modules return None).
        # Same pattern used in _load_slide_spec_module above.
        sys.modules["diagram_render"] = _dr
        _spec.loader.exec_module(_dr)
        # Diagram region: most of the body, leaving room for step caption + footer.
        region = (0.50, 1.30, 9.00, 3.10)
        # Brand tokens for color resolution
        tokens = _dr.load_brand_tokens()
        _dr.render_diagram(slide, diagram, region, tokens)
    except Exception as e:  # noqa: BLE001
        warnings.append(
            f"workflow_diagram on slide id={getattr(slide, 'slide_id', '?')}: "
            f"diagram render failed ({e}); slide will lack the diagram."
        )
    # Step caption band below — 3 boxes side-by-side under the diagram
    # (one per step_caption entry). 2026-04-27 fix #55: previous version
    # joined all 3 captions with double-space and crammed them into a
    # single 9.0x0.30in textbox, which ran off the right edge for any
    # caption longer than ~30 chars (live failure draft_2 slide 11).
    # Splitting into 3 columns gives each caption ~3.0in width and
    # 0.50in height for word-wrap.
    captions = content["step_caption"]
    n_captions = len(captions)
    if n_captions > 0:
        column_w = 9.0 / n_captions
        for i, caption in enumerate(captions):
            x = 0.50 + i * column_w
            # 2026-04-29 (v0.2.2): captions need word_wrap=True. Production
            # captions run 60-100 chars per step; without wrapping they
            # render as single overlong lines that bleed across column
            # boundaries (visible draft_10 slide 9: three captions visually
            # overlapping at the bottom). word_wrap=True lets each caption
            # wrap within its 3-in column. H=0.55 fits 2 lines at 11pt;
            # captions >100 chars will be clipped — content-side cap (~80
            # chars per step) is a v0.3 prompt iteration.
            _add_textbox(slide, caption,
                         x + 0.05, 4.50, column_w - 0.10, 0.55,
                         font_size_pt=11, color_rgb=GRAPHITE_GRAY_RGB,
                         word_wrap=True)
    if content.get("tool_version_footer"):
        _add_textbox(slide, content["tool_version_footer"],
                     0.30, 5.30, 9.40, 0.20,
                     font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_methods_summary(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    _set_placeholder_bullets(slide, 1, content["bullets"])
    # 2026-04-28 (v0.2.1): normAutofit at slide level so dense methods
    # content (5-7 bullets, 600-800 chars) shrinks to fit. The layout-
    # level body_pr autofit from LAYOUT_FIXES gets overridden when
    # python-pptx creates a fresh slide-level txBody during fill.
    _enable_normautofit(slide, 1)

    # 2026-04-26 fix #59: render tools_versions as a footer band.
    # Prior version dropped this structured data entirely, leaving only
    # a hardcoded "(see speaker notes)" hint. The schema field is named
    # `version` and the slide_compose prompt now (T2.7) requires real
    # version pins; even before that fix, surfacing whatever the model
    # produced is more honest than dropping it.
    tools_versions = content.get("tools_versions") or []
    if tools_versions:
        # Format as comma-separated "Tool ver" pairs:
        #   "RAST 2.0 · fastp 0.23 · DRAM 1.4"
        formatted = " · ".join(
            f"{tv.get('tool', '?')} {tv.get('version', '?')}"
            for tv in tools_versions
            if isinstance(tv, dict)
        )
        if formatted:
            _add_textbox(slide, formatted,
                         0.30, 5.18, 9.40, 0.28,
                         font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB)
            return  # tools_versions footer takes the speaker-notes hint slot

    # No tools_versions populated — fall back to the speaker-notes hint
    if content.get("see_notes_footer", True):
        _add_textbox(slide, "(see speaker notes for full detail)",
                     0.30, 5.20, 9.40, 0.30,
                     font_size_pt=10, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_concept_illustration(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # The body placeholder on the left is narrowed (per Adam's
    # build_master fix) but v0.1 schema has no body-text field for this
    # layout — remove the placeholder so its empty prompt doesn't show
    # alongside the AI-generated image.
    _remove_placeholder(slide, 1)
    path = _resolve_asset_path(content["image_path"], draft_dir, warnings,
                                "concept_illustration.image_path")
    if path:
        _add_picture(slide, path, *FIGURE_REGIONS["concept_illustration"])
    if content.get("caption"):
        _add_textbox(slide, content["caption"],
                     5.30, 5.05, 4.50, 0.25,
                     font_size_pt=11, color_rgb=GRAPHITE_GRAY_RGB)
    if content.get("ai_disclosure_footer", True):
        _add_textbox(slide, "AI-generated illustration", *AI_DISCLOSURE_BAND,
                     font_size_pt=8, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_cross_tenant_integration(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # Build the body content: tenant list, K-BERDL DBs, sibling projects
    lines: list[str] = []
    if content.get("no_signal_fallback"):
        lines.append("All data sourced from a single tenant.")
        lines.append("This project did not integrate across tenants.")
    else:
        if content.get("tenant_list"):
            lines.append(f"Tenants: {', '.join(content['tenant_list'])}")
        if content.get("kberdl_db_list"):
            lines.append(f"K-BERDL databases: {', '.join(content['kberdl_db_list'])}")
        if content.get("sibling_project_refs"):
            for ref in content["sibling_project_refs"]:
                lines.append(
                    f"From {ref['project_id']}: {ref['what_was_leveraged']}"
                )
    if lines:
        _set_placeholder_bullets(slide, 1, lines)
    # data_flow_diagram (optional) — same diagram-render stub story
    if content.get("data_flow_diagram"):
        warnings.append(
            f"cross_tenant_integration on slide id={slide.slide_id}: "
            f"data_flow_diagram rendering stubbed; v0.1.0-extractors-c."
        )


def _fill_implications(slide, content, draft_dir, warnings):
    _set_title(slide, content["title"])
    # bullets is list of {claim, evidence_pointer}
    lines = [
        f"{b['claim']}\n   ↪ {b['evidence_pointer']}"
        for b in content["bullets"]
    ]
    _set_placeholder_bullets(slide, 1, lines)


def _fill_acknowledgments(slide, content, draft_dir, warnings):
    # Title is hard-coded "Acknowledgments" (exempt from punchline rule per SPEC §6.1)
    _set_title(slide, "Acknowledgments")
    contributors = list(content["contributors"])
    # 2026-04-29 (v0.2.2): TBD soft-default. Live test draft_10 slide 25
    # showed "TBD - populated by production orchestrator" + "TBD" leaking
    # through. Acknowledgments are user-fill (no LLM authoring); when the
    # spec hasn't been edited post-draft, replace template placeholders
    # with a polite, presentable line so the slide doesn't ship with
    # "TBD" visible. The user can edit the spec or open the .pptx and
    # replace before presenting.
    if all(_is_tbd_placeholder(c) for c in contributors):
        contributors = ["Acknowledgments to be added before presentation."]
    else:
        # Filter out individual TBD entries; keep real contributors.
        contributors = [c for c in contributors if not _is_tbd_placeholder(c)]
        if not contributors:
            contributors = ["Acknowledgments to be added before presentation."]
    lines = contributors
    if content.get("tenant_attribution"):
        lines.append(content["tenant_attribution"])
    if content.get("code_repo_url"):
        lines.append(f"Code: {content['code_repo_url']}")
    _set_placeholder_bullets(slide, 1, lines)
    # funder_logos: render as a horizontal strip at the bottom
    if content.get("funder_logos"):
        n_logos = len(content["funder_logos"])
        slide_w_in = 10.0
        margin = 0.30
        gap = 0.20
        max_w = (slide_w_in - 2 * margin - (n_logos - 1) * gap) / max(n_logos, 1)
        max_h = 0.60
        for i, logo_path in enumerate(content["funder_logos"]):
            path = _resolve_asset_path(logo_path, draft_dir, warnings,
                                        f"acknowledgments.funder_logos[{i}]")
            if path:
                left = margin + i * (max_w + gap)
                _add_picture(slide, path, left, 4.95, max_w, max_h)


def _fill_references(slide, content, draft_dir, warnings):
    # Title is hard-coded "References" (exempt from punchline rule per SPEC §6.1)
    _set_title(slide, "References")
    _set_placeholder_bullets(slide, 1, content["refs_short"])
    # 2026-04-28 (v0.2.1): normAutofit at slide level. References at 8
    # entries × ~134 chars wraps to ~17 lines @ 18pt against a 12-line
    # cap; autofit shrinks the list to fit. Same rationale as
    # methods_summary above.
    _enable_normautofit(slide, 1)
    # AI-disclosure footer (always emitted)
    disclosure = content.get(
        "ai_disclosure",
        "Slides drafted with beril-presentation-maker; full bibliography "
        "in speaker notes."
    )
    _add_textbox(slide, disclosure, 0.30, 5.30, 9.40, 0.30,
                 font_size_pt=8, color_rgb=GRAPHITE_GRAY_RGB)


def _fill_qa_anticipated(slide, content, draft_dir, warnings):
    _set_title(slide, content["question"])
    body_lines = [content["answer_summary"]]
    if content.get("answer_detail"):
        body_lines.append("")
        body_lines.append(content["answer_detail"])
    body_lines.append("")
    body_lines.append(f"↪ {content['evidence_pointer']}")
    _set_placeholder_bullets(slide, 1, body_lines)
    # 2026-04-29 (v0.2.2): tighter normAutofit specifically for qa_anticipated.
    # v0.2.1's 80% fontScale wasn't aggressive enough — math: 80% × 18pt
    # × 1.2 leading × 4.00 in body × 9.32 in width ≈ 1400 chars capacity,
    # but production qa_prep produces 2000+ char 5-paragraph answers (~50%
    # overflow). 60% scale gives ~2000 chars at 10.8pt — readable at
    # projection. methods/refs stay at 80% (their content fits within that).
    # Companion fix: qa_prep.v1.md word-budget cap (~600 chars) lands in
    # v0.3+ as prompt iteration.
    _enable_normautofit(slide, 1, font_scale=60000, ln_spc_reduction=20000)


# Dispatcher
LAYOUT_HANDLERS = {
    "title":                    _fill_title,
    "section_divider":          _fill_section_divider,
    "big_idea":                 _fill_big_idea,
    "big_number":               _fill_big_number,
    "claim_evidence":           _fill_claim_evidence,
    "two_column_compare":       _fill_two_column_compare,
    "data_figure":              _fill_data_figure,
    "data_table":               _fill_data_table,
    "workflow_diagram":         _fill_workflow_diagram,
    "methods_summary":          _fill_methods_summary,
    "concept_illustration":     _fill_concept_illustration,
    "cross_tenant_integration": _fill_cross_tenant_integration,
    "implications":             _fill_implications,
    "acknowledgments":          _fill_acknowledgments,
    "references":               _fill_references,
    "qa_anticipated":           _fill_qa_anticipated,
}


# ---------------------------------------------------------------------------
# Top-level assemble entry point
# ---------------------------------------------------------------------------

def _build_poster_spec_from_slide_spec(spec: dict, draft_dir: Path):
    """Map a slide_spec.json (in poster mode) to a PosterSpec for poster_fill.
    Walks the slide list and pulls the relevant content fields.
    """
    from importlib import util as _util
    _spec = _util.spec_from_file_location("poster_fill", _THIS_DIR / "poster_fill.py")
    _pf = _util.module_from_spec(_spec)
    # Register before exec_module so @dataclass inside poster_fill works
    # (Python dataclasses walk sys.modules to resolve cls.__module__).
    # Same gotcha that bit diagram_render's load.
    sys.modules["poster_fill"] = _pf
    _spec.loader.exec_module(_pf)

    title = ""
    authors = ""
    affiliation = ""
    tl_dr = spec.get("throughline", {}).get("punchline", "") or ""
    methods_summary: list[str] = []
    figures: list[_pf.PosterFigure] = []
    cross_tenant_summary = ""
    implications: list[str] = []
    references_short: list[str] = []
    acknowledgments = ""
    funding = ""

    for slide in spec.get("slides", []):
        layout = slide.get("layout")
        content = slide.get("content", {}) or {}
        if layout == "title":
            title = content.get("title", "")
            authors = content.get("presenter", "")
            affiliation = content.get("affiliation", "")
        elif layout == "methods_summary":
            methods_summary = list(content.get("bullets", []) or [])
        elif layout in ("data_figure", "claim_evidence", "concept_illustration"):
            for key in ("figure", "image_path", "supporting_graphic"):
                if key in content and content[key]:
                    cap = (content.get("figure_caption")
                           or content.get("caption")
                           or "")
                    figures.append(_pf.PosterFigure(path=content[key],
                                                    caption=cap))
                    break
        elif layout == "cross_tenant_integration":
            parts = []
            if content.get("tenant_list"):
                parts.append(f"Tenants: {', '.join(content['tenant_list'])}")
            if content.get("kberdl_db_list"):
                parts.append(f"DBs: {', '.join(content['kberdl_db_list'])}")
            if content.get("sibling_project_refs"):
                refs = ", ".join(r["project_id"] for r in content["sibling_project_refs"])
                parts.append(f"Sibling projects: {refs}")
            if not parts and content.get("no_signal_fallback"):
                parts.append("All data sourced from a single tenant; no cross-tenant integration.")
            cross_tenant_summary = "; ".join(parts) if parts else (content.get("title", ""))
        elif layout == "implications":
            implications = [b.get("claim", "") if isinstance(b, dict) else str(b)
                            for b in content.get("bullets", []) or []]
        elif layout == "references":
            references_short = list(content.get("refs_short", []) or [])
        elif layout == "acknowledgments":
            ack_parts = list(content.get("contributors", []) or [])
            acknowledgments = " · ".join(ack_parts)
            funding = content.get("tenant_attribution", "") or ""

    return _pf, _pf.PosterSpec(
        title=title, authors=authors, affiliation=affiliation,
        tl_dr=tl_dr, methods_summary=methods_summary, figures=figures[:4],
        cross_tenant_summary=cross_tenant_summary,
        implications=implications, references_short=references_short,
        acknowledgments=acknowledgments, funding=funding,
    )


def assemble(slide_spec_path: str | Path,
             out_path: str | Path,
             *,
             master_path: str | Path | None = None,
             strict: bool = False) -> AssemblyResult:
    """Assemble a slide_spec.json into a .pptx file.

    Args:
      slide_spec_path: path to slide_spec.json (the contract source)
      out_path:        path to write slides.pptx
      master_path:     override the shipped master template (default:
                       references/templates/kbase-presentation-master.pptx)
      strict:          if True, raise on any warning (e.g., missing figure
                       file). Default False — warnings collected but rendered.

    Returns:
      AssemblyResult with `out_path`, `n_slides`, and `warnings`.

    Raises:
      AssemblyError on validation failure or missing master.
    """
    slide_spec_path = Path(slide_spec_path).resolve()
    out_path = Path(out_path).resolve()
    if not slide_spec_path.is_file():
        raise AssemblyError(f"slide_spec.json not found: {slide_spec_path}")

    spec = json.loads(slide_spec_path.read_text(encoding="utf-8"))

    # Pre-flight validation
    ss = _load_slide_spec_module()
    issues = ss.validate_slide_spec(spec)
    if issues:
        raise AssemblyError(
            f"slide_spec.json failed schema validation "
            f"({len(issues)} issue(s)):\n  "
            + "\n  ".join(i.format() for i in issues[:20])
            + ("\n  ..." if len(issues) > 20 else "")
        )

    # Poster mode dispatches to poster_fill (separate render path per
    # SPEC §12 / D-013). Posters skip the slide-by-slide handler loop.
    mode = spec.get("mode", "")
    if mode in ("poster-h", "poster-v"):
        # v0.3.2.1: pass the actual draft_N/, not draft_N/working/.
        actual_draft_dir = _derive_actual_draft_dir(slide_spec_path.parent)
        pf_module, poster_spec = _build_poster_spec_from_slide_spec(
            spec, actual_draft_dir)
        orientation = "horizontal" if mode == "poster-h" else "vertical"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            pf_module.fill_poster(
                poster_spec, out_path,
                orientation=orientation,
                draft_dir=actual_draft_dir,
            )
        except (FileNotFoundError, ValueError) as e:
            raise AssemblyError(f"poster_fill failed: {e}") from e
        return AssemblyResult(
            out_path=out_path, n_slides=1,
            warnings=[],
        )

    # Master (talk modes)
    master = Path(master_path) if master_path else default_master_path()
    if not master.is_file():
        raise AssemblyError(f"master template not found: {master}")
    prs = Presentation(master)

    # v0.3.2.1: slide_spec_path is `draft_N/working/slide_spec.json` in the
    # v0.3.1+ layout. Handlers expect the actual draft_N/ for figure path
    # resolution against project_dir = draft_N/../../. Walk up if we see the
    # working/ subdir.
    draft_dir = _derive_actual_draft_dir(slide_spec_path.parent)
    warnings: list[str] = []

    for slide_data in spec["slides"]:
        layout_name = slide_data["layout"]
        layout = _get_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        handler = LAYOUT_HANDLERS.get(layout_name)
        if handler is None:
            # Should be unreachable — slide_spec validator pins layout to vocab
            raise AssemblyError(
                f"no handler registered for layout {layout_name!r}"
            )
        handler(slide, slide_data["content"], draft_dir, warnings)

        if "speaker_notes" in slide_data:
            _set_speaker_notes(slide, slide_data["speaker_notes"])

    if strict and warnings:
        raise AssemblyError(
            f"--strict mode: {len(warnings)} warning(s):\n  "
            + "\n  ".join(warnings)
        )

    # Save
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    return AssemblyResult(
        out_path=out_path,
        n_slides=len(spec["slides"]),
        warnings=warnings,
    )


# ---------------------------------------------------------------------------
# PDF rendering (LibreOffice)
# ---------------------------------------------------------------------------

def render_pdf(pptx_path: Path) -> Path | None:
    """Render <pptx_path> to <pptx_path with .pdf suffix> via LibreOffice.

    Returns the .pdf path on success. Returns None if soffice is not on
    PATH (assembler emits a clear message; pptx-only output is still valid).
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        return None
    out_dir = pptx_path.parent
    cmd = [soffice, "--headless", "--convert-to", "pdf",
           "--outdir", str(out_dir), str(pptx_path)]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=120)
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        print(f"PDF render failed: {e}", file=sys.stderr)
        return None
    pdf = pptx_path.with_suffix(".pdf")
    return pdf if pdf.is_file() else None


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="assemble_pptx",
        description="Render slide_spec.json to slides.pptx (and optionally PDF).",
    )
    parser.add_argument("slide_spec",
                        help="Path to slide_spec.json (the contract input).")
    parser.add_argument("--out", required=True,
                        help="Output .pptx path.")
    parser.add_argument("--master",
                        help="Override master template path "
                             "(default: shipped kbase-presentation-master.pptx)")
    parser.add_argument("--format", choices=["pptx", "pdf"], default="pptx",
                        help="Output format. pdf invokes soffice; falls back "
                             "to pptx if LibreOffice not on PATH.")
    parser.add_argument("--strict", action="store_true",
                        help="Fail on any warning (e.g., missing figure file).")
    args = parser.parse_args(argv)

    try:
        result = assemble(args.slide_spec, args.out,
                          master_path=args.master, strict=args.strict)
    except AssemblyError as e:
        print(f"assemble_pptx: {e}", file=sys.stderr)
        return 2

    print(f"wrote {result.out_path} ({result.n_slides} slide(s))",
          file=sys.stderr)
    if result.warnings:
        print(f"{len(result.warnings)} warning(s):", file=sys.stderr)
        for w in result.warnings:
            print(f"  - {w}", file=sys.stderr)

    if args.format == "pdf":
        pdf = render_pdf(result.out_path)
        if pdf is None:
            print("PDF render unavailable (LibreOffice not found). "
                  "Open .pptx in PowerPoint/Keynote and export to PDF "
                  "manually.", file=sys.stderr)
        else:
            print(f"wrote {pdf}", file=sys.stderr)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

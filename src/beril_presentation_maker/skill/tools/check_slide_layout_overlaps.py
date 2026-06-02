#!/usr/bin/env python3
"""check_slide_layout_overlaps.py — deterministic bounding-box
overlap detector (v0.8.0 Tier G.10-A).

Pure geometry over the rendered .pptx: walks every shape on every
slide via python-pptx, computes pairwise overlap with a small
padding tolerance, and emits findings in adversarial-review shape.
Replaces the vision-LLM raster judgment of `element_overlap` /
`container_breach` / `footer_or_title_collision` — those signals
are now mechanical (cheaper, faster, deterministic).

Visual-QA's remaining job (post-G.10-A): `headline_body_mismatch`,
`illegible_scale`, AI-image-spoiler patterns. Cases where reading
the rendered raster genuinely adds signal vs. inspecting geometry.

Finding kinds:

- `text_box_overlap` — two text-bearing shapes overlap.
- `image_text_overlap` — an image shape overlaps a text-bearing
  shape (with footer-watermark allow-list per layout template).
- `footer_title_collision` — title placeholder overlaps a body
  placeholder (chrome eats narrative space).
- `container_breach` — a shape's bounding box extends past the
  slide's drawable area (P0; renders cut off in projection).

Severity:

- `container_breach`: P0 (the rendered slide is mechanically broken).
- All others: P1 (soft-warning; revise loop should resolve).

Padding tolerance: 36000 EMU ≈ 0.04 inch by default. Configurable
via `--pad-emu`. The threshold accounts for normal rendering
slop without false-positiving on aligned-but-not-overlapping
shapes.

Layout allow-list: some layout templates have intentional overlap
zones (e.g., a footer band sitting behind a watermark). The
allow-list (per-layout, indexed by template name) suppresses
findings whose overlap is entirely inside an allowed zone.
Default allow-list ships with the empty set (no templates have
known intentional overlap as of v0.8.0); operators extend the
YAML beside layout templates if needed.

Output: `audit/layout_overlaps.json` (`layout-overlaps.v1`)
plus optional markdown summary at `audit/layout_overlaps.md`.

Cascade integration: `review_cascade.py::_read_layout_overlaps`
reads the audit JSON (read-if-present pattern; cascade never
invokes this script). Findings lift into cascade Tier-1 as
`layout_overlaps:<finding-kind>` with the severity stamped here.

Test coverage: `tests/unit/test_check_slide_layout_overlaps.py`.

Refs: D-098 (the prior belt-and-suspenders pattern this mirrors);
V0_8_PUNCH_LIST.md Tier G.10 (the workstream this implements;
coupled with G.10-B autofit-commit + G.10-C overflow-finding).
"""
from __future__ import annotations

import argparse
import dataclasses
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

SCHEMA_VERSION = "layout-overlaps.v1"


# ---------------------------------------------------------------------------
# Geometry primitives (EMU; 914400 EMU = 1 inch)
# ---------------------------------------------------------------------------

_DEFAULT_PAD_EMU = 36000  # ~0.04 inch tolerance


@dataclass(frozen=True)
class Rect:
    """Axis-aligned bounding box in EMU."""
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    def area(self) -> int:
        return max(0, self.width) * max(0, self.height)


def overlaps(a: Rect, b: Rect, pad_emu: int = _DEFAULT_PAD_EMU) -> bool:
    """True if rectangles a and b overlap (after shrinking each by
    pad_emu/2 on every edge — equivalently, requiring at least
    pad_emu of true intersection on each axis).

    A pure separating-axis check: rectangles do NOT overlap iff
    one is entirely to the left/right or above/below the other.
    """
    if a.right - pad_emu <= b.left or b.right - pad_emu <= a.left:
        return False
    if a.bottom - pad_emu <= b.top or b.bottom - pad_emu <= a.top:
        return False
    return True


def intersection(a: Rect, b: Rect) -> Optional[Rect]:
    """Intersection rectangle, or None if disjoint (without
    padding — geometric truth)."""
    left = max(a.left, b.left)
    top = max(a.top, b.top)
    right = min(a.right, b.right)
    bottom = min(a.bottom, b.bottom)
    if right <= left or bottom <= top:
        return None
    return Rect(left, top, right - left, bottom - top)


def contained_in(inner: Rect, outer: Rect,
                 pad_emu: int = _DEFAULT_PAD_EMU) -> bool:
    """True if `inner` is entirely inside `outer` (within the
    padding tolerance on each side)."""
    return (inner.left >= outer.left - pad_emu
            and inner.top >= outer.top - pad_emu
            and inner.right <= outer.right + pad_emu
            and inner.bottom <= outer.bottom + pad_emu)


# ---------------------------------------------------------------------------
# Shape classification — text vs. image vs. chrome
# ---------------------------------------------------------------------------

# python-pptx shape_type codes (MSO_SHAPE_TYPE)
_SHAPE_PICTURE = 13       # MSO_SHAPE_TYPE.PICTURE
_SHAPE_PLACEHOLDER = 14   # MSO_SHAPE_TYPE.PLACEHOLDER
_SHAPE_TEXT_BOX = 17      # MSO_SHAPE_TYPE.TEXT_BOX


def shape_role(shape) -> str:
    """Classify a python-pptx shape as 'image', 'text', 'title', or
    'other'. The classification drives finding-kind routing
    (image-text vs. text-text vs. footer/title collision)."""
    stype = shape.shape_type
    if stype == _SHAPE_PICTURE:
        return "image"
    name = (getattr(shape, "name", "") or "").lower()
    # Title placeholders are special — they're chrome, not body
    # narrative. Detected by python-pptx placeholder_format.idx == 0
    # OR by name containing 'title'.
    if stype == _SHAPE_PLACEHOLDER:
        ph_format = getattr(shape, "placeholder_format", None)
        if ph_format is not None:
            idx = getattr(ph_format, "idx", None)
            if idx == 0:
                return "title"
        if "title" in name:
            return "title"
        return "text"
    if stype == _SHAPE_TEXT_BOX:
        if "title" in name:
            return "title"
        return "text"
    return "other"


def shape_rect(shape) -> Optional[Rect]:
    """Build a Rect from a python-pptx shape's left/top/width/height.
    Returns None if any coordinate is missing (defensive: some
    placeholders inherit position from layout and may not have
    explicit coords until rendered)."""
    try:
        left, top, width, height = (
            shape.left, shape.top, shape.width, shape.height
        )
    except AttributeError:
        return None
    if any(v is None for v in (left, top, width, height)):
        return None
    return Rect(int(left), int(top), int(width), int(height))


# ---------------------------------------------------------------------------
# Findings
# ---------------------------------------------------------------------------

@dataclass
class OverlapFinding:
    """One layout-overlap finding."""
    kind: str               # text_box_overlap | image_text_overlap |
                            # footer_title_collision | container_breach
    severity: str           # "P0" | "P1"
    slide_id: int           # 1-indexed slide number in the rendered deck
    layout_name: str        # python-pptx slide layout name
    shape_a: str            # name of the first shape
    shape_b: str            # name of the second shape (or "" for breach)
    overlap_area_emu: int   # area of the overlap rect (or breach-out area)
    overlap_fraction: float  # overlap area / min(a.area, b.area); 0..1
    message: str            # operator-visible one-line summary

    def to_dict(self) -> dict:
        return dataclasses.asdict(self)


@dataclass
class OverlapReport:
    """Aggregated per-deck report."""
    schema_version: str = SCHEMA_VERSION
    pptx_path: str = ""
    pad_emu: int = _DEFAULT_PAD_EMU
    findings: list[OverlapFinding] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "schema_version": self.schema_version,
            "pptx_path": self.pptx_path,
            "pad_emu": self.pad_emu,
            "findings": [f.to_dict() for f in self.findings],
        }


# ---------------------------------------------------------------------------
# Per-slide detection
# ---------------------------------------------------------------------------

def _name(shape) -> str:
    return str(getattr(shape, "name", "") or "<unnamed>")


def detect_container_breaches(
    shapes_and_rects: list[tuple[object, Rect]],
    slide_id: int,
    layout_name: str,
    slide_width: int,
    slide_height: int,
    pad_emu: int,
) -> list[OverlapFinding]:
    """Find shapes whose bounding box extends past the slide canvas.
    P0 because the rendered slide is mechanically broken: portions
    project as off-screen / cut off."""
    findings: list[OverlapFinding] = []
    canvas = Rect(0, 0, slide_width, slide_height)
    for shape, rect in shapes_and_rects:
        if contained_in(rect, canvas, pad_emu):
            continue
        # Compute breach magnitude as the area of `rect` outside the
        # canvas (intersection of rect minus canvas).
        inside = intersection(rect, canvas)
        inside_area = inside.area() if inside else 0
        breach_area = rect.area() - inside_area
        if breach_area <= 0:
            continue
        findings.append(OverlapFinding(
            kind="container_breach",
            severity="P0",
            slide_id=slide_id,
            layout_name=layout_name,
            shape_a=_name(shape),
            shape_b="",
            overlap_area_emu=breach_area,
            overlap_fraction=breach_area / max(rect.area(), 1),
            message=(
                f"shape {_name(shape)!r} extends past the slide canvas "
                f"({breach_area} EMU outside; "
                f"{breach_area / max(rect.area(), 1):.1%} of shape area)"
            ),
        ))
    return findings


def detect_pairwise_overlaps(
    shapes_and_rects: list[tuple[object, Rect]],
    slide_id: int,
    layout_name: str,
    pad_emu: int,
    allow_zones: list[Rect],
) -> list[OverlapFinding]:
    """Find pairs of shapes that overlap. Classifies the finding
    kind by role pair (text-text, image-text, title-text)."""
    findings: list[OverlapFinding] = []
    n = len(shapes_and_rects)
    for i in range(n):
        sa, ra = shapes_and_rects[i]
        for j in range(i + 1, n):
            sb, rb = shapes_and_rects[j]
            if not overlaps(ra, rb, pad_emu):
                continue
            isect = intersection(ra, rb)
            if isect is None:
                continue
            # Allow-list: if the overlap is entirely inside any
            # allowed zone, suppress.
            if any(contained_in(isect, z, pad_emu) for z in allow_zones):
                continue
            role_a = shape_role(sa)
            role_b = shape_role(sb)
            kind = _classify_pair(role_a, role_b)
            if kind is None:
                continue  # other/other — not interesting
            min_area = max(min(ra.area(), rb.area()), 1)
            fraction = isect.area() / min_area
            findings.append(OverlapFinding(
                kind=kind,
                severity="P1",
                slide_id=slide_id,
                layout_name=layout_name,
                shape_a=_name(sa),
                shape_b=_name(sb),
                overlap_area_emu=isect.area(),
                overlap_fraction=fraction,
                message=(
                    f"{kind}: {_name(sa)!r} ({role_a}) overlaps "
                    f"{_name(sb)!r} ({role_b}); "
                    f"{fraction:.1%} of smaller shape's area"
                ),
            ))
    return findings


def _classify_pair(role_a: str, role_b: str) -> Optional[str]:
    """Map a (role_a, role_b) pair to a finding kind. Returns None
    if the pair is uninteresting (other/other, etc.)."""
    roles = frozenset({role_a, role_b})
    if "image" in roles and ("text" in roles or "title" in roles):
        return "image_text_overlap"
    if "title" in roles and "text" in roles:
        return "footer_title_collision"
    if roles == frozenset({"text"}):
        return "text_box_overlap"
    if roles == frozenset({"title"}):
        # Two title-named shapes — rare; treat as collision.
        return "footer_title_collision"
    return None


# ---------------------------------------------------------------------------
# Allow-list loading
# ---------------------------------------------------------------------------

def load_allow_list(path: Optional[Path]) -> dict[str, list[Rect]]:
    """Load per-layout intentional-overlap zones from a JSON file.
    Shape:

        {
          "layout-name": [
            {"left": <emu>, "top": <emu>, "width": <emu>, "height": <emu>},
            ...
          ],
          ...
        }

    Returns an empty dict if the file doesn't exist (default
    behavior — no templates have known intentional overlap as of
    v0.8.0)."""
    if path is None or not path.is_file():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    if not isinstance(data, dict):
        return {}
    out: dict[str, list[Rect]] = {}
    for layout_name, zones in data.items():
        if not isinstance(zones, list):
            continue
        rects = []
        for z in zones:
            if not isinstance(z, dict):
                continue
            try:
                rects.append(Rect(
                    int(z["left"]), int(z["top"]),
                    int(z["width"]), int(z["height"])
                ))
            except (KeyError, TypeError, ValueError):
                continue
        if rects:
            out[layout_name] = rects
    return out


# ---------------------------------------------------------------------------
# Top-level check
# ---------------------------------------------------------------------------

def check_pptx(
    pptx_path: Path,
    pad_emu: int = _DEFAULT_PAD_EMU,
    allow_list: Optional[dict[str, list[Rect]]] = None,
) -> OverlapReport:
    """Walk every slide in the .pptx; emit findings. Lazy-imports
    python-pptx so import doesn't fail in environments where the
    binary isn't available (cascade-read pattern: the JSON is
    consumed, the checker isn't always invoked).
    """
    from pptx import Presentation  # local import per docstring
    prs = Presentation(str(pptx_path))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    allow_list = allow_list or {}

    report = OverlapReport(pptx_path=str(pptx_path), pad_emu=pad_emu)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_name = str(getattr(slide.slide_layout, "name", "") or "")
        # Collect every shape with a real bounding box.
        shapes_and_rects: list[tuple[object, Rect]] = []
        for sh in slide.shapes:
            rect = shape_rect(sh)
            if rect is None:
                continue
            shapes_and_rects.append((sh, rect))

        # 1. Container breaches first (P0, independent of pairs).
        report.findings.extend(detect_container_breaches(
            shapes_and_rects, slide_idx, layout_name,
            slide_width, slide_height, pad_emu,
        ))

        # 2. Pairwise overlaps.
        zones = allow_list.get(layout_name, [])
        report.findings.extend(detect_pairwise_overlaps(
            shapes_and_rects, slide_idx, layout_name, pad_emu, zones,
        ))

    return report


# ---------------------------------------------------------------------------
# Markdown summary
# ---------------------------------------------------------------------------

def render_markdown(report: OverlapReport) -> str:
    """Operator-friendly summary."""
    lines = ["# Layout-overlap report", ""]
    lines.append(f"**Source pptx:** `{report.pptx_path}`")
    lines.append(f"**Padding tolerance:** {report.pad_emu} EMU "
                 f"({report.pad_emu / 914400:.3f} in)")
    lines.append(f"**Total findings:** {len(report.findings)}")
    lines.append("")
    if not report.findings:
        lines.append("_No layout overlaps or container breaches detected._")
        return "\n".join(lines) + "\n"
    # Group by severity
    p0 = [f for f in report.findings if f.severity == "P0"]
    p1 = [f for f in report.findings if f.severity == "P1"]
    if p0:
        lines.append(f"## P0 — container breaches ({len(p0)})")
        lines.append("")
        for f in p0:
            lines.append(f"- slide {f.slide_id} ({f.layout_name}): "
                         f"{f.message}")
        lines.append("")
    if p1:
        lines.append(f"## P1 — overlaps ({len(p1)})")
        lines.append("")
        for f in p1:
            lines.append(f"- slide {f.slide_id} ({f.layout_name}): "
                         f"{f.message}")
        lines.append("")
    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv: Optional[list[str]] = None) -> int:
    ap = argparse.ArgumentParser(
        prog="check_slide_layout_overlaps.py",
        description=(
            "Deterministic bounding-box overlap detector for "
            "rendered .pptx decks (v0.8.0 Tier G.10-A)."
        ),
    )
    ap.add_argument(
        "draft_dir", type=Path,
        help="Draft directory containing deliverable/draft.pptx "
             "(audit/ lives alongside).",
    )
    ap.add_argument(
        "--pptx", type=Path, default=None,
        help="Override pptx path (default: <draft_dir>/deliverable/draft.pptx)",
    )
    ap.add_argument(
        "--out", type=Path, default=None,
        help="Override JSON output path (default: "
             "<draft_dir>/audit/layout_overlaps.json)",
    )
    ap.add_argument(
        "--out-md", type=Path, default=None,
        help="Override markdown output path (default: "
             "<draft_dir>/audit/layout_overlaps.md)",
    )
    ap.add_argument(
        "--pad-emu", type=int, default=_DEFAULT_PAD_EMU,
        help=f"Padding tolerance in EMU (default: {_DEFAULT_PAD_EMU} "
             f"≈ 0.04 inch). Lower → more sensitive.",
    )
    ap.add_argument(
        "--allow-list", type=Path, default=None,
        help="Optional JSON file with per-layout intentional-overlap "
             "zones (suppresses findings inside these zones).",
    )
    args = ap.parse_args(argv)

    pptx_path = (args.pptx
                 if args.pptx else args.draft_dir / "deliverable" / "draft.pptx")
    if not pptx_path.is_file():
        print(f"Error: pptx not found at {pptx_path}", file=sys.stderr)
        return 2

    audit_dir = args.draft_dir / "audit"
    audit_dir.mkdir(parents=True, exist_ok=True)
    out_json = args.out if args.out else audit_dir / "layout_overlaps.json"
    out_md = args.out_md if args.out_md else audit_dir / "layout_overlaps.md"

    allow_list = load_allow_list(args.allow_list)

    try:
        report = check_pptx(pptx_path, pad_emu=args.pad_emu,
                            allow_list=allow_list)
    except ImportError:
        print("Error: python-pptx not installed; cannot inspect deck",
              file=sys.stderr)
        return 2

    out_json.write_text(
        json.dumps(report.to_dict(), indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    out_md.write_text(render_markdown(report), encoding="utf-8")

    n_p0 = sum(1 for f in report.findings if f.severity == "P0")
    n_p1 = sum(1 for f in report.findings if f.severity == "P1")
    print(f"Layout overlap check: {len(report.findings)} finding(s) "
          f"({n_p0} P0, {n_p1} P1) → {out_json}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())

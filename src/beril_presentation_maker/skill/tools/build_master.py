#!/usr/bin/env python3
"""build_master.py — derive kbase-presentation-master.pptx from user-supplied .potx.

This is a build-time script (NOT runtime). It is run when:
- The KBase brand changes and Adam supplies a refreshed .potx.
- The slide-shape vocabulary in SPEC §6 / DECISIONS D-008 is updated.

Inputs (relative to repo root):
    reference/master-template-source/KBase 2026 and beyond.potx
        (gitignored; user-supplied; never redistributed)

Outputs:
    src/beril_presentation_maker/skill/references/templates/kbase-presentation-master.pptx
    src/beril_presentation_maker/skill/references/kbase-brand-tokens.json

Pipeline:
    1. Validate source .potx exists.
    2. Convert content-type .potx → .pptx (python-pptx can't open .potx directly).
    3. Open in python-pptx; verify expected source structure.
    4. Rename 15 chosen layouts in slide master 0 to our named vocabulary.
    5. Remove the 17 unused layouts from slide master 0 (XML + parts + rels).
    6. Remove slide master 1 entirely (Google Slides round-trip duplicate).
    7. Save derived master.
    8. Emit kbase-brand-tokens.json from canonical Style Guide values
       (NOT extracted from .potx — Style Guide is the source of truth per
       reference/kbase-style-extract.md).
    9. Print a build report (layout names confirmed, sizes, etc.).

Idempotent: running twice from the same .potx produces byte-identical output
(modulo XML attribute ordering; tested via unzipped-content-equality).

References:
    - SPEC §6, §14.1 — slide-shape vocabulary, master template
    - DECISIONS D-007, D-008 — derived master, closed vocabulary
    - reference/master-template-source-notes.md — full derivation contract
    - reference/kbase-style-extract.md — brand tokens canonical source
"""
from __future__ import annotations

import argparse
import hashlib
import io
import json
import shutil
import sys
import zipfile
from pathlib import Path
from typing import Iterable

# python-pptx
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Layout mapping: source-layout-name → our named-vocabulary name.
#
# Mapping was chosen by inspecting the .potx (see
# reference/master-template-source-notes.md §2). Each source layout's
# placeholder shape was matched to the closest fit for our 15 named
# vocabulary entries. Indices in the table below are advisory; the mapping
# matches by source name (more stable than ordinal index).
# ---------------------------------------------------------------------------

LAYOUT_RENAMES: dict[str, str] = {
    # Source name (from KBase 2026 .potx)         → our SPEC §6 name
    "TITLE_1_2":                                  "title",
    "CUSTOM_4":                                   "section_divider",
    "CUSTOM_1":                                   "big_idea",
    "CUSTOM_3_3":                                 "big_number",
    "TITLE_AND_BODY_1_1_1_1_1_1_1_1_1":           "claim_evidence",
    "CUSTOM_1_1":                                 "two_column_compare",
    "TITLE_AND_BODY_1":                           "data_figure",
    "TITLE_AND_BODY_1_1_1_1":                     "workflow_diagram",
    "TITLE_AND_BODY_1_1_1_1_1_1_1":               "methods_summary",
    "TITLE_AND_BODY_1_2":                         "concept_illustration",
    "TITLE_AND_BODY_1_1_1_1_1_1_1_1_2":           "cross_tenant_integration",
    "TITLE_AND_BODY_1_1_1":                       "implications",
    "CUSTOM_5":                                   "acknowledgments",
    "TITLE_AND_BODY_1_1_2":                       "references",
    "TITLE_AND_BODY_1_1_1_1_1_2_1":               "qa_anticipated",
}

# The 15 expected names after renaming. Used for validation.
NAMED_VOCABULARY: tuple[str, ...] = tuple(LAYOUT_RENAMES.values())
assert len(NAMED_VOCABULARY) == 15, "vocabulary must be exactly 15 layouts"
assert len(set(NAMED_VOCABULARY)) == 15, "vocabulary names must be unique"


# ---------------------------------------------------------------------------
# Layout fixes — applied AFTER renaming.
#
# Adam reviewed the v0.1 derived master visually (2026-04-26) and identified
# placeholder positions in the source .potx that don't match the visual
# contract of our named layouts. Each fix is encoded here as a dict the
# build script can apply via XML manipulation.
#
# Adding a new fix:
#   1. Adam visually edits the layout in PowerPoint and re-saves the sample
#      deck.
#   2. Diff old vs. new layout XML to identify the substantive changes.
#   3. Append a fix entry below with rationale.
#   4. Run build_master.py; verify the fix matches Adam's intent visually.
#
# EMU = English Metric Unit, 914400 EMU = 1 inch.
# ---------------------------------------------------------------------------

INCH = 914400  # EMU per inch

# Each LAYOUT_FIXES entry has a `rationale` and a list of `shape_edits`.
# Each shape_edit identifies a target by (placeholder type+idx) or
# (shape ordinal when no placeholder), and specifies optional changes
# to xfrm, bodyPr, lvl1pPr, and defRPr.
#
# Shape identification:
#   {"by_ph": ("title", None)}       — first placeholder of type "title"
#   {"by_ph": ("body", "1")}         — body placeholder with idx="1"
#   {"by_shape_index": 0}            — first <p:sp> in document order, no ph
#
# Change keys (all optional):
#   "xfrm":     {"off_x": int, "off_y": int, "ext_cx": int, "ext_cy": int}
#   "body_pr":  {"anchor": str, "auto_fit_kind": "noAutofit"|"normAutofit"|"spAutoFit"}
#   "lvl1_ppr": {"algn": str}
#   "def_rpr":  {"sz": int, "b": int}

LAYOUT_FIXES: dict[str, dict] = {
    # --- big_number ---------------------------------------------------------
    # Source .potx has TITLE as a top header strip. For a slide where the
    # NUMBER IS the message (27M genomes, 90% accuracy), the title belongs
    # in a large centered area at 66pt bold.
    "big_number": {
        "rationale": (
            "Headline statistic must be visually dominant. Source layout had "
            "the title as a top header strip; numbers like '27,000,000' need "
            "to be the visual focus, not a header."
        ),
        "shape_edits": [
            {
                "by_ph": ("title", None),
                "xfrm": {"off_x": 660902, "off_y": 923453,
                         "ext_cx": 7840301, "ext_cy": 3286408},
                "body_pr": {"anchor": "ctr", "auto_fit_kind": "noAutofit"},
                "lvl1_ppr": {"algn": "ctr"},
                "def_rpr": {"sz": 6600, "b": 1},
            },
        ],
    },

    # --- big_idea -----------------------------------------------------------
    # Round 1 (2026-04-26): title repositioned + 36pt + noAutofit.
    # Round 2 (2026-04-26 visual review): banner shrunk so the content area
    # has room.
    # Round 3 (2026-04-26 visual review): title was supposed to BE IN the
    # accent region (banner). Adam's original y=0.72 made sense in the
    # 2.69" banner; with the 0.91" banner the title needs y=0.16 to land
    # inside the accent. Below-banner title was wrong design.
    "big_idea": {
        "rationale": (
            "The single-sentence claim sits inside the accent banner — "
            "title-on-color treatment, like the source intent. Banner "
            "shrunk to 0.91\" so the body area below is available for a "
            "supporting graphic."
        ),
        "shape_edits": [
            # Banner: shrink from 2.69" to 0.91" (rev2 — Adam visual review)
            {
                "by_kind_index": ("sp", 0),
                "xfrm": {"off_x": -3000, "off_y": 0,
                         "ext_cx": 9150000, "ext_cy": 832919},
            },
            # Title: in the accent banner (y=0.16 fits inside 0.91" banner)
            {
                "by_ph": ("title", None),
                "xfrm": {"off_x": 137160, "off_y": 146304,
                         "ext_cx": 8522208, "ext_cy": 576072},
                "body_pr": {"anchor": "ctr", "auto_fit_kind": "noAutofit"},
                "def_rpr": {"sz": 3600},
            },
        ],
    },

    # --- section_divider ----------------------------------------------------
    # The substory punchline becomes a centered, vertically-middle band of
    # large 60pt centered text — the visual cue that we're pivoting.
    #
    # 2026-04-28 (v0.2.1 fix #1, draft_9 walk):
    #   Title was at off_x=-83050 (≈ -0.09 in), placing the leftmost
    #   characters off-canvas. Fixed to off_x=0. The leftover -83050
    #   appears to have been an alignment hack from the source .potx
    #   that survived earlier rounds. All three section divider slides
    #   (5, 10, 16 in draft_9) ship with text bleeding past the left
    #   edge as a result.
    "section_divider": {
        "rationale": (
            "Substory transitions need the punchline to BE the slide. "
            "Center band, 40pt centered, full slide width. Original "
            "2026-04-26 design used 60pt + noAutofit; that combination "
            "caused 220-260 char punchlines to overflow ~1.8-2.2x even "
            "after the v0.1.1-visual autofit fix (60pt × 80% = 48pt is "
            "still too big for 200+ chars). Lowered to 40pt so autofit "
            "at 80% gives 32pt — fits 14-word punchlines once T2.4 caps "
            "them. The universal title-autofit sweep (Step 5b in "
            "build_master) runs AFTER this fix and overrides body_pr "
            "to normAutofit + anchor=t. v0.2.1: off_x corrected from "
            "-83050 (off-canvas by 0.09 in) to 0."
        ),
        "shape_edits": [
            {
                "by_ph": ("title", None),
                "xfrm": {"off_x": 0, "off_y": 1934828,
                         "ext_cx": 9144000, "ext_cy": 1273844},
                "body_pr": {"auto_fit_kind": "normAutofit"},
                "lvl1_ppr": {"algn": "ctr"},
                "def_rpr": {"sz": 4000},
            },
        ],
    },

    # --- methods_summary ----------------------------------------------------
    # 2026-04-28 (v0.2.1 fix #2, draft_9 walk):
    #   Body placeholder is undersized for production-realistic methods
    #   content. Live test of presentation-maker v0.2.0 produced 5-7
    #   paragraphs of methods prose averaging ~600-700 chars across slides
    #   6, 11, 17 in draft_9 — overflowing the placeholder's 12-line cap
    #   at 18pt by ~2-3 wrapped lines each.
    #
    #   Fix: enable normAutofit so dense methods content shrinks to fit
    #   without losing structure. fontScale 80%, lnSpcReduction 20% — same
    #   defaults that work for section_divider. The slide_compose prompt
    #   continues to enforce a 5-10 bullet cap; this fix keeps the layout
    #   readable for the upper end of that range.
    "methods_summary": {
        "rationale": (
            "Methods slides typically carry 5-7 paragraphs of substantive "
            "methodological detail (~600-800 chars). The body placeholder "
            "as shipped (3.70 in tall, 18pt body) overflows by 2-3 lines "
            "on realistic content. Enable normAutofit so dense methods "
            "shrink to fit. Title autofit handled by the universal sweep."
        ),
        "shape_edits": [
            {
                "by_ph": ("body", "1"),
                "body_pr": {"auto_fit_kind": "normAutofit"},
            },
        ],
    },

    # --- qa_anticipated -----------------------------------------------------
    # 2026-04-28 (v0.2.1 fix #4, draft_9 walk):
    #   The qa_anticipated layout is the most catastrophically undersized
    #   of all 15 layouts as shipped. Live test produced:
    #     - 4-5 line questions (192-256 chars) crammed into a Title 2
    #       placeholder sized for 2 lines (H=0.63 in)
    #     - 5-paragraph answers (~2KB each, ~39 wrapped lines at 18pt)
    #       crammed into a body placeholder sized for 12 lines (H=3.82 in)
    #   Both placeholders were ~3x undersized.
    #
    #   Fix:
    #     1. Title 2 → H 0.63 → 1.00 in (allow 3-line questions readably,
    #        autofit shrinks for 4-5 line ones)
    #     2. Body Text Placeholder 1 → T 1.17 → 1.30 in (clear the taller
    #        title), H 3.82 → 4.00 in (max we can push without overlapping
    #        the bottom logos at 5.00 in)
    #     3. Body normAutofit so 5-paragraph answers shrink to fit
    #
    #   Even with these fixes, the qa_anticipated body should ideally hold
    #   1-3 paragraphs, not 5. The qa_prep.v1.md prompt should cap answer
    #   length; that's a v0.2.x prompt-iteration follow-up, not a master
    #   fix. Layout fix here makes the deck readable in the meantime.
    "qa_anticipated": {
        "rationale": (
            "qa_anticipated body was sized for 1-paragraph answers (~12 "
            "lines @ 18pt in 3.82 in); production qa_prep produces "
            "5-paragraph answers (~39 wrapped lines, 3x overflow). "
            "Title was sized for 2-line questions; production Q&A "
            "prompts produce 4-5 line questions (3-5x overflow). Enlarge "
            "both placeholders + enable body normAutofit. Companion fix "
            "to a future qa_prep.v1 word-budget cap."
        ),
        "shape_edits": [
            # Title 2 — taller (0.63 in → 1.00 in) to handle 3-line questions
            # readably; autofit handles 4-5 line cases.
            {
                "by_ph": ("title", None),
                "xfrm": {"off_x": 91440, "off_y": 128016,
                         "ext_cx": 8521700, "ext_cy": 914400},
            },
            # Body — push down (1.17 → 1.30) and grow (3.82 → 4.00); enable
            # normAutofit so 5-paragraph answers shrink.
            {
                "by_ph": ("body", "1"),
                "xfrm": {"off_x": 311700, "off_y": 1188720,
                         "ext_cx": 8521700, "ext_cy": 3657600},
                "body_pr": {"auto_fit_kind": "normAutofit"},
            },
        ],
    },

    # --- references ---------------------------------------------------------
    # 2026-04-28 (v0.2.1 fix #5, draft_9 walk):
    #   The references body was overflowing on draft_9 slide 26 — 8 ref
    #   entries averaging 134 chars each (~17 wrapped lines at 18pt) in
    #   a placeholder sized for 12 lines. Same enable-normAutofit fix as
    #   methods_summary.
    #
    #   Note: the AI-disclosure footer at 8pt (separate textbox added by
    #   _fill_references) is intentional — brand_tokens["sizes_pt"]
    #   ["ai_disclosure"] = 8. Walker's [TINY-FONT] flag on that
    #   textbox is a false positive against the brand spec; not a bug.
    "references": {
        "rationale": (
            "References body sized for 12 lines @ 18pt overflows on "
            "8 references averaging 134 chars each (~17 wrapped lines). "
            "Enable normAutofit so reference list shrinks to fit. "
            "AI-disclosure footer at 8pt is per brand spec (caption "
            "category), not a bug."
        ),
        "shape_edits": [
            {
                "by_ph": ("body", "1"),
                "body_pr": {"auto_fit_kind": "normAutofit"},
            },
        ],
    },

    # --- two_column_compare -------------------------------------------------
    # Source has columns squished into the bottom half; banner SP is 2.7"
    # tall (eats the top). Adam's edit shrinks the banner to 0.9" and
    # stretches both columns to full vertical extent.
    #
    # BUG FIX 2026-04-26 round 2: the banner is at drawable index 2 in source
    # (after two pic logos at indices 0–1). The original encoding used
    # by_shape_index: 0, which moved a LOGO to the top — leaving the actual
    # banner SP unshrunk. Switched to by_kind_index: ("sp", 0) which targets
    # the first NON-PLACEHOLDER <p:sp> regardless of pic/sp ordering.
    "two_column_compare": {
        "rationale": (
            "Source layout had columns occupying only the bottom half of "
            "the slide; the decorative banner ate the top. Both columns "
            "need full height for substantive comparison content."
        ),
        "shape_edits": [
            # Decorative top banner SP — shrink to a normal header strip.
            # by_kind_index: ("sp", 0) targets the first non-placeholder sp
            # (the banner) regardless of where logos sit in document order.
            {
                "by_kind_index": ("sp", 0),
                "xfrm": {"off_x": -3000, "off_y": 0,
                         "ext_cx": 9150000, "ext_cy": 832919},
            },
            # Left column body — move up and stretch
            {
                "by_ph": ("body", "1"),
                "xfrm": {"off_x": 112523, "off_y": 981671,
                         "ext_cx": 4386578, "ext_cy": 3479572},
            },
            # Right column body — move up and stretch
            {
                "by_ph": ("body", "2"),
                "xfrm": {"off_x": 4572000, "off_y": 981670,
                         "ext_cx": 4488950, "ext_cy": 3479572},
            },
        ],
    },

    # --- concept_illustration -----------------------------------------------
    # Adam's edits: pull a hidden decorative shape onto the slide edge, and
    # narrow the body so an AI-generated image has room beside it.
    "concept_illustration": {
        "rationale": (
            "AI-generated illustration slides need a visible visual anchor "
            "beside the title text. The source layout had a decorative shape "
            "hidden far off-screen; pulling it to the edge frames the body."
        ),
        "shape_edits": [
            # Hidden decorative shape — pull to slide edge
            {
                "by_shape_index": 0,
                "xfrm": {"off_x": -3, "off_y": 0,
                         "ext_cx": 9144003, "ext_cy": 5143501},
            },
            # Body placeholder — narrow to leave room for image on right
            {
                "by_ph": ("body", "1"),
                "xfrm": {"off_x": 311700, "off_y": 1185250,
                         "ext_cx": 4405155, "ext_cy": 3383700},
            },
        ],
    },

    # Other layouts may accumulate fixes here as visual review surfaces
    # issues. Keep each fix small and explicit — never a "while we're at
    # it" sweep.
}


# ---------------------------------------------------------------------------
# Brand tokens — canonical from KBase Style Guide June 2022.
# See reference/kbase-style-extract.md §2, §3, §4, §5.
# ---------------------------------------------------------------------------

BRAND_TOKENS: dict = {
    "version": "1.0",
    "source": "KBase Style Guidelines — PRINT & PRESENTATION, June 2022",
    "palette": {
        "primary": {
            "microbe_orange":  {"hex": "#F78E1E", "rgb": [247, 142, 30],
                                "tints": {"80": "#F9A455", "60": "#FBBA8C", "40": "#FDD0BB", "20": "#FFE5CB"}},
            "grass_green":     {"hex": "#5E9732", "rgb": [94, 151, 50],
                                "tints": {"80": "#7EAC5B", "60": "#9EC184", "40": "#BED5AD", "20": "#DFEAD6"}},
            "freshwater_blue": {"hex": "#007DC3", "rgb": [0, 125, 195],
                                "tints": {"80": "#3397CF", "60": "#66B1DB", "40": "#99CBE7", "20": "#CCE5F3"}},
            "golden_yellow":   {"hex": "#FFD200", "rgb": [255, 210, 0],
                                "tints": {"80": "#FFDB33", "60": "#FFE466", "40": "#FFED99", "20": "#FFF6CC"}},
            "spring_green":    {"hex": "#C1CD23", "rgb": [193, 205, 35],
                                "tints": {"80": "#CDD659", "60": "#DAE183", "40": "#E6EBAC", "20": "#F3F5D6"}},
            "ocean_blue":      {"hex": "#72CCD2", "rgb": [114, 204, 210],
                                "tints": {"80": "#8ED6DB", "60": "#AAE0E4", "40": "#C6EAED", "20": "#E3F5F6"}}
        },
        "secondary": {
            "cyanobacteria_teal": {"hex": "#009688", "rgb": [0, 150, 136]},
            "lupine_purple":      {"hex": "#66489D", "rgb": [102, 72, 157]},
            "frost_blue":         {"hex": "#C7DBEE", "rgb": [199, 219, 238]},
            "rainier_cherry_red": {"hex": "#D2232A", "rgb": [210, 35, 42]},
            "graphite_gray":      {"hex": "#9D9389", "rgb": [157, 146, 135]}
        },
        "neutral": {
            "white": {"hex": "#FFFFFF"},
            "black": {"hex": "#000000"}
        },
        "contrast_warnings": [
            {"pair": ["spring_green", "golden_yellow"],
             "reason": "KBase Style Guide §5: forbidden as contrasting colors."}
        ]
    },
    "typography": {
        "primary": {"family": "Oxygen", "weights": ["Regular", "Bold", "Italic", "Bold Italic"]},
        "fallback": {"family": "Calibri", "weights": ["Regular", "Light", "Italic", "Bold"]},
        "code": {"family": "Courier", "weights": ["Regular", "Bold"]},
        "sizes_pt": {
            "deck_title": 60,
            "section_divider": 48,
            "big_idea": 60,
            "big_number": 96,
            "content_title": 36,
            "body": 24,
            "caption": 14,
            "ai_disclosure": 8
        }
    },
    "logo": {
        "min_height_pt": 36,
        "clear_space_x_factor": 0.5
    },
    "contrast_minima_wcag": {
        "body_text_aa_ratio": 4.5,
        "body_text_aaa_ratio": 7.0,
        "large_text_aa_ratio": 3.0,
        "large_text_aaa_ratio": 4.5
    },
    "vocabulary": list(NAMED_VOCABULARY)
}


# ---------------------------------------------------------------------------
# I/O helpers
# ---------------------------------------------------------------------------

def repo_root() -> Path:
    """Walk up from this file until we find pyproject.toml at the repo root."""
    here = Path(__file__).resolve()
    for parent in [here, *here.parents]:
        if (parent / "pyproject.toml").is_file() and (parent / "SPEC.md").is_file():
            return parent
    raise RuntimeError(
        f"Could not locate repo root (looked for pyproject.toml + SPEC.md) starting from {here}"
    )


def potx_to_pptx_bytes(potx_path: Path) -> bytes:
    """Read a .potx and return the bytes of an equivalent .pptx (content type rewritten).

    python-pptx cannot open .potx directly. The only difference at the file
    level is the override declared in [Content_Types].xml: ...presentation.template
    vs ...presentation.presentation. Rewriting the content type produces a valid .pptx.
    """
    if not potx_path.is_file():
        raise FileNotFoundError(
            f"Source .potx not found at {potx_path}. "
            f"This file is gitignored and user-supplied; "
            f"see reference/master-template-source-notes.md §1 for sourcing details."
        )

    out = io.BytesIO()
    with zipfile.ZipFile(potx_path, "r") as zin, \
         zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                )
            zout.writestr(item, data)
    return out.getvalue()


# ---------------------------------------------------------------------------
# python-pptx XML helpers
# ---------------------------------------------------------------------------

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _set_layout_name(layout_part_xml: bytes, new_name: str) -> bytes:
    """Set <p:cSld name="..."> on a slide layout XML and return the updated bytes."""
    tree = etree.fromstring(layout_part_xml)
    csld = tree.find(qn("p:cSld"))
    if csld is None:
        raise ValueError("slideLayout XML missing <p:cSld> element")
    csld.set("name", new_name)
    return etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)


def _layout_name_from_xml(layout_part_xml: bytes) -> str:
    tree = etree.fromstring(layout_part_xml)
    csld = tree.find(qn("p:cSld"))
    return csld.get("name", "") if csld is not None else ""


def _drop_layout_refs_from_master_element(master_element, layout_rels_to_drop: set[str]) -> int:
    """Remove <p:sldLayoutId> entries from a slide master's <p:sldLayoutIdLst>
    whose r:id is in layout_rels_to_drop. Operates on the lxml element in-place.
    Returns count of refs removed.
    """
    lst = master_element.find(qn("p:sldLayoutIdLst"))
    if lst is None:
        return 0
    n = 0
    for layout_id in list(lst):
        rid = layout_id.get(qn("r:id"))
        if rid in layout_rels_to_drop:
            lst.remove(layout_id)
            n += 1
    return n


def _drop_master_ref_from_presentation_element(pres_element, master_rid: str) -> int:
    """Remove a <p:sldMasterId> with the given r:id from <p:sldMasterIdLst>.
    Operates on lxml element in-place. Returns count of refs removed.
    """
    lst = pres_element.find(qn("p:sldMasterIdLst"))
    if lst is None:
        return 0
    n = 0
    for master_id in list(lst):
        rid = master_id.get(qn("r:id"))
        if rid == master_rid:
            lst.remove(master_id)
            n += 1
    return n


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _drawables_in_sptree(layout_element):
    """Yield drawable children of <p:cSld>/<p:spTree> in document order.
    Drawable = <p:sp>, <p:pic>, <p:grpSp> (group shape), <p:graphicFrame>.
    """
    sptree = layout_element.find(f".//{{{P_NS}}}cSld/{{{P_NS}}}spTree")
    if sptree is None:
        return
    drawable_locals = {"sp", "pic", "grpSp", "graphicFrame"}
    for child in sptree:
        local = etree.QName(child).localname
        if local in drawable_locals:
            yield child


def _find_target_shape(layout_element, shape_edit: dict):
    """Resolve a shape_edit's target to a drawable element.

    Targeting modes:
      {"by_ph": ("title", None)}    — first <p:sp> whose ph.type=="title"
      {"by_ph": ("body", "1")}      — first <p:sp> whose ph.type=="body" and idx=="1"
      {"by_shape_index": N}         — Nth drawable in spTree document order
                                      (covers <p:sp>, <p:pic>, <p:grpSp>,
                                      <p:graphicFrame>)
      {"by_kind_index": ("sp", N)}  — Nth NON-PLACEHOLDER drawable of kind 'sp'
                                      (decorative shapes only — placeholders
                                      excluded). Use this to target banners,
                                      decorative bars, etc., independent of
                                      the document order of pics/sps which
                                      varies across source layouts.
      {"by_kind_index": ("pic", N)} — Nth drawable of kind 'pic'

    Returns the element. Raises ValueError if not found.

    Why by_kind_index exists: the source .potx has different drawable
    orderings across layouts (pic-first in some, sp-first in others). A
    naive `by_shape_index: 0` ends up targeting a logo in one layout and
    a banner in another. by_kind_index makes targeting robust against
    drawable-order shuffles by filtering on kind first. Lesson learned
    on 2026-04-26 when by_shape_index: 0 in two_column_compare moved a
    logo to (0,0) instead of shrinking the banner SP.
    """
    if "by_ph" in shape_edit:
        # Placeholder lookup is meaningful only for <p:sp> (placeholders live
        # in <p:nvSpPr>/<p:nvPr>/<p:ph>, which only sp elements have).
        ph_type, ph_idx = shape_edit["by_ph"]
        for sp in layout_element.iter(f"{{{P_NS}}}sp"):
            ph = sp.find(f".//{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
            if ph is None:
                continue
            this_type = ph.get("type", "body")
            this_idx = ph.get("idx", "0")
            if ph_type == "title":
                # 'title' placeholder may be type="title" or "ctrTitle";
                # also a placeholder with idx==0 implicitly.
                if this_type == "title" or this_type == "ctrTitle" or this_idx == "0":
                    return sp
            else:
                if this_type == ph_type and (ph_idx is None or this_idx == ph_idx):
                    return sp
        raise ValueError(f"target placeholder not found: by_ph={shape_edit['by_ph']}")

    if "by_shape_index" in shape_edit:
        n = shape_edit["by_shape_index"]
        drawables = list(_drawables_in_sptree(layout_element))
        if n >= len(drawables):
            raise ValueError(
                f"by_shape_index={n} out of range; layout has {len(drawables)} drawables"
            )
        return drawables[n]

    if "by_kind_index" in shape_edit:
        kind, n = shape_edit["by_kind_index"]
        drawables = list(_drawables_in_sptree(layout_element))
        matching = []
        for d in drawables:
            local = etree.QName(d).localname
            if local != kind:
                continue
            # For 'sp' kind, exclude placeholders — we target decorative shapes.
            if local == "sp":
                ph = d.find(f".//{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
                if ph is not None:
                    continue
            matching.append(d)
        if n >= len(matching):
            raise ValueError(
                f"by_kind_index=({kind!r}, {n}) out of range; "
                f"only {len(matching)} non-placeholder {kind!r} drawables found"
            )
        return matching[n]

    raise ValueError(
        f"shape_edit has no target spec "
        f"(by_ph, by_shape_index, or by_kind_index): {shape_edit}"
    )


def _apply_xfrm_change(sp, xfrm_change: dict) -> None:
    """Apply position/size changes to a shape's <a:xfrm>."""
    xfrm = sp.find(f".//{{{A_NS}}}xfrm")
    if xfrm is None:
        raise ValueError("shape has no <a:xfrm> (inherits from master); "
                         "cannot apply absolute position via build_master")
    off = xfrm.find(f"{{{A_NS}}}off")
    ext = xfrm.find(f"{{{A_NS}}}ext")
    if off is None or ext is None:
        raise ValueError("<a:xfrm> missing <a:off> or <a:ext>")
    if "off_x" in xfrm_change:
        off.set("x", str(xfrm_change["off_x"]))
    if "off_y" in xfrm_change:
        off.set("y", str(xfrm_change["off_y"]))
    if "ext_cx" in xfrm_change:
        ext.set("cx", str(xfrm_change["ext_cx"]))
    if "ext_cy" in xfrm_change:
        ext.set("cy", str(xfrm_change["ext_cy"]))


def _apply_body_pr_change(sp, body_pr_change: dict) -> None:
    """Apply anchor / auto-fit changes to <p:txBody><a:bodyPr>."""
    body_pr = sp.find(f".//{{{P_NS}}}txBody/{{{A_NS}}}bodyPr")
    if body_pr is None:
        raise ValueError("shape has no <p:txBody>/<a:bodyPr>")
    if "anchor" in body_pr_change:
        body_pr.set("anchor", body_pr_change["anchor"])
    if "auto_fit_kind" in body_pr_change:
        # Remove any existing autofit child(ren), then add the requested one.
        for tag in ("normAutofit", "noAutofit", "spAutoFit"):
            for child in list(body_pr):
                if child.tag == f"{{{A_NS}}}{tag}":
                    body_pr.remove(child)
        kind = body_pr_change["auto_fit_kind"]
        new_child = etree.SubElement(body_pr, f"{{{A_NS}}}{kind}")
        # PowerPoint requires explicit fontScale + lnSpcReduction on
        # normAutofit to actually shrink at render time. Bare normAutofit
        # is treated as "autofit-eligible" but no scale is computed,
        # leaving long text to overflow. Defaults are 80% font / 20% line
        # spacing reduction; PowerPoint shrinks further as needed.
        if kind == "normAutofit":
            # Allow caller override via body_pr_change["norm_autofit"] dict
            af_params = body_pr_change.get("norm_autofit", {})
            new_child.set("fontScale", str(af_params.get("fontScale", 80000)))
            new_child.set("lnSpcReduction",
                          str(af_params.get("lnSpcReduction", 20000)))


def _apply_lvl1_ppr_change(sp, lvl1_change: dict) -> None:
    """Apply paragraph-level changes to <a:lstStyle><a:lvl1pPr>."""
    lst_style = sp.find(f".//{{{P_NS}}}txBody/{{{A_NS}}}lstStyle")
    if lst_style is None:
        raise ValueError("shape has no <p:txBody>/<a:lstStyle>")
    lvl1 = lst_style.find(f"{{{A_NS}}}lvl1pPr")
    if lvl1 is None:
        raise ValueError("<a:lstStyle> missing <a:lvl1pPr>")
    if "algn" in lvl1_change:
        lvl1.set("algn", lvl1_change["algn"])


def _apply_def_rpr_change(sp, def_rpr_change: dict) -> None:
    """Apply default run-property changes (font size, bold) to <a:lvl1pPr>/<a:defRPr>."""
    lst_style = sp.find(f".//{{{P_NS}}}txBody/{{{A_NS}}}lstStyle")
    if lst_style is None:
        raise ValueError("shape has no <p:txBody>/<a:lstStyle>")
    lvl1 = lst_style.find(f"{{{A_NS}}}lvl1pPr")
    if lvl1 is None:
        raise ValueError("<a:lstStyle> missing <a:lvl1pPr>")
    def_rpr = lvl1.find(f"{{{A_NS}}}defRPr")
    if def_rpr is None:
        def_rpr = etree.SubElement(lvl1, f"{{{A_NS}}}defRPr")
    if "sz" in def_rpr_change:
        def_rpr.set("sz", str(def_rpr_change["sz"]))
    if "b" in def_rpr_change:
        def_rpr.set("b", str(def_rpr_change["b"]))


def _apply_shape_edit(layout_element, shape_edit: dict) -> None:
    """Apply a single shape_edit to a layout. Mutates layout_element in place."""
    sp = _find_target_shape(layout_element, shape_edit)
    if "xfrm" in shape_edit:
        _apply_xfrm_change(sp, shape_edit["xfrm"])
    if "body_pr" in shape_edit:
        _apply_body_pr_change(sp, shape_edit["body_pr"])
    if "lvl1_ppr" in shape_edit:
        _apply_lvl1_ppr_change(sp, shape_edit["lvl1_ppr"])
    if "def_rpr" in shape_edit:
        _apply_def_rpr_change(sp, shape_edit["def_rpr"])


def _apply_layout_fixes(prs, verbose: bool = True) -> list[str]:
    """Apply all LAYOUT_FIXES to the post-rename presentation. Returns list of
    layout names that received fixes.
    """
    applied: list[str] = []
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
    for layout_name, fix in LAYOUT_FIXES.items():
        if layout_name not in layouts:
            raise ValueError(
                f"LAYOUT_FIXES references '{layout_name}' but this layout "
                f"is not present in the renamed master. Check LAYOUT_RENAMES."
            )
        layout = layouts[layout_name]
        for shape_edit in fix["shape_edits"]:
            try:
                _apply_shape_edit(layout.element, shape_edit)
            except ValueError as e:
                raise ValueError(
                    f"LAYOUT_FIXES['{layout_name}'] shape_edit failed: {e}\n"
                    f"  shape_edit was: {shape_edit}"
                ) from e
        applied.append(layout_name)
        if verbose:
            print(f"[build_master] applied layout fix: {layout_name} "
                  f"({len(fix['shape_edits'])} shape edit(s))")
    return applied


# Layouts where LAYOUT_FIXES intentionally sets noAutofit on the title and
# we should NOT override (the title font size is part of the design intent).
# big_number: 96pt headline that must NEVER shrink (the number IS the slide).
# big_idea: 36pt single-line claim in the accent banner — designed to be tight.
# These layouts depend on the prompts capping title length; if titles overflow
# here, the prompt is the bug, not the master.
_LAYOUTS_WITH_INTENTIONAL_NO_AUTOFIT_TITLE = frozenset({
    "big_number",
    "big_idea",
})


def _ensure_title_autofit_universal(prs, verbose: bool = True) -> list[str]:
    """For every layout EXCEPT those in _LAYOUTS_WITH_INTENTIONAL_NO_AUTOFIT_TITLE,
    ensure the title placeholder's <a:bodyPr> has an explicit
    <a:normAutofit fontScale="80000" lnSpcReduction="20000"/> child + anchor="t".

    Why explicit fontScale: PowerPoint treats a bare `<a:normAutofit/>` as
    "autofit-eligible" but does NOT compute a scale at render time. Without
    fontScale set explicitly, long titles overflow the placeholder. With
    fontScale=80000 (80%) + lnSpcReduction=20000 (20%), PowerPoint applies
    those defaults at render and shrinks further as needed.

    Why anchor="t": text overflow grows downward only. Centered (anchor="ctr")
    overflow grows in BOTH directions, so a long title bleeds UPWARD into the
    logo region (live failure mode on draft_4 slide 1, 2026-04-26).

    Returns list of layout names that received the autofit sweep.
    """
    swept: list[str] = []
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name in _LAYOUTS_WITH_INTENTIONAL_NO_AUTOFIT_TITLE:
            continue
        # Find the title placeholder (type "title" or "ctrTitle"; idx may vary)
        title_sp = None
        for sp in layout.element.iter(f"{{{P_NS}}}sp"):
            ph = sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}nvPr/{{{P_NS}}}ph")
            if ph is None:
                continue
            ph_type = ph.get("type", "body")
            if ph_type in ("title", "ctrTitle"):
                title_sp = sp
                break
        if title_sp is None:
            continue
        body_pr = title_sp.find(
            f"{{{P_NS}}}txBody/{{{A_NS}}}bodyPr"
        )
        if body_pr is None:
            continue
        # Remove any existing autofit child
        for tag in ("normAutofit", "noAutofit", "spAutoFit"):
            for child in list(body_pr):
                if child.tag == f"{{{A_NS}}}{tag}":
                    body_pr.remove(child)
        # Add normAutofit with explicit attributes
        autofit = etree.SubElement(body_pr, f"{{{A_NS}}}normAutofit")
        autofit.set("fontScale", "80000")
        autofit.set("lnSpcReduction", "20000")
        # Set anchor=t so overflow grows downward, not into logo region
        body_pr.set("anchor", "t")
        swept.append(layout.name)
        if verbose:
            print(f"[build_master] title autofit set on layout: {layout.name}")
    return swept


def _drop_all_slide_refs_from_presentation_element(pres_element) -> int:
    """Remove ALL <p:sldId> entries from <p:sldIdLst>. The derived master ships with
    no slides (the assembler creates fresh decks each invocation).
    Returns count of refs removed.
    """
    lst = pres_element.find(qn("p:sldIdLst"))
    if lst is None:
        return 0
    n = 0
    for sld_id in list(lst):
        lst.remove(sld_id)
        n += 1
    return n


# ---------------------------------------------------------------------------
# Build pipeline
# ---------------------------------------------------------------------------

def build_master(source_potx: Path, dest_pptx: Path, brand_tokens_dest: Path,
                 verbose: bool = True) -> dict:
    """Run the full build pipeline. Returns a build-report dict."""
    report: dict = {
        "source": str(source_potx),
        "dest_master": str(dest_pptx),
        "dest_brand_tokens": str(brand_tokens_dest),
        "renamed_layouts": [],
        "deleted_layouts": [],
        "deleted_masters": [],
    }

    # Step 1: validate + convert content-type
    pptx_bytes = potx_to_pptx_bytes(source_potx)

    # Step 2: open with python-pptx for inspection
    prs = Presentation(io.BytesIO(pptx_bytes))
    if len(prs.slide_masters) == 0:
        raise ValueError("source has no slide masters")

    if verbose:
        print(f"[build_master] source has {len(prs.slide_masters)} slide masters, "
              f"master[0] has {len(prs.slide_masters[0].slide_layouts)} layouts")

    master0 = prs.slide_masters[0]

    # Step 3: identify layouts in master 0 to keep (by source name) vs delete
    keep_source_names = set(LAYOUT_RENAMES.keys())
    keep_indices: dict[str, int] = {}
    drop_indices: list[int] = []

    for li, layout in enumerate(master0.slide_layouts):
        if layout.name in keep_source_names:
            if layout.name in keep_indices:
                # Duplicate source name (shouldn't happen — but guard)
                drop_indices.append(li)
            else:
                keep_indices[layout.name] = li
        else:
            drop_indices.append(li)

    missing = keep_source_names - set(keep_indices.keys())
    if missing:
        raise ValueError(
            f"source master is missing expected layouts: {sorted(missing)}\n"
            f"available layout names: {sorted(l.name for l in master0.slide_layouts)}\n"
            f"this is a brand-version skew; update LAYOUT_RENAMES in build_master.py "
            f"or restore the missing layouts in the .potx."
        )

    if verbose:
        print(f"[build_master] keep {len(keep_indices)} layouts, "
              f"drop {len(drop_indices)} layouts from master[0]")

    # Step 4: rename kept layouts via XML
    rename_summary = []
    for source_name, new_name in LAYOUT_RENAMES.items():
        idx = keep_indices[source_name]
        layout = master0.slide_layouts[idx]
        # python-pptx exposes layout.element (lxml node). Set <p:cSld name="...">
        csld = layout.element.find(qn("p:cSld"))
        if csld is None:
            raise ValueError(f"layout '{source_name}' has no <p:cSld>")
        csld.set("name", new_name)
        rename_summary.append((source_name, new_name))

    report["renamed_layouts"] = [{"from": s, "to": d} for (s, d) in rename_summary]

    # Step 5: drop unused layouts from master 0
    # python-pptx doesn't expose layout deletion, so manipulate the package directly
    pkg = prs.part.package
    master0_part = master0.part

    # collect (source_name, rel_id, partname) for layouts to drop
    drops: list[tuple[str, str, str]] = []
    for li in drop_indices:
        layout = master0.slide_layouts[li]
        # find the relationship from master0 to this layout
        rel = next(
            (r for r in master0_part.rels.values()
             if r.target_part is layout.part), None)
        if rel is None:
            continue
        drops.append((layout.name, rel.rId, layout.part.partname))

    # remove them from master0's relationships
    # (python-pptx 1.x prunes orphaned parts at save time — verified empirically;
    # we don't need to touch package._parts directly. Dropping the rel is sufficient.)
    drop_rids = {rid for (_, rid, _) in drops}
    for (_, rid, _) in drops:
        if rid in master0_part.rels:
            master0_part.rels.pop(rid)

    # remove <p:sldLayoutId> references from master XML (in-place on element)
    _drop_layout_refs_from_master_element(master0_part._element, drop_rids)

    report["deleted_layouts"] = [{"name": n, "rId": r} for (n, r, _) in drops]

    # Step 5b: apply hand-tuned layout fixes (Adam's visual review).
    applied_fixes = _apply_layout_fixes(prs, verbose=verbose)
    report["applied_fixes"] = applied_fixes

    # Step 5b: universal title-autofit sweep — every layout except those
    # in _LAYOUTS_WITH_INTENTIONAL_NO_AUTOFIT_TITLE gets explicit
    # normAutofit + anchor=t on its title placeholder. Fixes the
    # 2026-04-26 visual review finding that 19/21 slides had title text
    # 1.6-5.0x oversize because PowerPoint doesn't auto-shrink with bare
    # normAutofit.
    autofit_swept = _ensure_title_autofit_universal(prs, verbose=verbose)
    report["title_autofit_swept"] = autofit_swept

    # Step 6: drop slide_masters[1:] (Google Slides round-trip duplicates)
    pres_part = prs.part
    if len(prs.slide_masters) > 1:
        masters_to_drop = list(prs.slide_masters)[1:]
        master_drop_rids: list[str] = []
        for extra in masters_to_drop:
            extra_part = extra.part
            # find rel from presentation → this master
            rel = next(
                (r for r in pres_part.rels.values()
                 if r.target_part is extra_part), None)
            if rel is None:
                continue
            master_drop_rids.append(rel.rId)
            # remove its layouts (cascade) — drop rels only, parts are pruned at save
            for layout in list(extra.slide_layouts):
                lay_rel = next(
                    (r for r in extra_part.rels.values()
                     if r.target_part is layout.part), None)
                if lay_rel is not None and lay_rel.rId in extra_part.rels:
                    extra_part.rels.pop(lay_rel.rId)
            # remove the master rel from the presentation
            if rel.rId in pres_part.rels:
                pres_part.rels.pop(rel.rId)
            report["deleted_masters"].append(
                {"partname": str(extra_part.partname), "rId": rel.rId}
            )

        # remove <p:sldMasterId> references from presentation.xml (in-place)
        for rid in master_drop_rids:
            _drop_master_ref_from_presentation_element(pres_part._element, rid)

    # Step 6b: drop ALL slides from the .potx — the derived master ships empty.
    # Source slides (32 in the Gazi example deck) are template fodder; the
    # assembler creates fresh decks per invocation.
    slide_rids_to_drop: list[str] = []
    if hasattr(prs, "slides"):
        for slide in list(prs.slides):
            slide_part = slide.part
            slide_rel = next(
                (r for r in pres_part.rels.values()
                 if r.target_part is slide_part), None)
            if slide_rel is None:
                continue
            slide_rids_to_drop.append(slide_rel.rId)
            if slide_rel.rId in pres_part.rels:
                pres_part.rels.pop(slide_rel.rId)
    n_slides_dropped = _drop_all_slide_refs_from_presentation_element(pres_part._element)
    report["deleted_slides_count"] = n_slides_dropped

    # Step 7: save derived master
    dest_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(dest_pptx)

    # Step 8: emit brand-tokens JSON
    brand_tokens_dest.parent.mkdir(parents=True, exist_ok=True)
    brand_tokens_dest.write_text(
        json.dumps(BRAND_TOKENS, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )

    # Step 9: hash the output for idempotency check
    dest_bytes = dest_pptx.read_bytes()
    report["dest_master_sha256"] = hashlib.sha256(dest_bytes).hexdigest()
    report["dest_master_size_bytes"] = len(dest_bytes)

    if verbose:
        print(f"[build_master] wrote {dest_pptx} "
              f"({len(dest_bytes):,} bytes; sha256={report['dest_master_sha256'][:12]}...)")
        print(f"[build_master] wrote {brand_tokens_dest}")

    return report


# ---------------------------------------------------------------------------
# Verification helpers (also used by test_build_master.py)
# ---------------------------------------------------------------------------

def verify_built_master(master_pptx: Path) -> dict:
    """Open the built master, verify all 15 named layouts present + no extras, and return
    a structured report.
    """
    prs = Presentation(master_pptx)
    n_masters = len(prs.slide_masters)
    if n_masters != 1:
        raise AssertionError(f"expected exactly 1 slide master, got {n_masters}")

    layouts = list(prs.slide_masters[0].slide_layouts)
    found_names = sorted(l.name for l in layouts)
    expected_names = sorted(NAMED_VOCABULARY)

    if found_names != expected_names:
        extra = sorted(set(found_names) - set(expected_names))
        missing = sorted(set(expected_names) - set(found_names))
        raise AssertionError(
            f"layout name mismatch:\n"
            f"  extra:   {extra}\n"
            f"  missing: {missing}"
        )

    # Each layout must have at least one TITLE placeholder (or be one of the
    # exempt layouts in SPEC §6.1 — references / acknowledgments / qa_anticipated
    # all still happen to have titles in our mapping, so no exemption needed).
    title_check = []
    for layout in layouts:
        has_title = any(
            ph.placeholder_format.idx == 0
            for ph in layout.placeholders
        )
        title_check.append({"layout": layout.name, "has_title_placeholder": has_title})

    return {
        "n_masters": n_masters,
        "n_layouts": len(layouts),
        "layout_names_sorted": found_names,
        "title_check": title_check,
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="build_master",
        description="Build kbase-presentation-master.pptx from user-supplied .potx.",
    )
    parser.add_argument(
        "--source", "-s",
        type=Path,
        default=None,
        help="Path to source .potx (default: <repo>/reference/master-template-source/KBase 2026 and beyond.potx)",
    )
    parser.add_argument(
        "--dest", "-d",
        type=Path,
        default=None,
        help="Path to derived .pptx output (default: <repo>/src/.../skill/references/templates/kbase-presentation-master.pptx)",
    )
    parser.add_argument(
        "--brand-tokens",
        type=Path,
        default=None,
        help="Path to kbase-brand-tokens.json output (default: <repo>/src/.../skill/references/kbase-brand-tokens.json)",
    )
    parser.add_argument(
        "--quiet", "-q", action="store_true", help="Suppress build progress output.",
    )
    parser.add_argument(
        "--verify-only", action="store_true",
        help="Don't rebuild; verify the existing derived master matches expected vocabulary.",
    )
    args = parser.parse_args(argv)

    root = repo_root()
    source = args.source or (root / "reference" / "master-template-source" /
                             "KBase 2026 and beyond.potx")
    dest = args.dest or (root / "src" / "beril_presentation_maker" / "skill" /
                         "references" / "templates" / "kbase-presentation-master.pptx")
    brand = args.brand_tokens or (root / "src" / "beril_presentation_maker" / "skill" /
                                  "references" / "kbase-brand-tokens.json")

    if args.verify_only:
        report = verify_built_master(dest)
        print(json.dumps(report, indent=2))
        return 0

    try:
        report = build_master(source, dest, brand, verbose=not args.quiet)
    except FileNotFoundError as e:
        print(f"build_master: {e}", file=sys.stderr)
        return 3
    except ValueError as e:
        print(f"build_master: {e}", file=sys.stderr)
        return 2

    # Verify after building.
    verify = verify_built_master(dest)
    report["verification"] = verify
    if not args.quiet:
        print(f"[build_master] verification: {verify['n_layouts']} layouts, "
              f"{verify['n_masters']} masters; all 15 named layouts present.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

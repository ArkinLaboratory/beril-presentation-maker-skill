#!/usr/bin/env python3
"""validate_deliverable.py — Cycle 1 pre-handoff deterministic gate.

Pure, READ-ONLY check over the produced deliverable (`deliverable/
draft.pptx`) + the working/audit artifacts. Emits findings under the
`deliverable-validation.v1` schema. Modeled on the existing
`content-overflow.v1` / `layout-overlaps.v1` Tier-1 cascade-finding
schemas (CRAFT-CONTRACT §2.1): per-finding `id` / `gate` / `severity`
/ `slide_id_or_target` / `message` / `remediation`. Telemetry-ready:
projectable fields (`gate`, `severity`, `remediation.kind`) are
tokens drawn from frozen vocabularies; nothing free-text in those
slots. NOT the cross-skill run-record contract (that's a later cycle).

Six gates (each maps to a caulobacter live-run defect → regression
fixture):

  G1 placeholder_or_leaked_template
       No TBD / "TBD - …" tokens in human-visible fields; no
       project-dir-name token in title/presenter (caught the
       `Lipida` typo); title + presenter populated from beril.yaml.
       Tightly scoped: only checks the title/presenter/affiliation
       slot on the title slide + contributors on acknowledgments.
       The deeper "no TBD anywhere" sweep would false-positive on
       legitimately TBD-marked figure placeholders mid-pipeline.

  G2 image_completeness
       Every `working/05_image_requests/<sid>_request.json` resolves
       to either an embedded picture in the deck (success) or a
       manifest reject/skip entry (explicit decline). This is the
       Cycle-1 widening of v1.1.1's image_gen_orchestrate
       assert-requests-resolved: that fired only at stage_image_gen
       tail; this fires at deliverable time and additionally
       confirms the PNG actually got into the rendered pptx (not
       just landed on disk). Also: if image_gen ran (manifest
       present) AND it produced zero approvals → advisory finding
       so a "no images at all" outcome surfaces.

  G3 slide_count_vs_budget
       Final slide count within MODE_SLIDE_BUDGETS for the
       persisted mode. Out-of-band is advisory only (renderers
       sometimes produce off-budget decks for good reasons);
       slightly low (within 10% of the floor) is a separate softer
       advisory ("budget edge"). This is the only gate that's
       intrinsically advisory — quantitative judgment, not yes/no.

  G4 mode_vs_user_intent
       The new Gate 4. Reads `audit/user_intent.json` (DP9b
       persistence layer). Asserts the `mode` recorded in
       `working/slide_spec.json`, `working/05_image_decisions.json`,
       `working/03_slides/qa_anticipated.json` ALL equal the
       user's persisted explicit mode pick. Cross-artifact agreement
       (which v1.1.1's mode_consistency check enforces) is NOT
       enough: the caulobacter run had every artifact uniformly
       agree on talk-30 while the user picked talk-45. Without
       user_intent.json there's no anchor to compare against.

  G5 figure_integrity (display vs. native aspect)
       For every embedded picture: display AR (pic.width /
       pic.height) ≈ native AR (image.size[0] / image.size[1]),
       within the ASPECT_SKEW_TOLERANCE band (5%). Guards the
       v1.1.1 DP4 fix from regressing — a future contributor who
       reintroduces the both-w-and-h add_picture call would have
       every figure trip this gate.

  G6 figure_path_resolution
       Every `content.figure` (and `content.supporting_graphic` /
       `content.image_path`) referenced by slide_spec.json
       resolves to a real file on disk via the same lookup logic
       assemble_pptx uses (draft_dir first, then project_dir).
       Catches stale spec entries where the curator dropped or
       moved a figure between compose and assemble.

Each finding carries a `remediation` block keyed to cost:

  remediation.kind ∈ {
      "auto",         # deterministic + ~free; finalize fixes + re-checks
      "targeted",     # cheap one-stage re-run; exact command emitted
      "advisory",     # surface, never block
  }
  remediation.action ∈ {                         # auto-routes (kind=auto)
      "reassemble",                  # rerun assemble only (no LLM)
      "populate_title_from_beril",   # read beril.yaml, splice into spec
      "strip_dirname_token",         # strip project-dir tokens from text
  }
  remediation.command (kind=targeted) — exact bash/python invocation.

DETECTION vs REMEDIATION. This module is the detection half ONLY.
It is pure: no filesystem mutation. The remediation half lives in
`finalize_deliverable.py`; the orchestrator invokes detection,
optionally remediation + re-detection, then halts.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path

# Sibling-tool imports. assemble_pptx is heavy — we only need its
# _resolve_asset_path helper + a few constants, so defer the import
# to inside the gate functions that need it.
_TOOLS_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(_TOOLS_DIR))


SCHEMA_VERSION = "deliverable-validation.v1"

# Severity vocabulary — frozen set; projectable telemetry token.
SEVERITY_P0 = "P0"        # blocks handoff (a remediable defect)
SEVERITY_P1 = "P1"        # standard finding (consider remediation)
SEVERITY_ADVISORY = "advisory"  # surfaced, never blocks
SEVERITIES = (SEVERITY_P0, SEVERITY_P1, SEVERITY_ADVISORY)

# Gate vocabulary — frozen set; projectable telemetry token.
GATES = (
    "placeholder_or_leaked_template",
    "image_completeness",
    "slide_count_vs_budget",
    "mode_vs_user_intent",
    "figure_integrity",
    "figure_path_resolution",
)

# Remediation-kind vocabulary — frozen; projectable telemetry token.
REMEDIATION_AUTO = "auto"
REMEDIATION_TARGETED = "targeted"
REMEDIATION_ADVISORY = "advisory"
REMEDIATION_KINDS = (REMEDIATION_AUTO, REMEDIATION_TARGETED, REMEDIATION_ADVISORY)

# Auto-action vocabulary (only meaningful when kind == "auto").
AUTO_REASSEMBLE = "reassemble"
AUTO_POPULATE_TITLE = "populate_title_from_beril"
AUTO_STRIP_DIRNAME_TOKEN = "strip_dirname_token"

# Figure aspect-ratio tolerance for G5. 5% slop accommodates EMU
# pixel rounding without false-positiving on pillarboxed figures
# (which preserve AR exactly by construction in v1.1.1+ _add_picture).
ASPECT_SKEW_TOLERANCE = 0.05


# ---------------------------------------------------------------------------
# Schema dataclasses
# ---------------------------------------------------------------------------


@dataclass
class Remediation:
    """How (if at all) finalize_deliverable should respond to a finding."""
    kind: str                          # REMEDIATION_KINDS
    action: str | None = None       # auto-action key (kind=auto)
    command: str | None = None      # exact one-stage cmd (kind=targeted)
    note: str | None = None         # operator-readable hint

    def to_dict(self) -> dict:
        return {
            "kind": self.kind,
            "action": self.action,
            "command": self.command,
            "note": self.note,
        }


@dataclass
class Finding:
    """One deliverable-validation finding."""
    id: str                                # stable per-finding id (gate:slot)
    gate: str                              # GATES
    severity: str                          # SEVERITIES
    slide_id_or_target: str | None      # 1-indexed slide id, or
                                           # 'spec'/'deck'/'manifest', or None
    message: str                           # human-readable, free-text
    remediation: Remediation

    def to_dict(self) -> dict:
        return {
            "id": self.id,
            "gate": self.gate,
            "severity": self.severity,
            "slide_id_or_target": self.slide_id_or_target,
            "message": self.message,
            "remediation": self.remediation.to_dict(),
        }


# ---------------------------------------------------------------------------
# Path helpers (mirror assemble_pptx + DraftPaths conventions)
# ---------------------------------------------------------------------------


def _derive_project_dir(draft_dir: Path) -> Path | None:
    """draft_N/ → project_dir (../../). None when the layout doesn't
    match `projects/<id>/talks/draft_N/`."""
    parts = draft_dir.resolve().parts
    if len(parts) < 4:
        return None
    if parts[-2] != "talks":
        return None
    return Path(*parts[:-2])


def _project_dir_token(draft_dir: Path) -> str:
    """The project directory name (e.g. `caulobacter_fur_lipida_loss`)
    used by G1 to detect leaked dir-name tokens in title/presenter."""
    project_dir = _derive_project_dir(draft_dir)
    return project_dir.name if project_dir else ""


def _resolve_figure_path(
    rel_or_abs: str, draft_dir: Path,
) -> Path | None:
    """Mirror of assemble_pptx._resolve_asset_path's lookup order
    (draft_dir first, then project_dir). Pure — no warning side
    effects. Returns the resolved path or None."""
    p = Path(rel_or_abs)
    if p.is_absolute():
        return p if p.is_file() else None
    cand_draft = draft_dir / p
    if cand_draft.is_file():
        return cand_draft
    project_dir = _derive_project_dir(draft_dir)
    if project_dir is not None:
        cand_project = project_dir / p
        if cand_project.is_file():
            return cand_project
    return None


def _load_slide_spec(draft_dir: Path) -> dict | None:
    """Read working/slide_spec.json (or fall back to flat path per
    the v1.1.1 assemble path-resolution fix). None on missing/malformed."""
    for candidate in (
        draft_dir / "working" / "slide_spec.json",
        draft_dir / "slide_spec.json",
    ):
        if not candidate.is_file():
            continue
        try:
            data = json.loads(candidate.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
        if isinstance(data, dict):
            return data
    return None


def _load_beril_yaml(draft_dir: Path) -> dict | None:
    """Read <project_dir>/beril.yaml. Returns the parsed mapping or None.

    Soft-import yaml — the package depends on it but we want a graceful
    None if it's missing (gate emits a finding via the absent-title
    path rather than crashing). Pure stdlib has no YAML reader."""
    project_dir = _derive_project_dir(draft_dir)
    if project_dir is None:
        return None
    yaml_path = project_dir / "beril.yaml"
    if not yaml_path.is_file():
        return None
    try:
        import yaml  # type: ignore[import-untyped]
    except ImportError:
        return None
    try:
        data = yaml.safe_load(yaml_path.read_text(encoding="utf-8"))
    except (OSError, yaml.YAMLError):  # type: ignore[attr-defined]
        return None
    return data if isinstance(data, dict) else None


# ---------------------------------------------------------------------------
# G1 placeholder_or_leaked_template
# ---------------------------------------------------------------------------

# TBD patterns we treat as placeholders. Match `TBD`, `TBD - …`, `TBD:`,
# `tbd` (case-insensitive), as whole-word OR prefix-of-line forms.
_TBD_RE = re.compile(r"\b[Tt][Bb][Dd]\b")


def _is_tbd(value: str | None) -> bool:
    """True iff `value` is None, empty/whitespace-only, or contains a
    TBD token. Used for required-string slots only."""
    if value is None:
        return True
    stripped = value.strip()
    if not stripped:
        return True
    return bool(_TBD_RE.search(stripped))


def _contains_dirname_token(
    value: str | None, dirname_token: str,
) -> bool:
    """True iff `value` contains the project directory name as a
    word (case-insensitive). The caulobacter `Lipida` typo came from
    `caulobacter_fur_lipida_loss` — substring-match catches the
    underscored form AND any tokenization the LLM does of it."""
    if not value or not dirname_token:
        return False
    # Normalize: lowercase + replace common separators (_ - .) with
    # spaces, then check token-by-token AND substring (an LLM might
    # have dropped 'lipida' verbatim from the dir-name's `_lipida_`
    # segment without otherwise mentioning the project).
    norm_value = value.lower()
    norm_token = dirname_token.lower()
    # Direct substring: catches `caulobacter_fur_lipida_loss` mentions.
    if norm_token in norm_value:
        return True
    # Token-by-token: catches per-segment leaks (`lipida` from
    # `caulobacter_fur_lipida_loss`).
    tokens = re.split(r"[_\-.\s]+", norm_token)
    # Filter short stopword-like tokens — single chars + 'fur'/'loss'
    # type 3-letter fragments are too noisy. The threshold is
    # generous; segments worth flagging (a project, organism, or
    # condition word) are usually ≥5 chars.
    distinctive = [t for t in tokens if len(t) >= 5]
    if not distinctive:
        return False
    value_words = set(re.split(r"[^A-Za-z0-9]+", norm_value))
    return any(t in value_words for t in distinctive)


def check_g1_placeholder_or_leaked_template(
    draft_dir: Path, spec: dict | None,
) -> list[Finding]:
    """G1: title + presenter present (not TBD/blank); no project
    directory name token in either; acknowledgments contributors not
    TBD. Title/presenter remediable from beril.yaml; dir-token
    remediable by stripping. Spec absent → return one P0 finding
    (G3/G4 will also fire; G1 still wants to be heard)."""
    if spec is None:
        return [Finding(
            id="g1:no_slide_spec",
            gate=GATES[0],
            severity=SEVERITY_P0,
            slide_id_or_target="spec",
            message=(
                "G1: slide_spec.json not found under draft_dir; cannot "
                "validate title/presenter fields."
            ),
            remediation=Remediation(
                kind=REMEDIATION_TARGETED,
                command=(
                    "beril-presentation-maker continue <draft_dir> "
                    "--resume-from merge"
                ),
                note=(
                    "Re-run the merge stage to produce slide_spec.json "
                    "from the per-substory fragments."
                ),
            ),
        )]

    findings: list[Finding] = []
    dirname_token = _project_dir_token(draft_dir)

    # Title slide is the first layout=='title' slide, by convention.
    title_slide = next(
        (s for s in spec.get("slides", [])
         if s.get("layout") == "title"),
        None,
    )
    if title_slide is None:
        findings.append(Finding(
            id="g1:title_slide_missing",
            gate=GATES[0],
            severity=SEVERITY_P0,
            slide_id_or_target="spec",
            message="G1: no layout='title' slide found in slide_spec.",
            remediation=Remediation(
                kind=REMEDIATION_ADVISORY,
                note=(
                    "The deck has no title slide; this is a deeper "
                    "merge/compose issue, not auto-remediable from the "
                    "validator."
                ),
            ),
        ))
    else:
        content = title_slide.get("content", {}) or {}
        slide_id = str(title_slide.get("position") or "title")

        # Title TBD / blank.
        title_val = content.get("title")
        if _is_tbd(title_val):
            findings.append(Finding(
                id=f"g1:title_tbd:{slide_id}",
                gate=GATES[0],
                severity=SEVERITY_P0,
                slide_id_or_target=slide_id,
                message=(
                    f"G1: title slide has TBD/blank title "
                    f"(got: {title_val!r})."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_AUTO,
                    action=AUTO_POPULATE_TITLE,
                    note=(
                        "Populate title from beril.yaml authors + the "
                        "throughline narrative; re-run assemble."
                    ),
                ),
            ))
        elif _contains_dirname_token(title_val, dirname_token):
            findings.append(Finding(
                id=f"g1:title_dirname_leak:{slide_id}",
                gate=GATES[0],
                severity=SEVERITY_P0,
                slide_id_or_target=slide_id,
                message=(
                    f"G1: title contains project-directory token "
                    f"(dirname={dirname_token!r}, title={title_val!r}). "
                    f"Likely an LLM hallucination of the project name "
                    f"as a science term."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_AUTO,
                    action=AUTO_STRIP_DIRNAME_TOKEN,
                    note=(
                        "Strip dir-name tokens from the title; "
                        "re-run assemble."
                    ),
                ),
            ))

        # Presenter TBD / blank.
        presenter_val = content.get("presenter")
        if _is_tbd(presenter_val):
            findings.append(Finding(
                id=f"g1:presenter_tbd:{slide_id}",
                gate=GATES[0],
                severity=SEVERITY_P0,
                slide_id_or_target=slide_id,
                message=(
                    f"G1: title slide has TBD/blank presenter "
                    f"(got: {presenter_val!r})."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_AUTO,
                    action=AUTO_POPULATE_TITLE,
                    note=(
                        "Populate presenter from beril.yaml authors[0] "
                        "(name + affiliation); re-run assemble."
                    ),
                ),
            ))

    # Acknowledgments contributors not all-TBD.
    ack_slide = next(
        (s for s in spec.get("slides", [])
         if s.get("layout") == "acknowledgments"),
        None,
    )
    if ack_slide is not None:
        ack_content = ack_slide.get("content", {}) or {}
        contributors = ack_content.get("contributors") or []
        ack_slide_id = str(ack_slide.get("position") or "acknowledgments")
        if not contributors or all(_is_tbd(c) for c in contributors):
            findings.append(Finding(
                id=f"g1:acknowledgments_tbd:{ack_slide_id}",
                gate=GATES[0],
                severity=SEVERITY_P1,
                slide_id_or_target=ack_slide_id,
                message=(
                    "G1: acknowledgments slide has no real contributors "
                    f"(got: {contributors!r})."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_AUTO,
                    action=AUTO_POPULATE_TITLE,  # same beril.yaml source
                    note=(
                        "Populate contributors from beril.yaml authors; "
                        "re-run assemble."
                    ),
                ),
            ))

    return findings


# ---------------------------------------------------------------------------
# G2 image_completeness
# ---------------------------------------------------------------------------


def check_g2_image_completeness(
    draft_dir: Path, spec: dict | None,
) -> list[Finding]:
    """G2: every image request resolves to an embedded PNG or an
    explicit manifest reject/skip. Widens v1.1.1's stage-tail check
    to the deliverable boundary + checks the embed (not just the
    PNG-on-disk). Reuses image_gen_orchestrate.find_unresolved_requests
    so the policy is identical.

    Three sub-checks, run independently (a deck can hit any combination):
      a) on-disk orphan requests (a request without a PNG and without
         a manifest reject/skip);
      b) zero-image advisory (image_gen ran + produced no approvals);
      c) spec-side unbound concept_illustration placeholders ({TBD}
         in image_path), regardless of whether image_gen ran.
    """
    findings: list[Finding] = []

    image_requests_dir = draft_dir / "working" / "05_image_requests"
    images_dir = draft_dir / "working" / "05_images"

    request_files: list[Path] = []
    if image_requests_dir.is_dir():
        request_files = sorted(image_requests_dir.glob("*_request.json"))

    # 1. Find on-disk orphans (request without PNG and without manifest
    #    reject/skip). Same routine v1.1.1 hotfix shipped. Only runs
    #    when there are request files at all.
    if request_files:
        try:
            import draft_paths as dp  # noqa: E402
            import image_gen_orchestrate as igo  # noqa: E402
            paths = dp.DraftPaths.from_draft_dir(draft_dir)
            orphan_sids = igo.find_unresolved_requests(paths)
        except Exception as exc:  # pragma: no cover — defensive
            return [Finding(
                id="g2:orchestrate_import_failed",
                gate=GATES[1],
                severity=SEVERITY_P1,
                slide_id_or_target="deck",
                message=(
                    f"G2: could not load image_gen_orchestrate helpers "
                    f"({exc!r}); skipping image-completeness check."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_ADVISORY,
                    note="Engine-side bug; not auto-remediable.",
                ),
            )]

        for sid in orphan_sids:
            findings.append(Finding(
                id=f"g2:orphan_request:{sid}",
                gate=GATES[1],
                severity=SEVERITY_P0,
                slide_id_or_target=sid,
                message=(
                    f"G2: image request for slide {sid} resolved neither "
                    f"to a PNG nor to a manifest reject/skip entry "
                    f"(image_requested_but_not_produced)."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_TARGETED,
                    command=(
                        f"beril-presentation-maker continue {draft_dir} "
                        f"--resume-from image_gen --auto-approve-images"
                    ),
                    note=(
                        "Re-run image_gen — the v1.1.1 stage-tail check "
                        "will now fail loud if generation silently misses "
                        "again; combined with this gate, that closes the "
                        "DP3 loop."
                    ),
                ),
            ))

    # 2. Advisory: image_gen ran with a non-empty manifest but produced
    #    zero approvals. Often legit (everything rejected) but worth
    #    surfacing on the user-visible summary.
    manifest_path = images_dir / "manifest.json"
    if manifest_path.is_file():
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            entries = manifest.get("entries") or []
            n_approved = sum(
                1 for e in entries if isinstance(e, dict)
                and e.get("approved") is True
            )
            if request_files and n_approved == 0:
                findings.append(Finding(
                    id="g2:zero_images_approved",
                    gate=GATES[1],
                    severity=SEVERITY_ADVISORY,
                    slide_id_or_target="deck",
                    message=(
                        f"G2: image_gen ran ({len(request_files)} "
                        f"request(s) authored) but zero images approved "
                        f"({len(entries)} manifest entries). Deck has no "
                        f"AI illustrations — may be intended."
                    ),
                    remediation=Remediation(
                        kind=REMEDIATION_ADVISORY,
                        note=(
                            "Inspect working/05_image_decisions.json + "
                            "audit/image_provenance.json for the per-slide "
                            "decline reason. To regen with relaxed style: "
                            "continue --resume-from image_gen "
                            "--image-allow-exploratory."
                        ),
                    ),
                ))
        except (OSError, json.JSONDecodeError):
            pass

    # 3. Spec-side: a slide with layout='concept_illustration' should
    #    have an image_path that resolves. If it's still {TBD} or
    #    points at a non-existent file, that's a missed embed even if
    #    the manifest path looked OK. (Cheaply detected here; G6 also
    #    covers it as a path-resolution finding, but G2's framing is
    #    image-specific.)
    if spec is not None:
        for slide in spec.get("slides", []):
            if slide.get("layout") != "concept_illustration":
                continue
            content = slide.get("content", {}) or {}
            image_path = content.get("image_path") or ""
            slide_id = str(slide.get("position") or slide.get("id") or "?")
            if image_path == "{TBD}" or not image_path:
                findings.append(Finding(
                    id=f"g2:unbound_image_placeholder:{slide_id}",
                    gate=GATES[1],
                    severity=SEVERITY_P0,
                    slide_id_or_target=slide_id,
                    message=(
                        f"G2: concept_illustration slide {slide_id} has "
                        f"unbound image_path={image_path!r}. The "
                        f"image-gen → bind step did not complete for "
                        f"this slide."
                    ),
                    remediation=Remediation(
                        kind=REMEDIATION_TARGETED,
                        command=(
                            f"beril-presentation-maker continue {draft_dir} "
                            f"--resume-from image_gen "
                            f"--auto-approve-images"
                        ),
                        note=(
                            "Re-run image_gen so mutate-fragment-bind "
                            "writes the resolved image_path."
                        ),
                    ),
                ))

    return findings


# ---------------------------------------------------------------------------
# G3 slide_count_vs_budget
# ---------------------------------------------------------------------------


def _load_mode_budgets() -> dict[str, tuple[int, int]]:
    """Read MODE_SLIDE_BUDGETS from validate_presentation. Avoids
    duplicating the (min, max) table — single source of truth."""
    try:
        import validate_presentation as vp  # noqa: E402
        return dict(vp.MODE_SLIDE_BUDGETS)
    except Exception:  # pragma: no cover — defensive
        return {}


def check_g3_slide_count_vs_budget(
    draft_dir: Path, spec: dict | None, resolved_mode: str | None,
) -> list[Finding]:
    """G3: final slide count within the mode budget. Advisory only —
    a deck slightly under/over budget is not a P0; we surface it so
    the operator can decide. The expected use case is the caulobacter
    37-slide outcome that lay at the talk-30 budget's high edge
    (under-sized for talk-45's [35,48] floor when revealed)."""
    if spec is None:
        return []  # G1 already complained
    n_slides = len(spec.get("slides") or [])
    if n_slides == 0:
        return [Finding(
            id="g3:empty_deck",
            gate=GATES[2],
            severity=SEVERITY_P0,
            slide_id_or_target="deck",
            message="G3: deck has zero slides.",
            remediation=Remediation(
                kind=REMEDIATION_ADVISORY,
                note="Deeper issue; re-run from substory_design.",
            ),
        )]

    budgets = _load_mode_budgets()
    mode = resolved_mode or spec.get("mode")
    if not mode or mode not in budgets:
        return [Finding(
            id="g3:unknown_mode",
            gate=GATES[2],
            severity=SEVERITY_ADVISORY,
            slide_id_or_target="deck",
            message=(
                f"G3: mode {mode!r} has no recorded slide budget; "
                f"deck has {n_slides} slides."
            ),
            remediation=Remediation(
                kind=REMEDIATION_ADVISORY,
                note="Add the mode to MODE_SLIDE_BUDGETS if intended.",
            ),
        )]

    lo, hi = budgets[mode]
    if lo <= n_slides <= hi:
        return []

    direction = "under" if n_slides < lo else "over"
    return [Finding(
        id=f"g3:slide_count_out_of_band:{direction}",
        gate=GATES[2],
        severity=SEVERITY_ADVISORY,
        slide_id_or_target="deck",
        message=(
            f"G3: deck has {n_slides} slides; mode {mode!r} budget is "
            f"[{lo},{hi}] ({direction}-band). May indicate the deck "
            f"was sized for a different mode (see G4) or the "
            f"compose stage under/over-produced."
        ),
        remediation=Remediation(
            kind=REMEDIATION_ADVISORY,
            note=(
                "If mode is wrong: start a fresh draft with --mode "
                "<intended>. If compose under-produced: inspect "
                "working/03_slides/ for missing fragments."
            ),
        ),
    )]


# ---------------------------------------------------------------------------
# G4 mode_vs_user_intent
# ---------------------------------------------------------------------------


def check_g4_mode_vs_user_intent(draft_dir: Path) -> list[Finding]:
    """G4: every artifact recording `mode` agrees with the user's
    persisted explicit pick (audit/user_intent.json). This is the
    Cycle-1 widening of the v1.1.1 cross-artifact check — the
    user_intent.json substrate (DP9b) is what makes "vs. user intent"
    possible at all. Missing user_intent.json → advisory (legacy
    drafts pre-v1.2.0); present-but-not-explicit → advisory (the
    user inherited a default and shouldn't be blamed for it)."""
    findings: list[Finding] = []

    try:
        import user_intent  # noqa: E402
    except Exception:  # pragma: no cover
        return [Finding(
            id="g4:user_intent_import_failed",
            gate=GATES[3],
            severity=SEVERITY_ADVISORY,
            slide_id_or_target="deck",
            message="G4: could not load user_intent helper.",
            remediation=Remediation(kind=REMEDIATION_ADVISORY),
        )]

    user_mode = user_intent.read_field(draft_dir, "mode")
    user_explicit = user_intent.field_was_explicit(draft_dir, "mode")

    if user_mode is None:
        return [Finding(
            id="g4:no_user_intent",
            gate=GATES[3],
            severity=SEVERITY_ADVISORY,
            slide_id_or_target="deck",
            message=(
                "G4: audit/user_intent.json missing — this draft was "
                "created before v1.2.0 or the user_intent write was "
                "skipped. Cannot validate mode against user intent."
            ),
            remediation=Remediation(
                kind=REMEDIATION_ADVISORY,
                note=(
                    "Legacy pre-v1.2.0 drafts won't have this file. "
                    "Future drafts get it automatically on the first "
                    "orchestrator entry."
                ),
            ),
        )]

    if not user_explicit:
        return [Finding(
            id="g4:user_mode_not_explicit",
            gate=GATES[3],
            severity=SEVERITY_ADVISORY,
            slide_id_or_target="deck",
            message=(
                f"G4: user_intent records mode={user_mode!r} but it was "
                f"NOT explicitly set (CLI default). Skipping the "
                f"intent-vs-artifact check; the v1.1.1 cross-artifact "
                f"consistency check still applies."
            ),
            remediation=Remediation(kind=REMEDIATION_ADVISORY),
        )]

    # Reuse mode_consistency for the artifact walk — same source of
    # truth, just compared against user_mode instead of an
    # orchestrator-supplied run_mode.
    try:
        import mode_consistency as mc  # noqa: E402
        mismatches = mc.check_mode_consistency(draft_dir, run_mode=user_mode)
    except Exception:  # pragma: no cover
        mismatches = []

    for msg in mismatches:
        # mc finding string looks like:
        #   "<artifact_label>: mode='X', expected run mode='Y'"
        # Extract the artifact label for the finding's slide_id_or_target slot.
        label = msg.split(":", 1)[0].strip() if ":" in msg else "artifact"
        findings.append(Finding(
            id=f"g4:mode_mismatch:{label}",
            gate=GATES[3],
            severity=SEVERITY_P0,
            slide_id_or_target=label,
            message=(
                f"G4: {msg}. User picked {user_mode!r} explicitly; the "
                f"artifact disagrees — this is the silent-uniformly-wrong "
                f"failure mode v1.1.1's cross-artifact check could not "
                f"catch (DP9b)."
            ),
            remediation=Remediation(
                kind=REMEDIATION_TARGETED,
                command=(
                    f"beril-presentation-maker continue {draft_dir} "
                    f"--resume-from <stage> --mode {user_mode}"
                ),
                note=(
                    "Re-run the affected stage with the correct --mode. "
                    "For 05_image_decisions.json: --resume-from image_gen. "
                    "For qa_anticipated.json: --resume-from qa_prep. "
                    "For slide_spec.json itself: requires a full re-merge "
                    "(--resume-from merge) — the spec is downstream of "
                    "throughline + substory_design, which were sized for "
                    "the wrong mode."
                ),
            ),
        ))

    return findings


# ---------------------------------------------------------------------------
# G5 figure_integrity (display vs. native aspect)
# ---------------------------------------------------------------------------


def check_g5_figure_integrity(draft_dir: Path) -> list[Finding]:
    """G5: every embedded picture's display AR matches the native
    image AR within tolerance. Walks the rendered pptx via python-pptx.
    Guards the v1.1.1/DP4 fix against regression. Pure: open + read."""
    findings: list[Finding] = []
    deck_path = draft_dir / "deliverable" / "draft.pptx"
    if not deck_path.is_file():
        # Fall back to legacy flat layout (assemble.py supports both).
        deck_path = draft_dir / "draft.pptx"
    if not deck_path.is_file():
        return [Finding(
            id="g5:no_deck",
            gate=GATES[4],
            severity=SEVERITY_P0,
            slide_id_or_target="deck",
            message=(
                f"G5: no rendered deck at "
                f"{draft_dir/'deliverable'/'draft.pptx'} or "
                f"{draft_dir/'draft.pptx'} — cannot check figure "
                f"aspect ratios."
            ),
            remediation=Remediation(
                kind=REMEDIATION_AUTO,
                action=AUTO_REASSEMBLE,
                note="Re-run assemble to produce the deck.",
            ),
        )]

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError:
        return [Finding(
            id="g5:pptx_unavailable",
            gate=GATES[4],
            severity=SEVERITY_ADVISORY,
            slide_id_or_target="deck",
            message="G5: python-pptx unavailable; cannot check picture aspect.",
            remediation=Remediation(kind=REMEDIATION_ADVISORY),
        )]

    try:
        prs = Presentation(str(deck_path))
    except Exception as exc:  # pragma: no cover
        return [Finding(
            id="g5:deck_unreadable",
            gate=GATES[4],
            severity=SEVERITY_P0,
            slide_id_or_target="deck",
            message=f"G5: cannot open deck ({exc!r}).",
            remediation=Remediation(kind=REMEDIATION_ADVISORY),
        )]

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            disp_w = shape.width
            disp_h = shape.height
            if disp_w <= 0 or disp_h <= 0:
                continue
            try:
                native_w_px, native_h_px = shape.image.size
            except Exception:  # pragma: no cover
                continue
            if native_w_px <= 0 or native_h_px <= 0:
                continue
            disp_ar = disp_w / disp_h
            native_ar = native_w_px / native_h_px
            ratio = disp_ar / native_ar
            if (1.0 - ASPECT_SKEW_TOLERANCE) <= ratio <= (1.0 + ASPECT_SKEW_TOLERANCE):
                continue
            findings.append(Finding(
                id=f"g5:aspect_skew:{slide_idx}:{shape.shape_id}",
                gate=GATES[4],
                severity=SEVERITY_P0,
                slide_id_or_target=str(slide_idx),
                message=(
                    f"G5: picture on slide {slide_idx} skewed: "
                    f"display AR {disp_ar:.3f}, native AR "
                    f"{native_ar:.3f}, ratio {ratio:.3f} (tolerance "
                    f"±{ASPECT_SKEW_TOLERANCE*100:.0f}%). DP4 hotfix "
                    f"contract: _add_picture must pass width-only "
                    f"and clamp height; a width-and-height call "
                    f"reintroduces this defect."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_AUTO,
                    action=AUTO_REASSEMBLE,
                    note=(
                        "Re-running assemble against a v1.1.1+ "
                        "_add_picture will produce a correct-aspect "
                        "render of the same spec. If the skew persists "
                        "after re-assemble, the regression is in "
                        "_add_picture itself."
                    ),
                ),
            ))
    return findings


# ---------------------------------------------------------------------------
# G6 figure_path_resolution
# ---------------------------------------------------------------------------

# Spec slot names that hold a figure/image path. Mirrors the slots
# assemble_pptx._resolve_asset_path is called against; if a new
# slot is added there, add it here too.
_FIGURE_SLOTS = ("figure", "supporting_graphic", "image_path")


def check_g6_figure_path_resolution(
    draft_dir: Path, spec: dict | None,
) -> list[Finding]:
    """G6: every spec figure path resolves on disk. Skips `{TBD}`
    placeholders (G2 owns those) and anything that's clearly not a
    figure path (e.g., an inline image_prompt blob)."""
    if spec is None:
        return []
    findings: list[Finding] = []
    for slide in spec.get("slides", []):
        content = slide.get("content", {}) or {}
        slide_id = str(slide.get("position") or slide.get("id") or "?")
        for slot in _FIGURE_SLOTS:
            value = content.get(slot)
            if not isinstance(value, str) or not value:
                continue
            if value == "{TBD}":
                continue
            resolved = _resolve_figure_path(value, draft_dir)
            if resolved is not None:
                continue
            findings.append(Finding(
                id=f"g6:unresolved_figure:{slide_id}:{slot}",
                gate=GATES[5],
                severity=SEVERITY_P0,
                slide_id_or_target=slide_id,
                message=(
                    f"G6: slide {slide_id} content.{slot}={value!r} "
                    f"did not resolve under draft_dir or project_dir."
                ),
                remediation=Remediation(
                    kind=REMEDIATION_ADVISORY,
                    note=(
                        "Inspect the spec entry: either the path is "
                        "wrong (fix in slide_spec.json) or the file "
                        "was moved/deleted (restore the file). Not "
                        "auto-remediable — the curator wrote the bad "
                        "reference; we can't guess the intended file."
                    ),
                ),
            ))
    return findings


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------


def validate(draft_dir: Path) -> list[Finding]:
    """Run all six gates against `draft_dir`. Pure: read-only over
    the filesystem. Returns the flat list of findings."""
    spec = _load_slide_spec(draft_dir)

    # Resolve the mode to use for G3 budget — prefer user_intent's
    # explicit pick (DP9b), fall back to slide_spec, then "unknown".
    resolved_mode: str | None = None
    try:
        import user_intent  # noqa: E402
        ui_mode = user_intent.read_field(draft_dir, "mode")
        if ui_mode is not None and user_intent.field_was_explicit(
                draft_dir, "mode"):
            resolved_mode = ui_mode
    except Exception:
        pass
    if resolved_mode is None and spec is not None:
        resolved_mode = spec.get("mode")

    findings: list[Finding] = []
    findings.extend(check_g1_placeholder_or_leaked_template(draft_dir, spec))
    findings.extend(check_g2_image_completeness(draft_dir, spec))
    findings.extend(check_g3_slide_count_vs_budget(
        draft_dir, spec, resolved_mode))
    findings.extend(check_g4_mode_vs_user_intent(draft_dir))
    findings.extend(check_g5_figure_integrity(draft_dir))
    findings.extend(check_g6_figure_path_resolution(draft_dir, spec))
    return findings


def write_findings(draft_dir: Path, findings: Iterable[Finding]) -> Path:
    """Persist findings to `audit/deliverable_validation.json`."""
    findings_list = list(findings)
    summary = _summarize(findings_list)
    audit_dir = draft_dir / "audit"
    audit_dir.mkdir(parents=True, exist_ok=True)
    out = audit_dir / "deliverable_validation.json"
    payload = {
        "schema_version": SCHEMA_VERSION,
        "draft_dir": str(draft_dir),
        "summary": summary,
        "findings": [f.to_dict() for f in findings_list],
    }
    out.write_text(
        json.dumps(payload, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    return out


def _summarize(findings: list[Finding]) -> dict:
    """Aggregate counters (telemetry-ready: keys are tokens)."""
    by_gate: dict[str, int] = {g: 0 for g in GATES}
    by_severity: dict[str, int] = {s: 0 for s in SEVERITIES}
    by_kind: dict[str, int] = {k: 0 for k in REMEDIATION_KINDS}
    for f in findings:
        by_gate[f.gate] = by_gate.get(f.gate, 0) + 1
        by_severity[f.severity] = by_severity.get(f.severity, 0) + 1
        by_kind[f.remediation.kind] = by_kind.get(f.remediation.kind, 0) + 1
    return {
        "total": len(findings),
        "by_gate": by_gate,
        "by_severity": by_severity,
        "by_remediation_kind": by_kind,
        "blocking": by_severity.get(SEVERITY_P0, 0),
    }


def readiness_exit_code(findings: list[Finding]) -> int:
    """Map findings to a readiness exit code:
      0 — clean OR only advisory findings (deliverable ready to hand off).
      1 — at least one P0 or P1 finding; remediation needed (the never-
          discard policy means the deliverable is STILL produced;
          finalize_deliverable.py runs auto-remediation + re-validates).
    """
    for f in findings:
        if f.severity in (SEVERITY_P0, SEVERITY_P1):
            return 1
    return 0


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def _cmd_check(args: argparse.Namespace) -> int:
    draft_dir = Path(args.draft_dir).resolve()
    if not draft_dir.is_dir():
        print(
            f"validate_deliverable: draft_dir not found: {draft_dir}",
            file=sys.stderr,
        )
        return 2
    findings = validate(draft_dir)
    out = write_findings(draft_dir, findings)
    summary = _summarize(findings)
    print(
        f"validate_deliverable: {summary['total']} finding(s) "
        f"(P0={summary['by_severity'].get(SEVERITY_P0, 0)}, "
        f"P1={summary['by_severity'].get(SEVERITY_P1, 0)}, "
        f"advisory={summary['by_severity'].get(SEVERITY_ADVISORY, 0)}); "
        f"wrote {out}",
        file=sys.stderr,
    )
    if args.print_findings and findings:
        for f in findings:
            print(
                f"  [{f.severity}] {f.gate} ({f.id}): {f.message}",
                file=sys.stderr,
            )
            note = f.remediation.note
            cmd = f.remediation.command
            if cmd:
                print(f"      → run: {cmd}", file=sys.stderr)
            elif note:
                print(f"      → {f.remediation.kind}: {note}", file=sys.stderr)
    return readiness_exit_code(findings)


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(
        prog="validate_deliverable",
        description=(
            "Cycle 1 pre-handoff deterministic gate. Six checks over "
            "the produced deliverable + working artifacts; emits "
            "findings under deliverable-validation.v1."
        ),
    )
    sub = p.add_subparsers(dest="cmd", required=True)

    p_chk = sub.add_parser(
        "check",
        help="Run all six gates; write audit/deliverable_validation.json.",
    )
    p_chk.add_argument("draft_dir")
    p_chk.add_argument(
        "--print-findings", action="store_true",
        help=(
            "Echo each finding to stderr in addition to writing the "
            "JSON (orchestrator-friendly default-off; useful at the "
            "shell prompt)."
        ),
    )
    p_chk.set_defaults(func=_cmd_check)

    args = p.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())

#!/usr/bin/env python3
"""poster_fill.py — fill the KBase poster templates with content.

Per SPEC §12, §15 + DECISIONS D-013. Posters are NOT animated, NOT
speaker-noted, NOT interactive. They are a single-slide deliverable
filled from a poster_spec dict. Two templates ship under
`references/templates/`:

  - kbase-poster-horizontal.pptx (48 × 36 in, landscape)
  - kbase-poster-vertical.pptx   (36 × 48 in, portrait)

Each template has pre-positioned title / authors / KBase logo / funding
shapes. This module:

  1. Replaces the "TITLE" text placeholder with the user's title.
  2. Replaces the "NAME 1, NAME 2, NAME 3 INSTITUTION..." placeholder
     with the user's authors + affiliation.
  3. Adds body-content panels for: TL;DR, Methods Summary, 2–4
     figures, Cross-tenant Integration, Implications, References,
     Acknowledgments.

Body layout depends on orientation:
  - Horizontal: 3-column flow (col 1: TL;DR + Methods; col 2: figures;
    col 3: Cross-tenant + Implications + Refs).
  - Vertical:   single column, sections stacked top-to-bottom.

CLI:

    python3 poster_fill.py <poster_spec.json> --out poster.pptx \\
        [--orientation horizontal|vertical]

Library:

    from poster_fill import fill_poster, PosterSpec
    spec = PosterSpec(...)
    fill_poster(spec, "out/poster.pptx", orientation="horizontal")
"""
from __future__ import annotations

import argparse
import dataclasses
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_THIS_DIR = Path(__file__).resolve().parent
TEMPLATE_DIR = _THIS_DIR.parent / "references" / "templates"
POSTER_TEMPLATES = {
    "horizontal": TEMPLATE_DIR / "kbase-poster-horizontal.pptx",
    "vertical":   TEMPLATE_DIR / "kbase-poster-vertical.pptx",
}

GRAPHITE_GRAY = RGBColor(0x9D, 0x93, 0x89)
FRESHWATER_BLUE = RGBColor(0x00, 0x7D, 0xC3)
GRASS_GREEN = RGBColor(0x5E, 0x97, 0x32)


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------

@dataclass
class PosterFigure:
    path: str            # absolute or relative to draft_dir
    caption: str = ""

    @classmethod
    def from_dict(cls, d: dict) -> "PosterFigure":
        return cls(path=d["path"], caption=d.get("caption", ""))


@dataclass
class PosterSpec:
    """All content needed to fill a KBase poster template.

    The orchestrator builds this from the talk's slide_spec by mapping:
      title / authors / affiliation        ← title slide content
      tl_dr                                 ← throughline.punchline
      methods_summary                       ← methods_summary slide bullets
      figures (2–4)                         ← curated figures (mode budget)
      cross_tenant_summary                  ← cross_tenant_integration slide
      implications                          ← implications slide bullets
      references_short                      ← references slide refs_short
      acknowledgments                       ← acknowledgments slide content
      funding                               ← acknowledgments tenant_attrib
    """
    title: str
    authors: str
    affiliation: str = ""
    tl_dr: str = ""
    methods_summary: list[str] = field(default_factory=list)
    figures: list[PosterFigure] = field(default_factory=list)
    cross_tenant_summary: str = ""
    implications: list[str] = field(default_factory=list)
    references_short: list[str] = field(default_factory=list)
    acknowledgments: str = ""
    funding: str = ""

    @classmethod
    def from_dict(cls, d: dict) -> "PosterSpec":
        return cls(
            title=d.get("title", ""),
            authors=d.get("authors", ""),
            affiliation=d.get("affiliation", ""),
            tl_dr=d.get("tl_dr", ""),
            methods_summary=list(d.get("methods_summary", []) or []),
            figures=[PosterFigure.from_dict(f) for f in d.get("figures", []) or []],
            cross_tenant_summary=d.get("cross_tenant_summary", ""),
            implications=list(d.get("implications", []) or []),
            references_short=list(d.get("references_short", []) or []),
            acknowledgments=d.get("acknowledgments", ""),
            funding=d.get("funding", ""),
        )


# ---------------------------------------------------------------------------
# Template-shape replacement
# ---------------------------------------------------------------------------

def _shape_text(shape) -> str:
    """Get the concatenated text content of a shape (for matching)."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text or ""


def _replace_shape_text(shape, new_text: str) -> bool:
    """Replace shape.text_frame.text with new_text. Returns True if shape
    has a text frame."""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    # Preserve the first paragraph's run formatting if possible
    tf.text = new_text
    return True


def _find_template_shape_by_text_substring(
    slide, substring: str,
):
    """Return the first shape whose text contains `substring`
    (case-sensitive, stripped). None if not found."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if substring in shape.text_frame.text:
            return shape
    return None


# ---------------------------------------------------------------------------
# Body-section drawing primitives
# ---------------------------------------------------------------------------

def _add_section(
    slide, *, title: str, body_lines: list[str] | str,
    left_in: float, top_in: float, width_in: float, height_in: float,
    title_pt: int = 36, body_pt: int = 18,
):
    """Add a section panel: heading + content text box. body_lines can be
    a list (rendered as bullets) or a string (rendered as paragraph)."""
    # Section heading
    head_h = title_pt / 72.0 * 1.4  # ~140% line height
    title_box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(head_h),
    )
    title_box.text_frame.text = title
    for para in title_box.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(title_pt)
            run.font.bold = True
            run.font.color.rgb = FRESHWATER_BLUE

    # Body
    body_box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in + head_h + 0.1),
        Inches(width_in), Inches(height_in - head_h - 0.1),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    if isinstance(body_lines, str):
        tf.text = body_lines
    elif isinstance(body_lines, list):
        if not body_lines:
            tf.text = "(empty)"
        else:
            tf.text = "• " + body_lines[0]
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = "• " + line
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(body_pt)
    return body_box


def _add_figure_panel(
    slide, *, figure: PosterFigure, draft_dir: Path,
    left_in: float, top_in: float, width_in: float, height_in: float,
    caption_pt: int = 14,
):
    """Add a figure panel: image + caption."""
    img_path = Path(figure.path)
    if not img_path.is_absolute():
        img_path = draft_dir / img_path
    caption_h = 0.6 if figure.caption else 0.0
    img_h = height_in - caption_h - 0.1
    if img_path.is_file():
        slide.shapes.add_picture(
            str(img_path),
            Inches(left_in), Inches(top_in),
            width=Inches(width_in), height=Inches(img_h),
        )
    else:
        # Missing-file placeholder
        ph = slide.shapes.add_textbox(
            Inches(left_in), Inches(top_in),
            Inches(width_in), Inches(img_h),
        )
        ph.text_frame.text = f"[missing: {figure.path}]"
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = GRAPHITE_GRAY
                run.font.size = Pt(18)
    if figure.caption:
        cap = slide.shapes.add_textbox(
            Inches(left_in), Inches(top_in + img_h + 0.1),
            Inches(width_in), Inches(caption_h),
        )
        cap.text_frame.text = figure.caption
        cap.text_frame.word_wrap = True
        for para in cap.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(caption_pt)
                run.font.color.rgb = GRAPHITE_GRAY


# ---------------------------------------------------------------------------
# Layouts
# ---------------------------------------------------------------------------

def _layout_horizontal(slide, spec: PosterSpec, draft_dir: Path) -> None:
    """3-column flow for 48×36 in landscape posters.

    Body region: y=7.5 to y=32.5 (25" tall), x=0.5 to x=47.5 (47" wide).
    Columns: 15"-wide each, 1" gaps.
    """
    # Column 1: TL;DR + Methods Summary
    col1_x = 0.5
    col_w = 15.0
    _add_section(
        slide, title="TL;DR", body_lines=spec.tl_dr,
        left_in=col1_x, top_in=7.5, width_in=col_w, height_in=11.5,
    )
    _add_section(
        slide, title="Methods", body_lines=spec.methods_summary,
        left_in=col1_x, top_in=20.0, width_in=col_w, height_in=12.0,
    )

    # Column 2: figures
    col2_x = 16.5
    if spec.figures:
        n = min(2, len(spec.figures))   # poster shows up to 2 figures in mid col
        panel_h = (25.0 - 0.5 * (n - 1)) / max(n, 1)
        for i, fig in enumerate(spec.figures[:n]):
            _add_figure_panel(
                slide, figure=fig, draft_dir=draft_dir,
                left_in=col2_x, top_in=7.5 + i * (panel_h + 0.5),
                width_in=col_w, height_in=panel_h,
            )

    # Column 3: Cross-tenant + Implications + References + Acks
    col3_x = 32.5
    _add_section(
        slide, title="Integration",
        body_lines=spec.cross_tenant_summary or "(no cross-tenant signal)",
        left_in=col3_x, top_in=7.5, width_in=col_w, height_in=7.0,
    )
    _add_section(
        slide, title="Implications",
        body_lines=spec.implications or ["(none)"],
        left_in=col3_x, top_in=15.0, width_in=col_w, height_in=8.0,
    )
    _add_section(
        slide, title="References",
        body_lines=spec.references_short or ["(none)"],
        left_in=col3_x, top_in=23.5, width_in=col_w, height_in=5.5,
    )
    if spec.acknowledgments:
        _add_section(
            slide, title="Acknowledgments",
            body_lines=spec.acknowledgments,
            left_in=col3_x, top_in=29.0, width_in=col_w, height_in=3.0,
            body_pt=14,
        )


def _layout_vertical(slide, spec: PosterSpec, draft_dir: Path) -> None:
    """Single-column flow for 36×48 in portrait posters.

    Body region: y=7.0 to y=45.5 (38.5" tall), x=0.5 to x=35.5 (35" wide).
    Sections stacked top-to-bottom.
    """
    x = 0.5
    w = 35.0
    y = 7.0

    _add_section(slide, title="TL;DR", body_lines=spec.tl_dr,
                 left_in=x, top_in=y, width_in=w, height_in=4.0)
    y += 4.5
    _add_section(slide, title="Methods", body_lines=spec.methods_summary,
                 left_in=x, top_in=y, width_in=w, height_in=4.0)
    y += 4.5

    # Figures: 1 or 2 inline
    if spec.figures:
        n = min(2, len(spec.figures))
        fig_h = 8.0
        if n == 1:
            _add_figure_panel(slide, figure=spec.figures[0], draft_dir=draft_dir,
                              left_in=x, top_in=y, width_in=w, height_in=fig_h)
        else:
            half_w = (w - 0.5) / 2
            for i, fig in enumerate(spec.figures[:2]):
                _add_figure_panel(slide, figure=fig, draft_dir=draft_dir,
                                  left_in=x + i * (half_w + 0.5),
                                  top_in=y, width_in=half_w, height_in=fig_h)
        y += fig_h + 0.5

    _add_section(slide, title="Integration",
                 body_lines=spec.cross_tenant_summary or "(no cross-tenant signal)",
                 left_in=x, top_in=y, width_in=w, height_in=4.0)
    y += 4.5
    _add_section(slide, title="Implications",
                 body_lines=spec.implications or ["(none)"],
                 left_in=x, top_in=y, width_in=w, height_in=4.0)
    y += 4.5
    _add_section(slide, title="References",
                 body_lines=spec.references_short or ["(none)"],
                 left_in=x, top_in=y, width_in=w, height_in=3.0)
    y += 3.5
    if spec.acknowledgments and y < 44.5:
        _add_section(slide, title="Acknowledgments",
                     body_lines=spec.acknowledgments,
                     left_in=x, top_in=y, width_in=w,
                     height_in=min(3.0, 45.5 - y),
                     body_pt=14)


# ---------------------------------------------------------------------------
# Top-level fill
# ---------------------------------------------------------------------------

def fill_poster(
    spec: PosterSpec,
    out_path: str | Path,
    *,
    orientation: str = "horizontal",
    template_path: str | Path | None = None,
    draft_dir: str | Path | None = None,
) -> Path:
    """Fill the KBase poster template with content from `spec`.

    Args:
      spec: PosterSpec.
      out_path: where to write the .pptx.
      orientation: "horizontal" (48×36) or "vertical" (36×48).
      template_path: override the shipped template.
      draft_dir: directory for resolving relative figure paths.

    Returns:
      Path to the saved .pptx.
    """
    if orientation not in POSTER_TEMPLATES:
        raise ValueError(
            f"unknown orientation {orientation!r}; "
            f"valid: {sorted(POSTER_TEMPLATES.keys())}"
        )
    tpl = Path(template_path) if template_path else POSTER_TEMPLATES[orientation]
    if not tpl.is_file():
        raise FileNotFoundError(f"poster template not found: {tpl}")

    out_path = Path(out_path).resolve()
    draft_dir = Path(draft_dir).resolve() if draft_dir else out_path.parent

    prs = Presentation(tpl)
    slide = prs.slides[0]

    # Replace title placeholder ("TITLE" exactly)
    title_shape = _find_template_shape_by_text_substring(slide, "TITLE")
    if title_shape is not None and spec.title:
        _replace_shape_text(title_shape, spec.title)

    # Replace authors placeholder. The template stores "NAME1, NAME2, NAME3"
    # (no space — runs are formatted separately for the superscript digit).
    # Match on "NAME1" or fall back to any shape containing "NAME".
    authors_shape = (_find_template_shape_by_text_substring(slide, "NAME1")
                     or _find_template_shape_by_text_substring(slide, "NAME 1")
                     or _find_template_shape_by_text_substring(slide, "INSTITUTION"))
    if authors_shape is not None and (spec.authors or spec.affiliation):
        text = spec.authors
        if spec.affiliation:
            text = f"{text}\n{spec.affiliation}"
        _replace_shape_text(authors_shape, text)

    # Replace funding placeholder (matches "Funding:" substring)
    funding_shape = _find_template_shape_by_text_substring(slide, "Funding:")
    if funding_shape is not None and spec.funding:
        _replace_shape_text(funding_shape, f"Funding: {spec.funding}")

    # Add body sections per orientation
    if orientation == "horizontal":
        _layout_horizontal(slide, spec, draft_dir)
    else:
        _layout_vertical(slide, spec, draft_dir)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(
        prog="poster_fill",
        description="Fill a KBase poster template (horizontal or vertical) "
                    "with content from a poster_spec JSON file.",
    )
    p.add_argument("poster_spec",
                   help="Path to poster_spec.json (PosterSpec.to_dict shape).")
    p.add_argument("--out", required=True,
                   help="Output .pptx path.")
    p.add_argument("--orientation", choices=["horizontal", "vertical"],
                   default="horizontal")
    p.add_argument("--template", help="Override the shipped template path.")
    p.add_argument("--draft-dir",
                   help="Directory for resolving relative figure paths "
                        "(default: parent of --out).")
    args = p.parse_args(argv)

    spec_path = Path(args.poster_spec).resolve()
    if not spec_path.is_file():
        print(f"poster_spec not found: {spec_path}", file=sys.stderr)
        return 2

    raw = json.loads(spec_path.read_text(encoding="utf-8"))
    spec = PosterSpec.from_dict(raw)
    try:
        out = fill_poster(spec, args.out, orientation=args.orientation,
                          template_path=args.template, draft_dir=args.draft_dir)
    except (FileNotFoundError, ValueError) as e:
        print(f"poster_fill: {e}", file=sys.stderr)
        return 2
    print(f"wrote {out}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

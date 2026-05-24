#!/usr/bin/env python3
"""diagram_render.py — render slide_spec diagram dicts to python-pptx shapes.

Per SPEC §8.2 + D-014: Tier 2 figure handling. Workflow / architecture /
ER / decision-tree diagrams are rendered as python-pptx NATIVE shapes
(AutoShapes + Connectors + TextBoxes), not as raster images. This keeps
the diagram editable in PowerPoint / Keynote / Slides downstream.

Inputs (one diagram, defined in slide_spec.py):

  diagram = {
    "kind": "boxes_and_arrows",
    "nodes": [
      {"id": "n1", "label": "Start", "shape": "rounded",
       "x": 0.5, "y": 1.0, "w": 1.5, "h": 0.8,
       "fill_color": "freshwater_blue", "text_color": "white"},
      ...
    ],
    "edges": [
      {"from": "n1", "to": "n2", "kind": "straight", "label": "..."},
      ...
    ],
  }

Coordinates are inches with (0, 0) at upper-left of `region` (passed by
caller — the body region of the slide). Brand colors are KBase palette
names; resolved against `kbase-brand-tokens.json`.

Library:

    from diagram_render import render_diagram
    render_diagram(slide, diagram, region=(0.5, 1.4, 9.0, 3.5),
                   brand_tokens=tokens)

Tests live at tests/unit/test_diagram_render.py.
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Emu, Inches, Pt


# ---------------------------------------------------------------------------
# Shape + connector mappings (mirrors slide_spec.DIAGRAM_NODE_SHAPES + DIAGRAM_EDGE_KINDS)
# ---------------------------------------------------------------------------

NODE_SHAPE_MAP: dict[str, int] = {
    "rectangle":      MSO_SHAPE.RECTANGLE,
    "rounded":        MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse":        MSO_SHAPE.OVAL,
    "parallelogram":  MSO_SHAPE.PARALLELOGRAM,
    "cylinder":       MSO_SHAPE.CAN,                 # database-y
    "callout":        MSO_SHAPE.RECTANGULAR_CALLOUT,
    "swimlane":       MSO_SHAPE.RECTANGLE,           # rendered no-fill
}

EDGE_KIND_MAP: dict[str, int] = {
    "straight": MSO_CONNECTOR.STRAIGHT,
    "elbow":    MSO_CONNECTOR.ELBOW,
    "curved":   MSO_CONNECTOR.CURVE,
}


# ---------------------------------------------------------------------------
# Defaults / fallbacks
# ---------------------------------------------------------------------------

# Default fill / text colors when a node doesn't specify (KBase palette).
DEFAULT_NODE_FILL = "freshwater_blue"
DEFAULT_NODE_TEXT = "white"
DEFAULT_SWIMLANE_BORDER = "graphite_gray"

# M4a Tier E round 2 (2026-05-23): connector lines and edge labels
# moved off the brand's `graphite_gray` (#9D9389 — a light tan-gray,
# documented brand identity for soft elements like swimlane borders)
# onto a slate-dark tone that holds contrast against cream backgrounds
# and the new (post-watermark-strip) flat slide background. Tier E
# round-1 visual-QA found the connector lines on slides 6/10/19 were
# essentially invisible against the watermark; the same color value is
# still washed out on the flat cream background. swimlane borders
# (DEFAULT_SWIMLANE_BORDER) intentionally keep graphite_gray — they
# ARE meant to recede.
from pptx.dml.color import RGBColor as _RGBColor
_CONNECTOR_RGB = _RGBColor(80, 75, 70)   # slate-dark; matches assemble_pptx.GRAPHITE_GRAY_RGB

# Hardcoded fallback hexes if brand_tokens not provided. Matches
# kbase-brand-tokens.json (KBase Style Guide June 2022).
_FALLBACK_HEX = {
    "microbe_orange":     "#F78E1E",
    "grass_green":        "#5E9732",
    "freshwater_blue":    "#007DC3",
    "golden_yellow":      "#FFD200",
    "spring_green":       "#C1CD23",
    "ocean_blue":         "#72CCD2",
    "cyanobacteria_teal": "#009688",
    "lupine_purple":      "#66489D",
    "frost_blue":         "#C7DBEE",
    "rainier_cherry_red": "#D2232A",
    "graphite_gray":      "#9D9389",
    "white":              "#FFFFFF",
    "black":              "#000000",
}


@dataclass
class _RenderedShape:
    """Internal — the python-pptx shape we created for a node, plus its
    geometry. Used to compute connector endpoints + the gap-based
    edge-label placement (M4a Tier E round 3, 2026-05-23: the round-2
    visual-QA found edge labels at center-to-center midpoints land ON
    the next node box when inter-node gap is < ~0.7in; the label pass
    now uses gap midpoints + node widths to size the textbox to fit
    the gap)."""
    node_id: str
    shape: Any                     # python-pptx shape object
    cx: float                      # center x (inches)
    cy: float                      # center y (inches)
    w: float = 0.0                 # node width (inches); 0 means unknown
    h: float = 0.0                 # node height (inches); 0 means unknown


# ---------------------------------------------------------------------------
# Brand-color resolution
# ---------------------------------------------------------------------------

def resolve_color(name_or_hex: str | None,
                  brand_tokens: dict | None,
                  default_token: str = DEFAULT_NODE_FILL) -> RGBColor:
    """Resolve a color reference (palette token or '#RRGGBB') to RGBColor.

    Resolution order:
      1. If name_or_hex starts with '#', parse as hex.
      2. If brand_tokens provided, look up by name in palette.primary
         then palette.secondary then palette.neutral.
      3. Fallback to _FALLBACK_HEX.
      4. If still unresolved, fall back to the default_token through the
         same chain.
    """
    candidates = [name_or_hex, default_token]
    for c in candidates:
        if c is None:
            continue
        if isinstance(c, str) and c.startswith("#") and len(c) == 7:
            return RGBColor.from_string(c[1:].upper())
        # Brand token lookup
        if brand_tokens is not None:
            palette = brand_tokens.get("palette", {}) if isinstance(brand_tokens, dict) else {}
            for group in ("primary", "secondary", "neutral"):
                grp = palette.get(group, {})
                if c in grp:
                    hexv = grp[c].get("hex", "")
                    if hexv.startswith("#") and len(hexv) == 7:
                        return RGBColor.from_string(hexv[1:].upper())
        # Hardcoded fallback
        if c in _FALLBACK_HEX:
            return RGBColor.from_string(_FALLBACK_HEX[c][1:].upper())
    # Final fallback: white
    return RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Coordinate transform (diagram → slide)
# ---------------------------------------------------------------------------

def _transform_coords(
    x_in: float, y_in: float,
    region: tuple[float, float, float, float],
) -> tuple[float, float]:
    """Pass through diagram coordinates as ABSOLUTE slide coords.

    2026-04-27 fix #78: previous version added region.top + region.left
    to (x, y), assuming coords were diagram-local (origin at region
    top-left). But:
      - slide_compose.v1.md tells the LLM 'Default region: x in [0.5,
        9.5], y in [1.4, 5.6]' — absolute slide coords.
      - repair_diagram_stubs.compute_linear_geometry produces absolute
        slide coords (CONTENT_LEFT=0.5, CONTENT_TOP=1.4 already
        baked in).
      - The renderer's region.top=1.30 was being ADDED to absolute y
        values like 5.0, producing slide y=6.30 (off-slide; slide
        height is 5.625in). Live failure draft_7 slide 11: 8 nodes
        with y up to 5.0 rendered partly off-slide.

    Treat region as the BOUNDS the diagram should fit within, not an
    origin offset. Coordinates pass through. The region tuple stays
    in the signature for future use (e.g., scaling-to-fit when a
    diagram is generated for a smaller region than the master's
    workflow_diagram body).
    """
    return x_in, y_in


# ---------------------------------------------------------------------------
# M4a Tier A — explicit-fontScale shrink-to-fit for diagram node labels
# ---------------------------------------------------------------------------
#
# LibreOffice does NOT honor a bare <a:normAutofit/> at render — the
# previous `tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE` wrote the
# bare form and the rendered text spilled past the box (ibd_phage_
# targeting draft_1 slides 6/10/19). The only mechanism LibreOffice
# computes is an explicit `<a:normAutofit fontScale="...">`. This
# helper writes that directly into the auto-shape's <a:bodyPr>.
#
# Parallel to assemble_pptx._fit_textbox (DQ3 — 60% fontScale floor);
# duplicated locally because diagram_render is loaded as a standalone
# module by assemble_pptx's importlib loader, so importing back into
# assemble_pptx is circular. If a third consumer arrives in Tier C,
# refactor both into a shared helper module.

# OOXML namespace — must match the one in assemble_pptx._ensure_slide_text_autofit
_DML_NS_DR = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PML_NS_DR = "http://schemas.openxmlformats.org/presentationml/2006/main"

NODE_FONTSCALE_FLOOR = 60000   # DQ3 — 60% of 14pt = 8.4pt at projection
NODE_FONTSCALE_FULL = 100000


def _apply_fontscale_to_shape(
    shape,
    content_chars: int,
    *,
    ladder: tuple[tuple[int, int], ...],
    full_below: int,
    ln_spc_reduction: int = 20000,
) -> int:
    """Write an explicit-fontScale normAutofit onto an auto-shape's bodyPr.

    Args:
      shape: a python-pptx auto-shape (the diagram node).
      content_chars: character count of the label.
      ladder: (char_cap, scale) pairs in ascending char order; the first
        cap >= chars wins. Falls through to NODE_FONTSCALE_FLOOR.
      full_below: chars at or below this render at 100% (no shrink).

    Returns the fontScale written.
    """
    from lxml import etree as _et

    if content_chars <= full_below:
        scale = NODE_FONTSCALE_FULL
    else:
        scale = NODE_FONTSCALE_FLOOR
        for cap, sc in ladder:
            if content_chars <= cap:
                scale = sc
                break

    # Find <p:txBody> under the shape (the auto-shape uses the same
    # namespace structure as a placeholder).
    sp_el = shape.element
    tx_body = sp_el.find(f"{{{_PML_NS_DR}}}txBody")
    if tx_body is None:
        tx_body = sp_el.find(f"{{{_DML_NS_DR}}}txBody")
    if tx_body is None:
        return scale
    body_pr = tx_body.find(f"{{{_DML_NS_DR}}}bodyPr")
    if body_pr is None:
        return scale
    for tag in ("normAutofit", "noAutofit", "spAutoFit"):
        for child in list(body_pr):
            if child.tag == f"{{{_DML_NS_DR}}}{tag}":
                body_pr.remove(child)
    af = _et.SubElement(body_pr, f"{{{_DML_NS_DR}}}normAutofit")
    af.set("fontScale", str(scale))
    af.set("lnSpcReduction", str(ln_spc_reduction))
    return scale


# ---------------------------------------------------------------------------
# Node rendering
# ---------------------------------------------------------------------------

def _render_node(
    slide,
    node: dict,
    region: tuple[float, float, float, float],
    brand_tokens: dict | None,
) -> _RenderedShape:
    """Add one node shape to the slide. Returns metadata for connector drawing."""
    shape_kind = node.get("shape", "rectangle")
    if shape_kind not in NODE_SHAPE_MAP:
        raise ValueError(
            f"unknown node shape: {shape_kind!r}; valid: {sorted(NODE_SHAPE_MAP.keys())}"
        )
    mso_shape = NODE_SHAPE_MAP[shape_kind]

    x, y = _transform_coords(float(node["x"]), float(node["y"]), region)
    w = float(node["w"])
    h = float(node["h"])

    sp = slide.shapes.add_shape(
        mso_shape,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )

    # Fill + outline
    fill_color_name = node.get("fill_color", DEFAULT_NODE_FILL)
    text_color_name = node.get("text_color", DEFAULT_NODE_TEXT)

    if shape_kind == "swimlane":
        # Swimlanes are containers — outline only, no fill.
        sp.fill.background()
        sp.line.color.rgb = resolve_color(
            DEFAULT_SWIMLANE_BORDER, brand_tokens, DEFAULT_SWIMLANE_BORDER,
        )
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = resolve_color(
            fill_color_name, brand_tokens, DEFAULT_NODE_FILL,
        )
        sp.line.fill.background()  # no outline on filled nodes by default

    # Label
    label = node.get("label", "")
    if label:
        tf = sp.text_frame
        tf.text = label
        # Contain the label inside the node box. word_wrap alone wraps
        # but does not shrink; the M3 attempt used `auto_size=
        # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE`, but LibreOffice does NOT
        # honor a bare <a:normAutofit/> — the rendered text still spills
        # past the box (ibd_phage_targeting draft_1 slides 6/10/19).
        # M4a Tier A replaces it with `_apply_fontscale_to_shape`, which
        # writes an explicit `<a:normAutofit fontScale="...">` — the
        # only shrink-to-fit LibreOffice computes at render. The ladder
        # is tuned for short node labels (~40 chars is the design size).
        tf.word_wrap = True
        _apply_fontscale_to_shape(sp, len(label),
                                  ladder=((40, 100000), (60, 90000),
                                          (90, 80000), (120, 70000)),
                                  full_below=40)
        # Color the text based on text_color_name
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = resolve_color(
                    text_color_name, brand_tokens, DEFAULT_NODE_TEXT,
                )
                run.font.size = Pt(14)

    return _RenderedShape(
        node_id=node["id"],
        shape=sp,
        cx=x + w / 2,
        cy=y + h / 2,
        w=w,
        h=h,
    )


# ---------------------------------------------------------------------------
# Edge rendering
# ---------------------------------------------------------------------------

def _render_edge_line(
    slide,
    edge: dict,
    rendered_by_id: dict[str, _RenderedShape],
    brand_tokens: dict | None,
) -> None:
    """Add the connector LINE only (no label) between two node shapes.

    Edge labels render in a separate third pass (see render_diagram) so
    they paint on top of the node fills. The previous implementation
    rendered line + label in the same edge pass; because edges run
    BEFORE nodes (M3 fix #54: node fills must occlude edge endpoints),
    labels ended up UNDER the node boxes — visible as missing labels on
    ibd_phage_targeting draft_1 slides 10/19 (M3 Tier-A-deferred item).
    """
    edge_kind = edge.get("kind", "straight")
    if edge_kind not in EDGE_KIND_MAP:
        raise ValueError(
            f"unknown edge kind: {edge_kind!r}; valid: {sorted(EDGE_KIND_MAP.keys())}"
        )
    mso_connector = EDGE_KIND_MAP[edge_kind]

    src_id = edge.get("from")
    dst_id = edge.get("to")
    if src_id not in rendered_by_id or dst_id not in rendered_by_id:
        # Skip edges referencing missing nodes — diagram_design prompt
        # should never produce these, but the validator catches it as
        # well via slide_spec validation.
        return

    src = rendered_by_id[src_id]
    dst = rendered_by_id[dst_id]

    # Use shape edge midpoints for cleaner endpoints than centers.
    # Crude heuristic: aim from src center toward dst center, but
    # snap to the nearest shape edge.
    line = slide.shapes.add_connector(
        mso_connector,
        Inches(src.cx), Inches(src.cy),
        Inches(dst.cx), Inches(dst.cy),
    )
    # Color the line slate-dark (M4a Tier E round 2). Was graphite_gray,
    # which washed out against the master's cream tone — see _CONNECTOR_RGB.
    line.line.color.rgb = _CONNECTOR_RGB


def _render_edge_label(
    slide,
    edge: dict,
    rendered_by_id: dict[str, _RenderedShape],
    brand_tokens: dict | None,
) -> None:
    """Add the edge LABEL textbox (no line) on top of the rendered nodes.

    Called from the third pass in render_diagram, AFTER nodes have
    painted. Splits the M3 fix #54 'edges-first, nodes-on-top' constraint
    so node fills occlude edge endpoints (clean visual) BUT labels still
    appear on top of the diagram instead of behind it (M4a Tier A3).

    M4a Tier E round 3 (2026-05-23): geometry now uses NODE-EDGE
    midpoint (not center-to-center) + node widths to size the textbox
    to fit the inter-node GAP. Round-2 visual-QA found the pre-round-3
    heuristic (0.7in label width centered on a center-to-center
    midpoint) put labels ON the next node box when the gap was
    < ~0.7in (slide 6 LDA→K=4 edge — label `max ARI at K=4`
    fragmented as `max AR` / `at K=4` because half of it spilled
    behind the K=4 box). The gap-based placement keeps the label
    inside the gap.

    For horizontal edges:  label spans the gap between src.right and
                            dst.left; width = max(0.4, gap - 0.10);
                            placed above the line (label_y < mid_y).
    For vertical edges:    same idea on Y axis.
    For diagonal edges:    fall back to the prior heuristic (the
                            label sits in open space anyway).
    """
    src_id = edge.get("from")
    dst_id = edge.get("to")
    if src_id not in rendered_by_id or dst_id not in rendered_by_id:
        return
    label = edge.get("label", "")
    if not label:
        return

    src = rendered_by_id[src_id]
    dst = rendered_by_id[dst_id]
    mid_x = (src.cx + dst.cx) / 2
    mid_y = (src.cy + dst.cy) / 2
    dy = abs(dst.cy - src.cy)
    dx = abs(dst.cx - src.cx)

    # Default textbox geometry (fallback when node widths unknown).
    label_w = 0.70
    label_h = 0.25

    if dy < 0.3 and src.w > 0 and dst.w > 0:
        # Horizontal edge — place label ABOVE the node row, sized to
        # the gap when the gap is wide enough for the label to read,
        # otherwise sized wider than the gap (overflows the gap but
        # sits above the node-tops and is therefore visible).
        #
        # Compute gap edges from node half-widths around their centers.
        if src.cx < dst.cx:
            src_right = src.cx + src.w / 2
            dst_left = dst.cx - dst.w / 2
        else:
            src_right = dst.cx + dst.w / 2
            dst_left = src.cx - src.w / 2
        gap = max(0.0, dst_left - src_right)
        # M4a Tier E round 4 (2026-05-24): raised the in-gap threshold
        # from 0.4 to 1.0in. Round-3 used >=0.4in and word_wrap=True,
        # which produced one-character-per-line vertical stacks on the
        # ibd_phage_targeting slide-6 lda_gmm→k4 edge (gap = exactly
        # 0.4in; label "max ARI at K=4" needs ~1.4in at 9pt to render
        # on one line, so wrapping into a 0.40in-wide box stacks every
        # character). Below 1.0in we place the label ABOVE the node
        # row at the connector midpoint with a 1.1in width — clears
        # the node-tops cleanly (y = mid_y - 0.40 = above node-top y =
        # cy - 0.45 by ~0.05in) and reads as one line.
        # Compute node-top y from src geometry (src.cy is center; src.h
        # is height). Both endpoints share y for horizontal edges, so
        # src_top == dst_top.
        node_top = src.cy - src.h / 2
        if gap >= 1.0:
            label_w = max(0.6, gap - 0.10)   # margin: 0.05in each side
            # x: center the textbox in the gap
            label_x = src_right + (gap - label_w) / 2
            # y: just above the node-top edge (the label sits in the
            # gap horizontally so the node-row vertical position is OK,
            # but go slightly above so the connector arrow head is
            # visible underneath).
            label_y = max(0.10, node_top - label_h - 0.05)
        else:
            # Narrow gap — sit the label fully ABOVE the node row so
            # horizontal overlap with the nodes does not matter (the
            # label is in the empty space above the workflow row).
            # M4a Tier E round 5 (2026-05-24): round-4's mid_y - 0.40
            # was INSIDE the node-row vertical extent on slides 6/10/19
            # (nodes h=0.9 → node_top=cy-0.45 → label_y=mid_y-0.40 is
            # only 0.05in above node_top, and label-bottom mid_y-0.15
            # falls inside the node — text rendered ON the node fill,
            # narrowly visible between adjacent nodes as the vertical
            # strip the visual-QA model reported).
            label_w = 1.10
            label_x = mid_x - label_w / 2
            label_y = max(0.10, node_top - label_h - 0.05)
    elif dx < 0.3 and src.h > 0 and dst.h > 0:
        # Vertical edge — place label in the inter-node gap to the right.
        if src.cy < dst.cy:
            src_bottom = src.cy + src.h / 2
            dst_top = dst.cy - dst.h / 2
        else:
            src_bottom = dst.cy + dst.h / 2
            dst_top = src.cy - src.h / 2
        gap_y = max(0.0, dst_top - src_bottom)
        if gap_y >= 0.3:
            label_h = max(0.20, min(0.40, gap_y - 0.05))
            label_x = mid_x + 0.10
            label_y = src_bottom + (gap_y - label_h) / 2
        else:
            label_x = mid_x + 0.10
            label_y = mid_y - 0.15
    else:
        # Diagonal edge (or w/h unknown) — prior offset heuristic.
        label_x = mid_x - 0.20
        label_y = mid_y - 0.30
    tb = slide.shapes.add_textbox(
        Inches(label_x), Inches(label_y),
        Inches(label_w), Inches(label_h),
    )
    tb.text_frame.text = label
    # Tier E round 4 (2026-05-24): word_wrap REMOVED. Round-3 used
    # word_wrap=True, but combined with a narrow gap-sized box on
    # slide-6 (0.4in wide), python-pptx fell back to char-by-char
    # wrap (no word fit) and rendered every label character on its
    # own line — visible as a vertical character stack. The new
    # round-4 geometry sizes the box wide enough (1.1in min for
    # narrow gaps, gap-width-minus-margin for wide gaps) that
    # single-line rendering is correct; word_wrap=False lets the
    # textbox draw cleanly without falling back to char-wrap if a
    # corner case slips past.
    tb.text_frame.word_wrap = False
    for paragraph in tb.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(9)
            # Slate-dark for contrast against the (post-watermark-strip)
            # cream background. Bumped 2026-05-23 (M4a Tier E round 2).
            run.font.color.rgb = _CONNECTOR_RGB


# ---------------------------------------------------------------------------
# Top-level render
# ---------------------------------------------------------------------------

def render_diagram(
    slide,
    diagram: dict,
    region: tuple[float, float, float, float],
    brand_tokens: dict | None = None,
) -> None:
    """Render a slide_spec diagram into a python-pptx slide.

    Args:
      slide: python-pptx Slide object.
      diagram: slide_spec diagram dict (kind, nodes, edges).
      region: (left, top, width, height) in inches — the slide region
              the diagram fits into. Diagram-local (x, y) coords map to
              region.left + x, region.top + y.
      brand_tokens: parsed kbase-brand-tokens.json (or None for hardcoded
                    fallback hexes).

    Raises:
      ValueError on unknown shape / edge kind. (Schema validator should
      have caught these before render.)
    """
    if diagram.get("kind") != "boxes_and_arrows":
        raise ValueError(
            f"unsupported diagram kind: {diagram.get('kind')!r} "
            f"(only 'boxes_and_arrows' supported in v0.1; "
            f"see SPEC §6 / DECISIONS D-028 for v0.2 plans)"
        )

    # Render order matters for visual quality (python-pptx z-order = paint
    # order; later shapes on top). M4a Tier A3 keeps the M3 fix #54
    # constraint AND fixes the M3-deferred edge-label-behind-nodes issue
    # (ibd_phage_targeting draft_1 slides 10/19):
    #
    #   Pass 1: edge LINES (computed from node centers, no shapes yet).
    #   Pass 2: NODE shapes — fills occlude edge endpoints (the lines
    #           appear to "enter" the node edge, clean visual).
    #   Pass 3: edge LABELS — paint on top of nodes so labels are never
    #           hidden behind a node box even when the inter-node gap is
    #           tight (slide 10 in the M3 smoke). The 2026-04-27 fix #54
    #           offset heuristic still keeps labels in the gap.
    nodes = diagram.get("nodes", []) or []
    edges = diagram.get("edges", []) or []

    # Pre-compute geometry (no slide mutation yet). Carries width + height
    # alongside centers so the line pass can build _RenderedShape proxies
    # with the same w/h that the actual nodes will have — the third pass
    # (edge labels) needs node widths to size the label textbox to the
    # inter-node gap (M4a Tier E round 3).
    geom_by_id: dict[str, tuple[float, float, float, float]] = {}
    for node in nodes:
        node_id = node.get("id", "")
        x, y = _transform_coords(float(node["x"]), float(node["y"]), region)
        w = float(node["w"])
        h = float(node["h"])
        geom_by_id[node_id] = (x + w / 2, y + h / 2, w, h)

    # Pass 1: render edge LINES using computed centers
    edge_proxy: dict[str, _RenderedShape] = {
        nid: _RenderedShape(node_id=nid, shape=None,
                            cx=cx, cy=cy, w=w, h=h)
        for nid, (cx, cy, w, h) in geom_by_id.items()
    }
    for edge in edges:
        _render_edge_line(slide, edge, edge_proxy, brand_tokens)

    # Pass 2: render node shapes on top of edges
    rendered_by_id: dict[str, _RenderedShape] = {}
    for node in nodes:
        rs = _render_node(slide, node, region, brand_tokens)
        rendered_by_id[rs.node_id] = rs

    # Pass 3: render edge LABELS on top of nodes
    for edge in edges:
        _render_edge_label(slide, edge, rendered_by_id, brand_tokens)


# ---------------------------------------------------------------------------
# Brand-tokens loader
# ---------------------------------------------------------------------------

def load_brand_tokens(path: Path | str | None = None) -> dict | None:
    """Convenience: load kbase-brand-tokens.json from the shipped
    references dir (default) or a user-supplied path."""
    if path is None:
        here = Path(__file__).resolve().parent
        path = here.parent / "references" / "kbase-brand-tokens.json"
    p = Path(path)
    if not p.is_file():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return None


# ---------------------------------------------------------------------------
# CLI (smoke testing)
# ---------------------------------------------------------------------------

def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="diagram_render",
        description="Render a slide_spec diagram dict into a one-slide pptx "
                    "for visual smoke testing.",
    )
    parser.add_argument("diagram_json",
                        help="Path to a JSON file containing a single diagram dict.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    parser.add_argument("--brand-tokens",
                        help="Path to kbase-brand-tokens.json "
                             "(default: shipped one).")
    parser.add_argument("--region", default="0.5,1.4,9.0,3.5",
                        help="Comma-separated 4-tuple in inches "
                             "(left,top,width,height). Default: 0.5,1.4,9.0,3.5")
    args = parser.parse_args(argv)

    diagram = json.loads(Path(args.diagram_json).read_text(encoding="utf-8"))
    region = tuple(float(x) for x in args.region.split(","))
    if len(region) != 4:
        print("--region must be 4 comma-separated floats", file=sys.stderr)
        return 2
    tokens = load_brand_tokens(args.brand_tokens)

    # Build a minimal pptx
    from pptx import Presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    render_diagram(slide, diagram, region, tokens)
    prs.save(args.out)
    print(f"wrote {args.out}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
